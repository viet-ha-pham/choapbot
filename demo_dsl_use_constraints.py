"""
PML/PSL Renderer v1.0
=====================

Adds support for:
- presentation title slide
- section header slide
- content slides
- title + image slide layout
- text + image two-column layouts
- grid layouts: grid-3, grid-4, grid-5, grid-6, grid-4x2
- stair-progress layout for step-by-step progress/process slides
- stacked-stairs layout: overlapping, shrinking, one-side-aligned steps
- timeline layout: horizontal axis with alternating milestones
- icon bullets: shared icon and per-item icon bullets
- text overflow handling: wrap, shrink, ellipsis/clip with max_lines/min_size
- hyperlinks: link blocks and url/href on text/bullet/card/milestone items
- tables: title-table layout with headers/rows, HTML table and PPTX table rendering
- numbered columns layouts: numbered-columns-3..6 with prominent circular step numbers and bullets inside cards
- auto pagination for long BulletList/TableBox objects with overflow: paginate
- content-aware pagination: splits when estimated rendered text height exceeds object box height
- two image layout: two-images
- hero image layout: hero-image (large image + caption)
- image caption layout: image-caption (image + caption only, no title/subtitle)
- footer images: footer_image / footer_images on content slides
- footer text: footer_text or footer.text on content slides
- background images: background_image or background.image on slides/layouts
- image fit modes: stretch, contain, cover/crop, original
- text styles: bold, italic, underline, align left/center/right/justify
- PCL: Presentation Constraint Language for slide/layout/object constraints

Install:
    pip install python-pptx

Run:
    python demo_dsl_tables.py

Outputs:
    demo_output.html
    demo_output.pptx
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import html
import re

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    PPTXPresentation = None
    Inches = Pt = RGBColor = PP_ALIGN = MSO_SHAPE = None


# -----------------------------------------------------------------------------
# Data models
# -----------------------------------------------------------------------------

@dataclass
class Slide:
    title: str
    layout: str = "title-bullets"
    intent: Optional[str] = None
    blocks: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Section:
    title: str
    header_layout: str = "section-header"
    header: Dict[str, Any] = field(default_factory=dict)
    slides: List[Slide] = field(default_factory=list)


@dataclass
class PresentationDoc:
    title: str
    meta: Dict[str, Any] = field(default_factory=dict)
    style_file: Optional[str] = None
    constraint_file: Optional[str] = None
    cover_layout: str = "title-slide"
    cover: Dict[str, Any] = field(default_factory=dict)
    sections: List[Section] = field(default_factory=list)


@dataclass
class Theme:
    name: str
    page: Dict[str, Any] = field(default_factory=dict)
    colors: Dict[str, str] = field(default_factory=dict)
    fonts: Dict[str, str] = field(default_factory=dict)
    presentation_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    section_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)


@dataclass
class ConstraintSet:
    """PCL model.

    PCL stores constraint blocks by scope, then apply_constraints() merges those
    blocks onto render objects before pagination/rendering.

    Supported scopes:
      default:         applies to every slide
      presentation.*:  applies to presentation/title-slide layouts
      section.*:       applies to section-header layouts
      layout.*:        applies to content slide layouts
      slide.*:         applies to slide title exact match
    """
    name: str
    defaults: Dict[str, Any] = field(default_factory=dict)
    presentation_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    section_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    slides: Dict[str, Dict[str, Any]] = field(default_factory=dict)



# -----------------------------------------------------------------------------
# Low-level indentation parser
# -----------------------------------------------------------------------------

@dataclass
class Line:
    indent: int
    text: str
    lineno: int


def strip_comment(line: str) -> str:
    """Remove comments outside quotes."""
    in_quote = False
    quote_char = ""
    out = []

    for ch in line:
        if ch in ('"', "'"):
            if not in_quote:
                in_quote = True
                quote_char = ch
            elif quote_char == ch:
                in_quote = False
        if ch == "#" and not in_quote:
            break
        out.append(ch)

    return "".join(out).rstrip()


def tokenize_lines(text: str) -> List[Line]:
    lines: List[Line] = []
    for i, raw in enumerate(text.splitlines(), start=1):
        no_comment = strip_comment(raw)
        if not no_comment.strip():
            continue
        indent = len(no_comment) - len(no_comment.lstrip(" "))
        if "\t" in raw[:indent]:
            raise SyntaxError(f"Line {i}: tabs are not allowed for indentation")
        lines.append(Line(indent=indent, text=no_comment.strip(), lineno=i))
    return lines


def parse_scalar(value: str) -> Any:
    value = value.strip()

    if value == "":
        return ""

    if (value.startswith('"') and value.endswith('"')) or (
        value.startswith("'") and value.endswith("'")
    ):
        return value[1:-1]

    if value.lower() == "true":
        return True
    if value.lower() == "false":
        return False
    if value.lower() in {"null", "none"}:
        return None

    if value.startswith("[") and value.endswith("]"):
        inner = value[1:-1].strip()
        if not inner:
            return []
        return [parse_scalar(part.strip()) for part in inner.split(",")]

    try:
        if "." in value:
            return float(value)
        return int(value)
    except ValueError:
        return value


def split_key_value(text: str) -> Tuple[str, Optional[str]]:
    if ":" not in text:
        return text, None
    key, val = text.split(":", 1)
    return key.strip(), val.strip()


def parse_generic_block(lines: List[Line], start: int, indent: int) -> Tuple[Any, int]:
    result: Dict[str, Any] = {}
    items: List[Any] = []
    mode: Optional[str] = None
    i = start

    while i < len(lines):
        line = lines[i]
        if line.indent < indent:
            break
        if line.indent > indent:
            raise SyntaxError(
                f"Line {line.lineno}: unexpected indentation. Expected {indent}, got {line.indent}"
            )

        text = line.text

        if text.startswith("- "):
            if mode is None:
                mode = "list"
            elif mode != "list":
                raise SyntaxError(f"Line {line.lineno}: cannot mix list and mapping")

            item_text = text[2:].strip()

            # Supports list scalars:
            #   - item
            # Supports list mappings:
            #   - heading: A
            #     text: B
            # Supports nested named mappings:
            #   - image:
            #       src: a.png
            if item_text.endswith(":"):
                key = item_text[:-1].strip()
                child, new_i = parse_generic_block(lines, i + 1, indent + 2)
                items.append({key: child})
                i = new_i
            elif ":" in item_text:
                key, val = split_key_value(item_text)
                item: Dict[str, Any] = {key: parse_scalar(val or "")}
                new_i = i + 1
                if new_i < len(lines) and lines[new_i].indent > indent:
                    child, new_i = parse_generic_block(lines, new_i, lines[new_i].indent)
                    if isinstance(child, dict):
                        item.update(child)
                    else:
                        item["children"] = child
                items.append(item)
                i = new_i
            else:
                items.append(parse_scalar(item_text))
                i += 1
            continue

        if mode is None:
            mode = "dict"
        elif mode != "dict":
            raise SyntaxError(f"Line {line.lineno}: cannot mix mapping and list")

        key, val = split_key_value(text)
        if val is None:
            text_lines = []
            while i < len(lines) and lines[i].indent == indent:
                text_lines.append(lines[i].text)
                i += 1
            return "\n".join(text_lines), i

        if val != "":
            result[key] = parse_scalar(val)
            i += 1
        else:
            if i + 1 >= len(lines) or lines[i + 1].indent <= indent:
                result[key] = {}
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                result[key] = child
                i = new_i

    if mode == "list":
        return items, i
    return result, i


# -----------------------------------------------------------------------------
# PML / PSL parser
# -----------------------------------------------------------------------------

PRESENTATION_RE = re.compile(r'^presentation\s+"(.+)"\s*:$')
SECTION_RE = re.compile(r'^section\s+"(.+)"\s*:$')
SLIDE_RE = re.compile(r'^slide\s+"(.+)"\s*:$')
THEME_RE = re.compile(r'^theme\s+"(.+)"\s*:$')
PCL_RE = re.compile(r'^(?:pcl|constraints)\s+"(.+)"\s*:$')


def parse_pml(text: str) -> PresentationDoc:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PML")

    m = PRESENTATION_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PML must start with: presentation "Title":')

    doc = PresentationDoc(title=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        text = line.text
        key, val = split_key_value(text)

        if text == "meta:":
            block, i = parse_generic_block(lines, i + 1, 4)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: meta must be a mapping")
            doc.meta = block
            continue

        if key == "cover_layout":
            doc.cover_layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if text == "cover:":
            block, i = parse_generic_block(lines, i + 1, 4)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: cover must be a mapping")
            doc.cover = block
            continue

        if text.startswith("use style:"):
            _, value = split_key_value(text)
            doc.style_file = parse_scalar(value or "")
            i += 1
            continue

        if text.startswith("use constraints:"):
            _, value = split_key_value(text)
            doc.constraint_file = parse_scalar(value or "")
            i += 1
            continue

        sec_match = SECTION_RE.match(text)
        if sec_match:
            section, i = parse_section(lines, i)
            doc.sections.append(section)
            continue

        raise SyntaxError(f"Line {line.lineno}: unknown PML statement: {text}")

    return doc


def parse_section(lines: List[Line], start: int) -> Tuple[Section, int]:
    m = SECTION_RE.match(lines[start].text)
    assert m is not None
    section = Section(title=m.group(1))
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= lines[start].indent:
            break
        if line.indent != lines[start].indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected section child indent")

        key, val = split_key_value(line.text)

        if key == "header_layout":
            section.header_layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if line.text == "header:":
            block, i = parse_generic_block(lines, i + 1, line.indent + 2)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: section header must be a mapping")
            section.header = block
            continue

        slide_match = SLIDE_RE.match(line.text)
        if slide_match:
            slide, i = parse_slide(lines, i)
            section.slides.append(slide)
        else:
            raise SyntaxError(f"Line {line.lineno}: section only accepts header/header_layout/slide blocks")

    return section, i


def parse_slide(lines: List[Line], start: int) -> Tuple[Slide, int]:
    m = SLIDE_RE.match(lines[start].text)
    assert m is not None
    slide = Slide(title=m.group(1))
    base_indent = lines[start].indent
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= base_indent:
            break
        if line.indent != base_indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected slide child indent")

        key, val = split_key_value(line.text)

        if key == "layout":
            slide.layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if key == "intent":
            slide.intent = str(parse_scalar(val or ""))
            i += 1
            continue

        if val == "":
            if i + 1 >= len(lines) or lines[i + 1].indent <= line.indent:
                slide.blocks[key] = ""
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                slide.blocks[key] = child
                i = new_i
        else:
            slide.blocks[key] = parse_scalar(val or "")
            i += 1

    return slide, i


def parse_psl(text: str) -> Theme:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PSL")

    m = THEME_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PSL must start with: theme "Name":')

    theme = Theme(name=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        key, val = split_key_value(line.text)
        if val != "":
            raise SyntaxError(f"Line {line.lineno}: top-level PSL keys must be blocks")

        block, i = parse_generic_block(lines, i + 1, 4)
        if not isinstance(block, dict):
            raise SyntaxError(f"Line {line.lineno}: top-level PSL block must be a mapping")

        if key == "page":
            theme.page = block
        elif key == "colors":
            theme.colors = block
        elif key == "fonts":
            theme.fonts = block
        elif key.startswith("presentation."):
            name = key.split(".", 1)[1]
            theme.presentation_layouts[name] = block
        elif key.startswith("section."):
            name = key.split(".", 1)[1]
            theme.section_layouts[name] = block
        elif key.startswith("slide."):
            name = key.split(".", 1)[1]
            theme.layouts[name] = block
        else:
            theme.layouts[key] = block

    return theme




def parse_pcl(text: str) -> ConstraintSet:
    """Parse PCL: Presentation Constraint Language.

    Example:

        pcl "safe-layout":
          default:
            bullets:
              overflow: paginate
              min_size: 14
            image:
              mode: contain

          layout "title-bullets":
            bullets:
              max_lines: 8

          slide "Danh sách dài":
            bullets:
              overflow: paginate

    PCL does not replace PSL. PSL describes intended style; PCL describes
    safety/validity constraints that are applied after IR construction.
    """
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PCL")

    m = PCL_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PCL must start with: pcl "Name": or constraints "Name":')

    constraints = ConstraintSet(name=m.group(1))
    i = 1
    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        text = line.text
        key, val = split_key_value(text)
        if val not in ("", None):
            raise SyntaxError(f"Line {line.lineno}: top-level PCL keys must be blocks")

        block, i = parse_generic_block(lines, i + 1, 4)
        if not isinstance(block, dict):
            raise SyntaxError(f"Line {line.lineno}: top-level PCL block must be a mapping")

        if key == "default":
            constraints.defaults = block
            continue

        layout_match = re.match(r'^layout\s+"(.+)"$', key)
        slide_match = re.match(r'^slide\s+"(.+)"$', key)
        pres_match = re.match(r'^presentation\s+"(.+)"$', key)
        section_match = re.match(r'^section\s+"(.+)"$', key)

        if layout_match:
            constraints.layouts[layout_match.group(1)] = block
        elif slide_match:
            constraints.slides[slide_match.group(1)] = block
        elif pres_match:
            constraints.presentation_layouts[pres_match.group(1)] = block
        elif section_match:
            constraints.section_layouts[section_match.group(1)] = block
        elif key.startswith("layout."):
            constraints.layouts[key.split(".", 1)[1]] = block
        elif key.startswith("slide."):
            constraints.slides[key.split(".", 1)[1]] = block
        elif key.startswith("presentation."):
            constraints.presentation_layouts[key.split(".", 1)[1]] = block
        elif key.startswith("section."):
            constraints.section_layouts[key.split(".", 1)[1]] = block
        else:
            constraints.layouts[key] = block

    return constraints

# -----------------------------------------------------------------------------
# Render IR construction
# -----------------------------------------------------------------------------

def resolve_token(value: Any, theme: Theme) -> Any:
    if isinstance(value, str):
        if value in theme.colors:
            return theme.colors[value]
        if value in theme.fonts:
            return theme.fonts[value]
    return value


def style_value(style: Dict[str, Any], key: str, theme: Theme, default: Any = None) -> Any:
    return resolve_token(style.get(key, default), theme)


def style_value_from_item(item: Dict[str, Any], style: Dict[str, Any], key: str, theme: Theme, default: Any = None) -> Any:
    """Resolve a style value with per-item override first, then layout style.

    Use this for card-like structures where a cell/step/milestone can override
    the layout's default fill/border, for example:

      cells:
        - heading: Ingest
          text: Read source files
          fill: "#E0F2FE"
          border: primary

    `fill` is for card/shape background; `background` remains slide-level.
    """
    if isinstance(item, dict) and key in item:
        return resolve_token(item.get(key), theme)
    return style_value(style, key, theme, default)


def border_value_from_item(item: Dict[str, Any], style: Dict[str, Any], theme: Theme, default: Any = None) -> Any:
    """Resolve card border with aliases: border, border_color, stroke."""
    for key in ("border", "border_color", "stroke"):
        if isinstance(item, dict) and key in item:
            return resolve_token(item.get(key), theme)
    for key in ("border", "border_color", "stroke"):
        if key in style:
            return resolve_token(style.get(key), theme)
    return default


def normalize_text_block(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, dict):
        if "text" in value:
            return str(value["text"])
        if "heading" in value:
            return str(value["heading"])
    return str(value)



def parse_bool(value: Any, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    if isinstance(value, (int, float)):
        return bool(value)
    text = str(value).strip().lower()
    if text in {"true", "yes", "y", "1", "on"}:
        return True
    if text in {"false", "no", "n", "0", "off"}:
        return False
    return default


def normalize_align(value: Any, default: str = "left") -> str:
    align = str(value or default).strip().lower()
    aliases = {
        "middle": "center",
        "centre": "center",
        "justified": "justify",
        "both": "justify",
    }
    align = aliases.get(align, align)
    if align not in {"left", "center", "right", "justify"}:
        return default
    return align


def normalize_overflow(value: Any, default: str = "wrap") -> str:
    mode = str(value or default).strip().lower()
    aliases = {
        "auto": "wrap",
        "truncate": "ellipsis",
        "elide": "ellipsis",
        "hidden": "clip",
        "cut": "clip",
        "fit": "shrink",
        "auto_shrink": "shrink",
        "auto_paginate": "paginate",
        "paginate_auto": "paginate",
        "page": "paginate",
        "paging": "paginate",
    }
    mode = aliases.get(mode, mode)
    if mode not in {"wrap", "shrink", "ellipsis", "clip", "paginate"}:
        return default
    return mode


def safe_int(value: Any, default: int) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def safe_float(value: Any, default: float) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def text_overflow_props(style: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text overflow policy from style.

    Supported PSL keys:
      overflow: wrap | shrink | ellipsis | clip | paginate
      max_lines: 3
      min_size: 14
      auto_shrink: true

    Notes:
    - HTML uses CSS wrapping/clamping where possible.
    - PPTX uses an approximate line estimator because python-pptx does not provide
      actual rendered text measurement.
    """
    overflow = normalize_overflow(style.get("overflow", "wrap"))
    if parse_bool(style.get("auto_shrink", False)):
        overflow = "shrink"
    return {
        "overflow": overflow,
        "max_lines": safe_int(style.get("max_lines", 0), 0),
        "min_size": safe_int(style.get("min_size", 12), 12),
        "auto_shrink": parse_bool(style.get("auto_shrink", overflow == "shrink")),
    }


def text_style_props(style: Dict[str, Any], default_align: str = "left") -> Dict[str, Any]:
    """Extract common text style flags from a PSL style mapping."""
    return {
        "bold": parse_bool(style.get("bold", False)),
        "italic": parse_bool(style.get("italic", False)),
        "underline": parse_bool(style.get("underline", False)),
        "align": normalize_align(style.get("align", style.get("text_align", default_align)), default_align),
        **text_overflow_props(style),
    }

def get_url(value: Any, style: Optional[Dict[str, Any]] = None) -> str:
    """Extract a hyperlink URL from PML value or PSL style.

    Supported keys: url, href, link.
    """
    style = style or {}
    if isinstance(value, dict):
        for key in ("url", "href", "link"):
            if value.get(key):
                return str(value[key])
    for key in ("url", "href", "link"):
        if style.get(key):
            return str(style[key])
    return ""


def normalize_link_block(value: Any) -> Tuple[str, str]:
    """Return (text, url) for a PML link block.

    Supported PML:
      link:
        text: "Open docs"
        url: "https://example.com"

      link: "https://example.com"
    """
    if isinstance(value, dict):
        url = get_url(value)
        text = str(value.get("text", value.get("label", value.get("title", url))))
        return text, url
    url = str(value)
    return url, url


def make_text_object(kind: str, value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [60, 40])
    url = get_url(value, style)
    obj = {
        "type": "TextBox",
        "role": kind,
        "text": normalize_text_block(value),
        "url": url,
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 80),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 28),
        "color": style_value(style, "color", theme, "#000000"),
        **text_style_props(style),
    }
    if url and "underline" not in style:
        obj["underline"] = True
    return obj


def make_link_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    text, url = normalize_link_block(value)
    obj = make_text_object("link", {"text": text, "url": url}, style, theme)
    if url:
        obj["url"] = url
        obj["underline"] = parse_bool(style.get("underline", True), True)
    return obj


def normalize_table(value: Any) -> Tuple[List[str], List[List[str]]]:
    """Normalize a PML table block into (headers, rows).

    Supported PML:
      table:
        headers: [Giai đoạn, Việc chính, Kết quả]
        rows:
          - [Q1, Khởi tạo, Phạm vi dự án]
          - [Q2, Chuẩn hóa, Dataset sạch]

      table:
        columns:
          - Giai đoạn
          - Việc chính
        rows:
          - cells: [Q1, Khởi tạo]
          - cells: [Q2, Chuẩn hóa]

      table:
        rows:
          - [A, B]
          - [C, D]
    """
    if not isinstance(value, dict):
        return [], []

    headers_raw = value.get("headers", value.get("columns", []))
    if headers_raw is None:
        headers_raw = []
    if not isinstance(headers_raw, list):
        headers = [str(headers_raw)]
    else:
        headers = [str(x) for x in headers_raw]

    rows_raw = value.get("rows", [])
    if not isinstance(rows_raw, list):
        rows_raw = [rows_raw]

    rows: List[List[str]] = []
    for row in rows_raw:
        if isinstance(row, dict):
            cells = row.get("cells", row.get("values", row.get("row", [])))
            if not isinstance(cells, list):
                cells = [cells]
            rows.append([str(x) for x in cells])
        elif isinstance(row, list):
            rows.append([str(x) for x in row])
        else:
            rows.append([str(row)])

    # Make all rows rectangular.
    col_count = max([len(headers)] + [len(r) for r in rows] + [0])
    if col_count == 0:
        return headers, rows
    if headers and len(headers) < col_count:
        headers = headers + [""] * (col_count - len(headers))
    normalized_rows = []
    for row in rows:
        if len(row) < col_count:
            row = row + [""] * (col_count - len(row))
        normalized_rows.append(row[:col_count])
    return headers[:col_count], normalized_rows


def make_table_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [70, 150])
    headers, rows = normalize_table(value)

    # Merge style from PSL with table-level policy from PML.
    # This lets authors write `overflow: paginate` directly inside `table:`.
    table_cfg = value if isinstance(value, dict) else {}
    effective_style = dict(style)
    for key in ("overflow", "max_rows_per_slide", "rows_per_slide"):
        if key in table_cfg:
            effective_style[key] = table_cfg[key]

    return {
        "type": "TableBox",
        "headers": headers,
        "rows": rows,
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1140),
        "h": style.get("height", 430),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 18),
        "color": style_value(style, "color", theme, "#111827"),
        "header_fill": style_value(style, "header_fill", theme, "#003A8C"),
        "header_color": style_value(style, "header_color", theme, "#FFFFFF"),
        "cell_fill": style_value(style, "cell_fill", theme, "#FFFFFF"),
        "alt_fill": style_value(style, "alt_fill", theme, "#F8FAFC"),
        "border_color": style_value(style, "border_color", theme, "#CBD5E1"),
        "align": normalize_align(style.get("align", "left")),
        "header_align": normalize_align(style.get("header_align", style.get("align", "center"))),
        "bold": parse_bool(style.get("bold", False)),
        "header_bold": parse_bool(style.get("header_bold", True), True),
        "overflow": normalize_overflow(effective_style.get("overflow", "wrap")),
        "max_rows_per_slide": safe_int(effective_style.get("max_rows_per_slide", effective_style.get("rows_per_slide", 0)), 0),
    }


BULLET_ICON_MAP = {
    "check": "✓",
    "tick": "✓",
    "done": "✓",
    "arrow": "→",
    "right": "→",
    "star": "★",
    "diamond": "◆",
    "square": "■",
    "dot": "●",
    "circle": "●",
    "dash": "—",
    "plus": "+",
    "pin": "📌",
    "gear": "⚙",
    "warning": "⚠",
    "none": "",
}


def resolve_bullet_icon(icon: Any, default: str = "") -> str:
    if icon is None:
        return default
    icon_text = str(icon).strip()
    if not icon_text:
        return default
    return BULLET_ICON_MAP.get(icon_text.lower(), icon_text)


def normalize_bullet_items(value: Any, default_icon: str = "") -> List[Dict[str, Any]]:
    """Normalize bullets to [{icon, text}].

    Supported PML:
      bullets:
        - A
        - B

      bullets:
        icon: check
        items:
          - A
          - B

      bullets:
        - icon: check
          text: A
        - icon: warning
          text: B
    """
    if isinstance(value, dict):
        default_icon = resolve_bullet_icon(value.get("icon"), default_icon)
        raw_items = value.get("items", value.get("bullets", []))
        if not isinstance(raw_items, list):
            raw_items = [raw_items]
    elif isinstance(value, list):
        raw_items = value
    else:
        raw_items = [value]

    items: List[Dict[str, str]] = []
    for raw in raw_items:
        if isinstance(raw, dict):
            icon = resolve_bullet_icon(raw.get("icon"), default_icon)
            text = raw.get("text", raw.get("label", raw.get("title", raw.get("heading", ""))))
            if text == "" and len(raw) == 1:
                # Handles parser output like {'text': '...'} or any single-key mapping fallback.
                text = next(iter(raw.values()))
            items.append({"icon": icon, "text": str(text), "url": get_url(raw)})
        else:
            items.append({"icon": default_icon, "text": str(raw), "url": ""})
    return items


def make_bullet_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [90, 160])

    # Merge style from PSL with bullet-level policy from PML.
    # PML should override PSL for behavioral flags such as overflow/pagination.
    bullet_cfg = value if isinstance(value, dict) else {}
    effective_style = dict(style)
    for key in (
        "overflow", "max_lines", "min_size", "auto_shrink",
        "max_items_per_slide", "items_per_slide",
        "icon", "bullet_icon", "icon_color", "icon_size", "icon_gap",
    ):
        if key in bullet_cfg:
            effective_style[key] = bullet_cfg[key]

    default_icon = resolve_bullet_icon(effective_style.get("icon", effective_style.get("bullet_icon")), "")
    icon_color = style_value(effective_style, "icon_color", theme, style_value(effective_style, "color", theme, "#000000"))
    return {
        "type": "BulletList",
        "items": normalize_bullet_items(value, default_icon),
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 420),
        "font": style_value(effective_style, "font", theme, "Aptos"),
        "size": effective_style.get("size", 24),
        "color": style_value(effective_style, "color", theme, "#000000"),
        "line_gap": effective_style.get("line_gap", 8),
        "icon_color": icon_color,
        "icon_size": effective_style.get("icon_size", effective_style.get("size", 24)),
        "icon_gap": effective_style.get("icon_gap", 12),
        "max_items_per_slide": safe_int(
            effective_style.get("max_items_per_slide")
            or effective_style.get("items_per_slide"),
            0,
        ),
        **text_style_props(effective_style),
    }


def normalize_fit_mode(value: Any, default: str = "stretch") -> str:
    """Normalize image placement modes.

    Modes:
      stretch  - resize image to exactly fill the box; aspect ratio may change.
      contain  - fit entire image inside the box; aspect ratio is preserved.
      cover    - fill the box and crop overflow; aspect ratio is preserved.
      crop     - alias of cover.
      original - use original pixel size, centered in the box when possible.
    """
    mode = str(value or default).strip().lower()
    aliases = {
        "resize": "stretch",
        "fill": "stretch",
        "fit": "contain",
        "scale": "contain",
        "crop": "cover",
        "none": "original",
    }
    mode = aliases.get(mode, mode)
    if mode not in {"stretch", "contain", "cover", "original"}:
        return default
    return mode


def make_image_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    """
    Supported PML forms:

      image: "path/to/image.png"

      image:
        src: "path/to/image.png"
        alt: "description"
        mode: contain   # stretch | contain | cover/crop | original

    The image style in PSL may also define `mode`. PML overrides PSL.
    """
    pos = style.get("position", [700, 150])

    if isinstance(value, dict):
        src = str(value.get("src", value.get("path", "")))
        alt = str(value.get("alt", ""))
        mode = normalize_fit_mode(value.get("mode", style.get("mode", "stretch")))
    else:
        src = str(value)
        alt = ""
        mode = normalize_fit_mode(style.get("mode", "stretch"))

    return {
        "type": "Image",
        "src": src,
        "alt": alt,
        "mode": mode,
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 500),
        "h": style.get("height", 320),
    }



def normalize_background_image(value: Any) -> Dict[str, Any]:
    """Normalize background image declarations.

    Supported forms:

      background_image: "bg.png"

      background_image:
        src: "bg.png"
        opacity: 0.25
        mode: cover

      background:
        image: "bg.png"
        opacity: 0.25
        mode: cover
    """
    if value is None or value == "":
        return {}
    if isinstance(value, str):
        return {"src": value}
    if isinstance(value, dict):
        if "image" in value and "src" not in value:
            return {
                "src": str(value.get("image", "")),
                "opacity": value.get("opacity", 1.0),
                "mode": normalize_fit_mode(value.get("mode", "cover"), default="cover"),
            }
        return {
            "src": str(value.get("src", value.get("path", ""))),
            "opacity": value.get("opacity", 1.0),
            "mode": normalize_fit_mode(value.get("mode", "cover"), default="cover"),
        }
    return {"src": str(value)}


def make_background_image_spec(blocks: Dict[str, Any], layout_style: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Resolve background image from PML first, then PSL layout defaults.

    Priority:
      1. slide/cover/header block: background_image
      2. slide/cover/header block: background.image
      3. layout style: background_image
      4. layout style: background.image
    """
    candidates: List[Any] = []

    if isinstance(blocks, dict):
        if "background_image" in blocks:
            candidates.append(blocks.get("background_image"))
        bg_block = blocks.get("background")
        if isinstance(bg_block, dict) and "image" in bg_block:
            candidates.append(bg_block)

    if isinstance(layout_style, dict):
        if "background_image" in layout_style:
            candidates.append(layout_style.get("background_image"))
        bg_style = layout_style.get("background")
        if isinstance(bg_style, dict) and "image" in bg_style:
            candidates.append(bg_style)

    for candidate in candidates:
        spec = normalize_background_image(candidate)
        if spec.get("src"):
            spec.setdefault("opacity", 1.0)
            spec["mode"] = normalize_fit_mode(spec.get("mode", "cover"), default="cover")
            return spec

    return None

def make_footer_image_objects(slide_blocks: Dict[str, Any], layout_style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """Build footer image objects for a slide.

    Supported PML forms:

      footer_image:
        src: "logo.png"
        alt: "Logo"

      footer_images:
        - src: "logo_a.png"
          alt: "Logo A"
        - src: "logo_b.png"
          alt: "Logo B"

    Style can be declared per slide layout:

      slide.title-bullets:
        footer_image:
          position: [1080, 660]
          width: 120
          height: 36

        footer_images:
          position: [900, 665]
          width: 80
          height: 30
          gap: 16
    """
    objects: List[Dict[str, Any]] = []

    if "footer_image" in slide_blocks:
        style = layout_style.get("footer_image", layout_style.get("footer", {}))
        default_style = {
            "position": style.get("position", [1080, 660]),
            "width": style.get("width", 120),
            "height": style.get("height", 36),
        }
        obj = make_image_object(slide_blocks["footer_image"], default_style, theme)
        obj["role"] = "footer_image"
        objects.append(obj)

    if "footer_images" in slide_blocks:
        images = slide_blocks["footer_images"]
        if not isinstance(images, list):
            images = [images]

        style = layout_style.get("footer_images", layout_style.get("footer", {}))
        pos = style.get("position", [900, 665])
        width = style.get("width", 80)
        height = style.get("height", 30)
        gap = style.get("gap", 16)

        for idx, image in enumerate(images):
            item_style = {
                "position": [pos[0] + idx * (width + gap), pos[1]],
                "width": width,
                "height": height,
            }
            obj = make_image_object(normalize_image_item(image), item_style, theme)
            obj["role"] = "footer_image"
            objects.append(obj)

    return objects



def make_footer_text_objects(slide_blocks: Dict[str, Any], layout_style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """Build footer text objects for a slide.

    Supported PML forms:

      footer_text:
        Nguồn: Viettel AI Lab | Internal Use Only

      footer_text: "Internal Use Only"

      footer:
        text: "Confidential"
        align: center

    Supported PSL forms:

      slide.title-bullets:
        footer_text:
          position: [80, 685]
          width: 900
          height: 24
          font: body
          size: 14
          color: muted

        footer:
          text_position: [80, 685]
          text_width: 900
          text_height: 24
          font: body
          size: 14
          color: muted
    """
    footer_text: Any = None
    footer_block = slide_blocks.get("footer")

    if "footer_text" in slide_blocks:
        footer_text = slide_blocks["footer_text"]
    elif isinstance(footer_block, dict) and "text" in footer_block:
        footer_text = footer_block.get("text")

    if footer_text is None or footer_text == "":
        return []

    text = normalize_text_block(footer_text)
    footer_style = layout_style.get("footer", {})
    style = dict(layout_style.get("footer_text", {}))

    # Allow a generic footer block in PSL to provide defaults.
    if "position" not in style:
        style["position"] = footer_style.get("text_position", footer_style.get("position", [80, 685]))
    if "width" not in style:
        style["width"] = footer_style.get("text_width", 900)
    if "height" not in style:
        style["height"] = footer_style.get("text_height", 24)
    if "font" not in style:
        style["font"] = footer_style.get("font", "body")
    if "size" not in style:
        style["size"] = footer_style.get("size", 14)
    if "color" not in style:
        style["color"] = footer_style.get("color", "muted")

    # Optional quick alignment override in PML. This is simple positioning on a 1280px canvas.
    align = None
    if isinstance(footer_block, dict):
        align = footer_block.get("align")
    if isinstance(align, str):
        width = style.get("width", 900)
        y = style.get("position", [80, 685])[1]
        if align.lower() == "center":
            style["position"] = [(1280 - width) / 2, y]
        elif align.lower() == "right":
            style["position"] = [1280 - width - 80, y]
        elif align.lower() == "left":
            style["position"] = [80, y]

    obj = make_text_object("footer_text", text, style, theme)
    obj["role"] = "footer_text"
    return [obj]


def make_column_objects(role: str, value: Any, style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """
    Build objects for a left/right column.

    Supported forms:

      left:
        heading: Main points
        bullets:
          - A
          - B

      right:
        image:
          src: "diagram.png"
          alt: "Diagram"

      left:
        text: Free paragraph text

    A column can combine heading + bullets + image, but the common use case is
    text in one column and image in the other column.
    """
    pos = style.get("position", [70, 150])
    width = style.get("width", 520)
    height = style.get("height", 360)
    objects: List[Dict[str, Any]] = []

    text_y = pos[1]

    if isinstance(value, dict):
        if "heading" in value:
            objects.append({
                "type": "TextBox",
                "role": f"{role}.heading",
                "text": str(value["heading"]),
                "x": pos[0],
                "y": text_y,
                "w": width,
                "h": style.get("heading_height", 50),
                "font": style_value(style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": style.get("heading_size", 24),
                "color": style_value(style, "heading_color", theme, theme.colors.get("primary", "#000000")),
                **text_style_props(style),
            })
            text_y += style.get("heading_height", 50) + style.get("heading_gap", 10)

        if "text" in value:
            objects.append({
                "type": "TextBox",
                "role": f"{role}.text",
                "text": str(value["text"]),
                "x": pos[0],
                "y": text_y,
                "w": width,
                "h": style.get("text_height", 220),
                "font": style_value(style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": style.get("size", 20),
                "color": style_value(style, "color", theme, theme.colors.get("text", "#000000")),
                **text_style_props(style),
            })
            text_y += style.get("text_height", 220) + style.get("text_gap", 12)

        if "bullets" in value:
            objects.append({
                "type": "BulletList",
                "items": [str(x) for x in value["bullets"]],
                "x": pos[0],
                "y": text_y,
                "w": width,
                "h": style.get("bullets_height", 300),
                "font": style_value(style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": style.get("size", 20),
                "color": style_value(style, "color", theme, theme.colors.get("text", "#000000")),
                "line_gap": style.get("line_gap", 6),
                **text_style_props(style),
            })

        if "image" in value:
            image_style = dict(style)
            image_style["position"] = style.get("image_position", pos)
            image_style["width"] = style.get("image_width", width)
            image_style["height"] = style.get("image_height", height)
            objects.append(make_image_object(value["image"], image_style, theme))

    else:
        objects.append({
            "type": "TextBox",
            "role": role,
            "text": str(value),
            "x": pos[0],
            "y": pos[1],
            "w": width,
            "h": height,
            "font": style_value(style, "font", theme, theme.fonts.get("body", "Aptos")),
            "size": style.get("size", 20),
            "color": style_value(style, "color", theme, theme.colors.get("text", "#000000")),
            **text_style_props(style),
        })

    return objects


def normalize_image_item(value: Any) -> Dict[str, Any]:
    """Normalize an image entry for the two-images layout."""
    if isinstance(value, dict):
        if "image" in value and isinstance(value["image"], dict):
            return value["image"]
        return value
    return {"src": str(value), "alt": ""}


def make_two_image_objects(slide_blocks: Dict[str, Any], style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """
    Build a native two-images layout.

    Supported PML forms:

      slide "Compare":
        layout: two-images
        images:
          - src: "before.png"
            alt: "Before"
          - src: "after.png"
            alt: "After"

      slide "Compare":
        layout: two-images
        left:
          image:
            src: "before.png"
        right:
          image:
            src: "after.png"
    """
    objects: List[Dict[str, Any]] = []

    if "images" in slide_blocks:
        images = slide_blocks["images"]
        if not isinstance(images, list):
            images = [images]
        images = images[:2]
        side_names = ["left", "right"]
        for idx, img in enumerate(images):
            side = side_names[idx]
            side_style = style.get(side, {})
            image_style = dict(side_style)
            image_style["position"] = side_style.get("image_position", side_style.get("position", [70 + idx * 610, 160]))
            image_style["width"] = side_style.get("image_width", side_style.get("width", 540))
            image_style["height"] = side_style.get("image_height", side_style.get("height", 380))
            objects.append(make_image_object(normalize_image_item(img), image_style, theme))
        return objects

    for side in ["left", "right"]:
        if side not in slide_blocks:
            continue
        value = slide_blocks[side]
        if isinstance(value, dict) and "image" in value:
            side_style = style.get(side, {})
            image_style = dict(side_style)
            image_style["position"] = side_style.get("image_position", side_style.get("position", [70, 160]))
            image_style["width"] = side_style.get("image_width", side_style.get("width", 540))
            image_style["height"] = side_style.get("image_height", side_style.get("height", 380))
            objects.append(make_image_object(value["image"], image_style, theme))

            if "caption" in value:
                cap_style = dict(side_style)
                cap_style["position"] = side_style.get("caption_position", [image_style["position"][0], image_style["position"][1] + image_style["height"] + 14])
                cap_style["width"] = side_style.get("caption_width", image_style["width"])
                cap_style["height"] = side_style.get("caption_height", 40)
                cap_style["size"] = side_style.get("caption_size", 16)
                cap_style["color"] = side_style.get("caption_color", theme.colors.get("muted", "#666666"))
                objects.append(make_text_object(f"{side}.caption", value["caption"], cap_style, theme))

    return objects


def normalize_cell(value: Any) -> Dict[str, Any]:
    """Normalize a grid cell to a dict with optional heading/text/bullets/image."""
    if isinstance(value, dict):
        if len(value) == 1 and "cell" in value and isinstance(value["cell"], dict):
            return value["cell"]
        return value
    return {"text": str(value)}


def make_grid_objects(value: Any, style: Dict[str, Any], theme: Theme, layout_name: str = "grid") -> List[Dict[str, Any]]:
    """
    Build card-like objects for grid layouts.

    Supported PML:
      cells:
        - heading: Thu thập
          text: Lấy dữ liệu từ nhiều nguồn
        - heading: Làm sạch
          bullets:
            - Chuẩn hóa
            - Khử trùng lặp
        - heading: Sơ đồ
          image:
            src: diagram.png
            alt: Diagram

    Supported PSL:
      slide.grid-6:
        grid:
          position: [70, 150]
          width: 1140
          height: 500
          columns: 3
          gap: 24
    """
    cells = value if isinstance(value, list) else [value]
    cells = [normalize_cell(cell) for cell in cells]

    grid_style = style.get("grid", style)
    pos = grid_style.get("position", [70, 150])
    total_w = grid_style.get("width", 1140)
    total_h = grid_style.get("height", 500)
    gap = grid_style.get("gap", 24)

    # Native grid layout presets. Explicit PSL values still override defaults.
    # grid-4x2 means 4 columns × 2 rows = up to 8 cells.
    grid_presets = {
        "grid-3": {"columns": 3, "rows": 1},
        "grid-4": {"columns": 2, "rows": 2},
        "grid-5": {"columns": 3, "rows": 2},
        "grid-6": {"columns": 3, "rows": 2},
        "grid-4x2": {"columns": 4, "rows": 2},
    }
    preset = grid_presets.get(layout_name, {"columns": 2, "rows": None})
    columns = int(grid_style.get("columns", preset["columns"]))
    columns = max(1, columns)

    explicit_rows = grid_style.get("rows", preset.get("rows"))
    if explicit_rows is None:
        rows = max(1, (len(cells) + columns - 1) // columns)
    else:
        rows = max(1, int(explicit_rows))

    cell_w = (total_w - gap * (columns - 1)) / columns
    cell_h = (total_h - gap * (rows - 1)) / rows

    objects: List[Dict[str, Any]] = []
    for idx, cell in enumerate(cells):
        row = idx // columns
        col = idx % columns
        x = pos[0] + col * (cell_w + gap)
        y = pos[1] + row * (cell_h + gap)

        objects.append({
            "type": "Card",
            "role": "grid.cell",
            "x": x,
            "y": y,
            "w": cell_w,
            "h": cell_h,
            "fill": style_value_from_item(cell, grid_style, "fill", theme, "#FFFFFF"),
            "border": border_value_from_item(cell, grid_style, theme, "#D1D5DB"),
            "radius": grid_style.get("radius", grid_style.get("border_radius", 16)),
        })

        inner_pad = grid_style.get("padding", 18)
        text_x = x + inner_pad
        text_y = y + inner_pad
        text_w = cell_w - 2 * inner_pad

        if "heading" in cell:
            objects.append({
                "type": "TextBox",
                "role": "grid.heading",
                "text": str(cell["heading"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": grid_style.get("heading_height", 34),
                "font": style_value(grid_style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": grid_style.get("heading_size", 20),
                "color": style_value(grid_style, "heading_color", theme, theme.colors.get("primary", "#000000")),
                **text_style_props(grid_style),
            })
            text_y += grid_style.get("heading_height", 34) + grid_style.get("heading_gap", 8)

        if "text" in cell:
            objects.append({
                "type": "TextBox",
                "role": "grid.text",
                "text": str(cell["text"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": grid_style.get("text_height", max(45, cell_h - (text_y - y) - inner_pad)),
                "font": style_value(grid_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": grid_style.get("size", 16),
                "color": style_value(grid_style, "color", theme, theme.colors.get("text", "#000000")),
                **text_style_props(grid_style),
            })
            text_y += grid_style.get("text_height", 70) + grid_style.get("text_gap", 8)

        if "bullets" in cell:
            bullets = cell["bullets"] if isinstance(cell["bullets"], list) else [cell["bullets"]]
            objects.append({
                "type": "BulletList",
                "items": [str(x) for x in bullets],
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": grid_style.get("bullets_height", max(45, cell_h - (text_y - y) - inner_pad)),
                "font": style_value(grid_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": grid_style.get("size", 15),
                "color": style_value(grid_style, "color", theme, theme.colors.get("text", "#000000")),
                "line_gap": grid_style.get("line_gap", 4),
                **text_style_props(grid_style),
            })
            text_y += grid_style.get("bullets_height", 80) + grid_style.get("text_gap", 8)

        if "image" in cell:
            image_style = dict(grid_style)
            image_style["position"] = [text_x, text_y]
            image_style["width"] = text_w
            image_style["height"] = max(50, cell_h - (text_y - y) - inner_pad)
            objects.append(make_image_object(cell["image"], image_style, theme))

    return objects




def normalize_numbered_column(value: Any) -> Dict[str, Any]:
    """Normalize a numbered-column item to a dict with heading/text/bullets."""
    if isinstance(value, dict):
        if len(value) == 1 and "column" in value and isinstance(value["column"], dict):
            return value["column"]
        return value
    return {"text": str(value)}


def make_numbered_columns_objects(slide_blocks: Dict[str, Any], style: Dict[str, Any], theme: Theme, layout_name: str) -> List[Dict[str, Any]]:
    """
    Build a 3–6 column layout where each card has a prominent circular number
    floating above it.

    Supported PML:

  
    slide "Danh sách dài tự ngắt":
      layout: title-bullets

      title:
        Danh sách dài tự ngắt sang slide sau

      bullets:
        icon: check
        overflow: paginate
          items:
          - Thu thập dữ liệu từ nhiều nguồn nội bộ
          - Chuẩn hóa định dạng đầu vào
          - Loại bỏ bản ghi trùng lặp
          - Gắn metadata cho từng tài liệu
          - Tạo embedding và chỉ mục truy xuất
          - Mở rộng ngữ cảnh bằng đồ thị quan hệ
          - Sinh tóm tắt theo chuẩn báo cáo
          - Kiểm tra chất lượng đầu ra
          - Xuất dashboard và slide trình bày
          - Lưu log phục vụ kiểm toán

    slide "Quy trình 4 bước":
        layout: numbered-columns-4
        title:
          Quy trình xử lý
        columns:
          - heading: Thu thập
            text: Lấy dữ liệu từ nhiều nguồn.
          - heading: Làm sạch
            bullets:
              - Chuẩn hóa
              - Khử trùng lặp

    Supported layouts:
      numbered-columns-3, numbered-columns-4, numbered-columns-5, numbered-columns-6

    Supported PSL:
      slide.numbered-columns-4:
        numbered_columns:
          position: [70, 170]
          width: 1140
          height: 380
          gap: 24
          circle_size: 58
          circle_offset_y: -34
          fill: "#FFFFFF"
          border: "#CBD5E1"
          circle_fill: primary
          circle_color: white
    """
    raw_items = slide_blocks.get("columns", slide_blocks.get("cells", slide_blocks.get("items", [])))
    if not isinstance(raw_items, list):
        raw_items = [raw_items]
    items = [normalize_numbered_column(item) for item in raw_items]

    numbered_style = style.get("numbered_columns", style.get("columns", style))
    pos = numbered_style.get("position", [70, 170])
    total_w = float(numbered_style.get("width", 1140))
    total_h = float(numbered_style.get("height", 380))
    gap = float(numbered_style.get("gap", 24))

    m = re.search(r"(\d+)$", layout_name)
    default_cols = int(m.group(1)) if m else len(items)
    columns = int(numbered_style.get("columns", default_cols))
    columns = max(3, min(6, columns))
    items = items[:columns]

    card_w = (total_w - gap * (columns - 1)) / columns
    card_h = total_h
    circle_size = float(numbered_style.get("circle_size", 58))
    circle_offset_y = float(numbered_style.get("circle_offset_y", -34))
    pad = float(numbered_style.get("padding", 18))

    objects: List[Dict[str, Any]] = []
    for idx, item in enumerate(items):
        x = float(pos[0]) + idx * (card_w + gap)
        y = float(pos[1])
        cx = x + (card_w - circle_size) / 2
        cy = y + circle_offset_y

        objects.append({
            "type": "Card",
            "role": "numbered-column.card",
            "x": x,
            "y": y,
            "w": card_w,
            "h": card_h,
            "fill": style_value_from_item(item, numbered_style, "fill", theme, "#FFFFFF"),
            "border": border_value_from_item(item, numbered_style, theme, "#D1D5DB"),
            "radius": numbered_style.get("radius", numbered_style.get("border_radius", 18)),
        })

        objects.append({
            "type": "Circle",
            "role": "numbered-column.circle",
            "x": cx,
            "y": cy,
            "w": circle_size,
            "h": circle_size,
            "fill": style_value_from_item(item, numbered_style, "circle_fill", theme, theme.colors.get("primary", "#2563EB")),
            "border": style_value(numbered_style, "circle_border", theme, "#FFFFFF"),
        })

        objects.append({
            "type": "TextBox",
            "role": "numbered-column.number",
            "text": str(item.get("number", idx + 1)),
            "x": cx,
            "y": cy + float(numbered_style.get("number_y_adjust", 7)),
            "w": circle_size,
            "h": circle_size,
            "font": style_value(numbered_style, "number_font", theme, theme.fonts.get("heading", "Aptos Display")),
            "size": numbered_style.get("number_size", 24),
            "color": style_value(numbered_style, "circle_color", theme, "#FFFFFF"),
            "bold": True,
            "italic": False,
            "underline": False,
            "align": "center",
            "overflow": "shrink",
            "min_size": 14,
        })

        text_x = x + pad
        text_y = y + pad + float(numbered_style.get("top_padding_extra", 22))
        text_w = card_w - 2 * pad

        if "heading" in item:
            objects.append({
                "type": "TextBox",
                "role": "numbered-column.heading",
                "text": str(item["heading"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": numbered_style.get("heading_height", 46),
                "font": style_value(numbered_style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": numbered_style.get("heading_size", 20),
                "color": style_value(numbered_style, "heading_color", theme, theme.colors.get("primary", "#2563EB")),
                "bold": parse_bool(numbered_style.get("heading_bold", True)),
                "italic": parse_bool(numbered_style.get("heading_italic", False)),
                "underline": parse_bool(numbered_style.get("heading_underline", False)),
                "align": normalize_align(numbered_style.get("heading_align", "center")),
                "overflow": numbered_style.get("heading_overflow", numbered_style.get("overflow", "shrink")),
                "max_lines": numbered_style.get("heading_max_lines", 2),
                "min_size": numbered_style.get("min_size", 13),
            })
            text_y += float(numbered_style.get("heading_height", 46)) + float(numbered_style.get("heading_gap", 8))

        if "text" in item:
            objects.append({
                "type": "TextBox",
                "role": "numbered-column.text",
                "text": str(item["text"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": numbered_style.get("text_height", max(60, card_h - (text_y - y) - pad)),
                "font": style_value(numbered_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": numbered_style.get("size", 15),
                "color": style_value(numbered_style, "color", theme, theme.colors.get("text", "#111827")),
                "bold": parse_bool(numbered_style.get("bold", False)),
                "italic": parse_bool(numbered_style.get("italic", False)),
                "underline": parse_bool(numbered_style.get("underline", False)),
                "align": normalize_align(numbered_style.get("align", "center")),
                "overflow": numbered_style.get("overflow", "shrink"),
                "max_lines": numbered_style.get("max_lines", 5),
                "min_size": numbered_style.get("min_size", 12),
            })
            text_y += float(numbered_style.get("text_height", 95)) + float(numbered_style.get("text_gap", 8))

        if "bullets" in item:
            # Numbered-column cards support both simple bullets and icon bullets.
            #
            # Simple form:
            #   bullets:
            #     - Normalize
            #     - Validate
            #
            # Shared icon form:
            #   bullets:
            #     icon: check
            #     items:
            #       - Normalize
            #       - Validate
            #
            # Per-item icon form:
            #   bullets:
            #     - icon: gear
            #       text: Extract
            #     - icon: star
            #       text: Summarize
            bullets = item["bullets"]
            bullet_style = dict(numbered_style)
            bullet_style["position"] = [text_x, text_y]
            bullet_style["width"] = text_w
            bullet_style["height"] = max(55, card_h - (text_y - y) - pad)
            bullet_style["size"] = numbered_style.get("bullet_size", numbered_style.get("size", 14))
            bullet_style["align"] = numbered_style.get("bullet_align", "left")
            bullet_style["line_gap"] = numbered_style.get("bullet_line_gap", numbered_style.get("line_gap", 4))

            # Optional card-level override:
            #   - heading: ...
            #     bullet_icon: check
            #     bullets:
            #       - ...
            if "bullet_icon" in item:
                bullet_style["icon"] = item["bullet_icon"]
            elif "icon" in item and not isinstance(bullets, dict):
                bullet_style["icon"] = item["icon"]

            objects.append(make_bullet_object(bullets, bullet_style, theme))

    return objects

def normalize_step(value: Any) -> Dict[str, Any]:
    """Normalize a stair-progress step to a dict with heading/text/bullets/icon."""
    if isinstance(value, dict):
        if len(value) == 1 and "step" in value and isinstance(value["step"], dict):
            return value["step"]
        return value
    return {"text": str(value)}


def make_stair_progress_objects(slide_blocks: Dict[str, Any], style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """
    Build a stair-step progress layout.

    Supported PML:

      slide "Roadmap":
        layout: stair-progress
        title:
          Lộ trình triển khai
        steps:
          - heading: Khởi tạo
            text: Xác định phạm vi và nguồn dữ liệu.
          - heading: Chuẩn hóa
            text: Làm sạch và hợp nhất dữ liệu.
          - heading: Tự động hóa
            bullets:
              - Xây pipeline
              - Theo dõi lỗi

    Supported PSL:

      slide.stair-progress:
        stair:
          position: [80, 170]
          width: 1080
          height: 420
          step_width: 245
          step_height: 120
          x_step: 205
          y_step: 55
          fill: secondary
          border: primary
    """
    raw_steps = slide_blocks.get("steps", slide_blocks.get("items", []))
    if not isinstance(raw_steps, list):
        raw_steps = [raw_steps]
    steps = [normalize_step(step) for step in raw_steps]

    stair_style = style.get("stair", style)
    pos = stair_style.get("position", [80, 170])
    max_steps = int(stair_style.get("max_steps", 6))
    steps = steps[:max_steps]

    step_w = float(stair_style.get("step_width", 250))
    step_h = float(stair_style.get("step_height", 118))
    x_step = float(stair_style.get("x_step", 190))
    y_step = float(stair_style.get("y_step", 55))
    pad = float(stair_style.get("padding", 16))

    objects: List[Dict[str, Any]] = []
    for idx, step in enumerate(steps):
        x = float(pos[0]) + idx * x_step
        y = float(pos[1]) + idx * y_step

        # Card background for each stair.
        objects.append({
            "type": "Card",
            "role": "stair.step",
            "x": x,
            "y": y,
            "w": step_w,
            "h": step_h,
            "fill": style_value_from_item(step, stair_style, "fill", theme, "#FFFFFF"),
            "border": border_value_from_item(step, stair_style, theme, theme.colors.get("primary", "#2563EB")),
            "radius": stair_style.get("radius", stair_style.get("border_radius", 16)),
        })

        # Step number.
        number_text = str(step.get("number", idx + 1))
        objects.append({
            "type": "TextBox",
            "role": "stair.number",
            "text": number_text,
            "x": x + pad,
            "y": y + pad,
            "w": float(stair_style.get("number_width", 34)),
            "h": float(stair_style.get("number_height", 30)),
            "font": style_value(stair_style, "number_font", theme, theme.fonts.get("heading", "Aptos Display")),
            "size": stair_style.get("number_size", 20),
            "color": style_value(stair_style, "number_color", theme, theme.colors.get("primary", "#2563EB")),
            "bold": True,
            "italic": False,
            "underline": False,
            "align": "center",
        })

        text_x = x + pad + float(stair_style.get("number_width", 34)) + float(stair_style.get("number_gap", 8))
        text_w = step_w - (text_x - x) - pad
        text_y = y + pad

        if "heading" in step:
            objects.append({
                "type": "TextBox",
                "role": "stair.heading",
                "text": str(step["heading"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": stair_style.get("heading_height", 32),
                "font": style_value(stair_style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": stair_style.get("heading_size", 18),
                "color": style_value(stair_style, "heading_color", theme, theme.colors.get("text", "#111827")),
                "bold": parse_bool(stair_style.get("heading_bold", True)),
                "italic": parse_bool(stair_style.get("heading_italic", False)),
                "underline": parse_bool(stair_style.get("heading_underline", False)),
                "align": normalize_align(stair_style.get("heading_align", "left")),
            })
            text_y += float(stair_style.get("heading_height", 32)) + float(stair_style.get("heading_gap", 4))

        if "text" in step:
            objects.append({
                "type": "TextBox",
                "role": "stair.text",
                "text": str(step["text"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": stair_style.get("text_height", step_h - (text_y - y) - pad),
                "font": style_value(stair_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": stair_style.get("size", 14),
                "color": style_value(stair_style, "color", theme, theme.colors.get("muted", "#6B7280")),
                **text_style_props(stair_style),
            })
            text_y += float(stair_style.get("text_height", 48)) + float(stair_style.get("text_gap", 4))

        if "bullets" in step:
            bullets = step["bullets"] if isinstance(step["bullets"], list) else [step["bullets"]]
            objects.append({
                "type": "BulletList",
                "role": "stair.bullets",
                "items": [str(x) for x in bullets],
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": stair_style.get("bullets_height", step_h - (text_y - y) - pad),
                "font": style_value(stair_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": stair_style.get("size", 13),
                "color": style_value(stair_style, "color", theme, theme.colors.get("muted", "#6B7280")),
                "line_gap": stair_style.get("line_gap", 3),
                **text_style_props(stair_style),
            })

    return objects


def make_stacked_stairs_objects(slide_blocks: Dict[str, Any], style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """
    Build a stacked stair layout where steps are placed on top of each other,
    shrink gradually, and align to one side.

    Supported PML:

      slide "Mô hình trưởng thành":
        layout: stacked-stairs
        title:
          Các tầng năng lực
        steps:
          - heading: Nền tảng dữ liệu
            text: Thu thập và chuẩn hóa dữ liệu.
          - heading: Truy xuất tri thức
            text: Xây dựng index và graph context.

    Supported PSL:

      slide.stacked-stairs:
        stacked_stairs:
          position: [120, 170]
          base_width: 980
          step_height: 86
          shrink: 90
          overlap: 24
          align_side: left   # left | right
          fill: secondary
          border: primary
    """
    raw_steps = slide_blocks.get("steps", slide_blocks.get("items", []))
    if not isinstance(raw_steps, list):
        raw_steps = [raw_steps]
    steps = [normalize_step(step) for step in raw_steps]

    stair_style = style.get("stacked_stairs", style.get("stacked-stairs", style.get("stairs", style)))
    pos = stair_style.get("position", [120, 170])
    max_steps = int(stair_style.get("max_steps", 6))
    steps = steps[:max_steps]

    x0 = float(pos[0])
    y0 = float(pos[1])
    base_w = float(stair_style.get("base_width", stair_style.get("width", 980)))
    step_h = float(stair_style.get("step_height", 86))
    shrink = float(stair_style.get("shrink", 90))
    overlap = float(stair_style.get("overlap", 22))
    pad_x = float(stair_style.get("padding_x", stair_style.get("padding", 22)))
    pad_y = float(stair_style.get("padding_y", 14))
    align_side = str(stair_style.get("align_side", stair_style.get("align", "left"))).lower()
    if align_side not in {"left", "right"}:
        align_side = "left"

    objects: List[Dict[str, Any]] = []
    for idx, step in enumerate(steps):
        w = max(float(stair_style.get("min_width", 360)), base_w - idx * shrink)
        h = step_h
        x = x0 if align_side == "left" else x0 + (base_w - w)
        y = y0 + idx * (step_h - overlap)

        objects.append({
            "type": "Card",
            "role": "stacked_stair.step",
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "fill": style_value_from_item(step, stair_style, "fill", theme, theme.colors.get("secondary", "#E6F0FF")),
            "border": border_value_from_item(step, stair_style, theme, theme.colors.get("primary", "#2563EB")),
            "radius": stair_style.get("radius", stair_style.get("border_radius", 16)),
        })

        number_w = float(stair_style.get("number_width", 42))
        number_text = str(step.get("number", idx + 1))
        objects.append({
            "type": "TextBox",
            "role": "stacked_stair.number",
            "text": number_text,
            "x": x + pad_x,
            "y": y + pad_y,
            "w": number_w,
            "h": h - 2 * pad_y,
            "font": style_value(stair_style, "number_font", theme, theme.fonts.get("heading", "Aptos Display")),
            "size": stair_style.get("number_size", 22),
            "color": style_value(stair_style, "number_color", theme, theme.colors.get("primary", "#2563EB")),
            "bold": True,
            "italic": False,
            "underline": False,
            "align": "center",
        })

        text_x = x + pad_x + number_w + float(stair_style.get("number_gap", 14))
        text_w = max(100, w - (text_x - x) - pad_x)
        text_y = y + pad_y

        if "heading" in step:
            objects.append({
                "type": "TextBox",
                "role": "stacked_stair.heading",
                "text": str(step["heading"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": stair_style.get("heading_height", 30),
                "font": style_value(stair_style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": stair_style.get("heading_size", 18),
                "color": style_value(stair_style, "heading_color", theme, theme.colors.get("text", "#111827")),
                "bold": parse_bool(stair_style.get("heading_bold", True)),
                "italic": parse_bool(stair_style.get("heading_italic", False)),
                "underline": parse_bool(stair_style.get("heading_underline", False)),
                "align": normalize_align(stair_style.get("heading_align", "left")),
            })
            text_y += float(stair_style.get("heading_height", 30)) + float(stair_style.get("heading_gap", 2))

        if "text" in step:
            objects.append({
                "type": "TextBox",
                "role": "stacked_stair.text",
                "text": str(step["text"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": max(24, h - (text_y - y) - pad_y),
                "font": style_value(stair_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": stair_style.get("size", 14),
                "color": style_value(stair_style, "color", theme, theme.colors.get("muted", "#6B7280")),
                **text_style_props(stair_style),
            })

        if "bullets" in step:
            bullets = step["bullets"] if isinstance(step["bullets"], list) else [step["bullets"]]
            objects.append({
                "type": "BulletList",
                "role": "stacked_stair.bullets",
                "items": [str(x) for x in bullets],
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": max(24, h - (text_y - y) - pad_y),
                "font": style_value(stair_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": stair_style.get("size", 13),
                "color": style_value(stair_style, "color", theme, theme.colors.get("muted", "#6B7280")),
                "line_gap": stair_style.get("line_gap", 3),
                **text_style_props(stair_style),
            })

    return objects



def normalize_milestone(item: Any) -> Dict[str, Any]:
    """Normalize a timeline milestone to a dict.

    Supported item fields:
      - date / label / heading / text / bullets
      - number is optional; if omitted, the item index is used.
    """
    if isinstance(item, dict):
        return dict(item)
    return {"heading": str(item)}


def make_linebox(x: float, y: float, w: float, h: float, fill: Any, role: str = "timeline.line") -> Dict[str, Any]:
    return {
        "type": "LineBox",
        "role": role,
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "fill": fill,
    }


def make_circle(x: float, y: float, size: float, fill: Any, border: Any, role: str = "timeline.marker") -> Dict[str, Any]:
    return {
        "type": "Circle",
        "role": role,
        "x": x,
        "y": y,
        "w": size,
        "h": size,
        "fill": fill,
        "border": border,
    }


def make_timeline_objects(slide_blocks: Dict[str, Any], style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """
    Build a horizontal timeline layout.

    Supported PML:

      slide "Lộ trình triển khai":
        layout: timeline
        title:
          Lộ trình triển khai
        milestones:
          - date: Q1
            heading: Khởi tạo
            text: Xác định phạm vi và dữ liệu nguồn.
          - date: Q2
            heading: Chuẩn hóa
            text: Làm sạch và chuẩn hóa dữ liệu.

    Also accepts `steps:` or `items:` instead of `milestones:`.

    Supported PSL:

      slide.timeline:
        timeline:
          position: [90, 155]
          width: 1100
          height: 470
          axis_y: 360
          line_height: 5
          marker_size: 22
          connector_height: 58
          card_width: 210
          card_height: 116
          alternate: true
          fill: "#FFFFFF"
          border: primary
          line_color: primary
          marker_fill: primary
          marker_border: "#FFFFFF"
    """
    raw = slide_blocks.get("milestones", slide_blocks.get("steps", slide_blocks.get("items", [])))
    if not isinstance(raw, list):
        raw = [raw]
    milestones = [normalize_milestone(item) for item in raw]

    timeline_style = style.get("timeline", style)
    max_items = int(timeline_style.get("max_items", timeline_style.get("max_steps", 7)))
    milestones = milestones[:max_items]
    if not milestones:
        return []

    pos = timeline_style.get("position", [90, 155])
    x0 = float(pos[0])
    y0 = float(pos[1])
    total_w = float(timeline_style.get("width", 1100))
    total_h = float(timeline_style.get("height", 470))
    axis_y = float(timeline_style.get("axis_y", y0 + total_h / 2))
    line_h = float(timeline_style.get("line_height", 5))
    marker_size = float(timeline_style.get("marker_size", 22))
    connector_h = float(timeline_style.get("connector_height", 58))
    card_w = float(timeline_style.get("card_width", 210))
    card_h = float(timeline_style.get("card_height", 116))
    card_pad = float(timeline_style.get("padding", 14))
    gap = float(timeline_style.get("gap", 24))
    alternate = parse_bool(timeline_style.get("alternate", True))
    start_side = str(timeline_style.get("start_side", "top")).lower()

    line_color = style_value(timeline_style, "line_color", theme, theme.colors.get("primary", "#2563EB"))
    marker_fill = style_value(timeline_style, "marker_fill", theme, theme.colors.get("primary", "#2563EB"))
    marker_border = style_value(timeline_style, "marker_border", theme, "#FFFFFF")
    card_fill = style_value(timeline_style, "fill", theme, "#FFFFFF")
    card_border = style_value(timeline_style, "border", theme, theme.colors.get("primary", "#2563EB"))

    objects: List[Dict[str, Any]] = []

    # Main horizontal axis.
    objects.append(make_linebox(x0, axis_y - line_h / 2, total_w, line_h, line_color, "timeline.axis"))

    if len(milestones) == 1:
        xs = [x0 + total_w / 2]
    else:
        inset = float(timeline_style.get("inset", 20))
        usable_w = max(1.0, total_w - 2 * inset)
        xs = [x0 + inset + i * usable_w / (len(milestones) - 1) for i in range(len(milestones))]

    for idx, item in enumerate(milestones):
        center_x = xs[idx]
        side_top = (idx % 2 == 0) if alternate else (start_side != "bottom")
        if start_side == "bottom" and alternate:
            side_top = not side_top

        marker_x = center_x - marker_size / 2
        marker_y = axis_y - marker_size / 2

        # Keep cards inside the timeline bounding box when possible.
        card_x = max(x0, min(center_x - card_w / 2, x0 + total_w - card_w))
        if side_top:
            card_y = axis_y - connector_h - card_h
            connector_y = card_y + card_h
            connector_h_actual = axis_y - connector_y - marker_size / 2
        else:
            connector_y = axis_y + marker_size / 2
            connector_h_actual = connector_h
            card_y = connector_y + connector_h_actual

        # Clamp vertically but keep connector visually reasonable.
        card_y = max(y0, min(card_y, y0 + total_h - card_h))
        if side_top:
            connector_y = card_y + card_h
            connector_h_actual = max(8, axis_y - marker_size / 2 - connector_y)
        else:
            connector_y = axis_y + marker_size / 2
            connector_h_actual = max(8, card_y - connector_y)

        # Connector line.
        objects.append(make_linebox(center_x - line_h / 2, connector_y, line_h, connector_h_actual, line_color, "timeline.connector"))

        # Milestone marker.
        objects.append(make_circle(marker_x, marker_y, marker_size, marker_fill, marker_border, "timeline.marker"))

        # Card body.
        objects.append({
            "type": "Card",
            "role": "timeline.card",
            "x": card_x,
            "y": card_y,
            "w": card_w,
            "h": card_h,
            "fill": style_value_from_item(item, timeline_style, "fill", theme, card_fill),
            "border": border_value_from_item(item, timeline_style, theme, card_border),
            "radius": timeline_style.get("radius", timeline_style.get("border_radius", 16)),
        })

        text_x = card_x + card_pad
        text_y = card_y + card_pad
        text_w = card_w - 2 * card_pad

        date_text = item.get("date", item.get("label", item.get("number", idx + 1)))
        if date_text is not None and str(date_text).strip():
            objects.append({
                "type": "TextBox",
                "role": "timeline.date",
                "text": str(date_text),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": timeline_style.get("date_height", 22),
                "font": style_value(timeline_style, "date_font", theme, theme.fonts.get("body", "Aptos")),
                "size": timeline_style.get("date_size", 13),
                "color": style_value(timeline_style, "date_color", theme, theme.colors.get("primary", "#2563EB")),
                "bold": parse_bool(timeline_style.get("date_bold", True)),
                "italic": parse_bool(timeline_style.get("date_italic", False)),
                "underline": parse_bool(timeline_style.get("date_underline", False)),
                "align": normalize_align(timeline_style.get("date_align", "left")),
            })
            text_y += float(timeline_style.get("date_height", 22)) + float(timeline_style.get("date_gap", 3))

        if "heading" in item:
            objects.append({
                "type": "TextBox",
                "role": "timeline.heading",
                "text": str(item["heading"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": timeline_style.get("heading_height", 28),
                "font": style_value(timeline_style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": timeline_style.get("heading_size", 17),
                "color": style_value(timeline_style, "heading_color", theme, theme.colors.get("text", "#111827")),
                "bold": parse_bool(timeline_style.get("heading_bold", True)),
                "italic": parse_bool(timeline_style.get("heading_italic", False)),
                "underline": parse_bool(timeline_style.get("heading_underline", False)),
                "align": normalize_align(timeline_style.get("heading_align", "left")),
            })
            text_y += float(timeline_style.get("heading_height", 28)) + float(timeline_style.get("heading_gap", 4))

        if "text" in item:
            objects.append({
                "type": "TextBox",
                "role": "timeline.text",
                "text": str(item["text"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": max(26, card_y + card_h - card_pad - text_y),
                "font": style_value(timeline_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": timeline_style.get("size", 13),
                "color": style_value(timeline_style, "color", theme, theme.colors.get("muted", "#6B7280")),
                **text_style_props(timeline_style),
            })

        if "bullets" in item:
            bullets = item["bullets"] if isinstance(item["bullets"], list) else [item["bullets"]]
            objects.append({
                "type": "BulletList",
                "role": "timeline.bullets",
                "items": [str(x) for x in bullets],
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": max(26, card_y + card_h - card_pad - text_y),
                "font": style_value(timeline_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": timeline_style.get("size", 12),
                "color": style_value(timeline_style, "color", theme, theme.colors.get("muted", "#6B7280")),
                "line_gap": timeline_style.get("line_gap", 3),
                **text_style_props(timeline_style),
            })

    return objects

def build_text_objects_from_block(
    block: Dict[str, Any],
    layout_style: Dict[str, Any],
    theme: Theme,
    allowed_keys: Optional[List[str]] = None,
) -> List[Dict[str, Any]]:
    objects: List[Dict[str, Any]] = []
    keys = allowed_keys or list(block.keys())

    for key in keys:
        if key not in block:
            continue
        style = layout_style.get(key, {})
        objects.append(make_text_object(key, block[key], style, theme))

    return objects



# -----------------------------------------------------------------------------
# Auto pagination helpers
# -----------------------------------------------------------------------------

def clone_obj(obj: Dict[str, Any]) -> Dict[str, Any]:
    """Shallow-copy an IR object; list fields used by paginated objects are copied."""
    out = dict(obj)
    if "items" in out and isinstance(out["items"], list):
        out["items"] = [dict(x) if isinstance(x, dict) else x for x in out["items"]]
    if "rows" in out and isinstance(out["rows"], list):
        out["rows"] = [list(r) if isinstance(r, list) else r for r in out["rows"]]
    if "headers" in out and isinstance(out["headers"], list):
        out["headers"] = list(out["headers"])
    return out


def chunk_list(items: List[Any], chunk_size: int) -> List[List[Any]]:
    if chunk_size <= 0:
        return [items]
    return [items[i:i + chunk_size] for i in range(0, len(items), chunk_size)] or [[]]


def approx_chars_per_line(width_px: float, font_size_px: float, factor: float = 0.55) -> int:
    """Approximate how many average Latin/Vietnamese characters fit in a line.

    This is intentionally conservative. python-pptx cannot measure real rendered text,
    so pagination uses a predictable heuristic instead of waiting for PowerPoint to
    clip overflowing text.
    """
    font_size_px = max(1.0, float(font_size_px or 12))
    usable_width = max(10.0, float(width_px or 100))
    return max(1, int(usable_width / (font_size_px * factor)))


def estimate_wrapped_lines(text: Any, width_px: float, font_size_px: float) -> int:
    """Estimate wrapped line count for a text string."""
    s = str(text or "")
    if not s:
        return 1
    chars_per_line = approx_chars_per_line(width_px, font_size_px)
    total = 0
    for raw_line in s.splitlines() or [""]:
        line = raw_line.strip()
        if not line:
            total += 1
            continue
        words = line.split()
        current = 0
        used = 1
        for word in words:
            w = len(word)
            if current == 0:
                current = w
            elif current + 1 + w <= chars_per_line:
                current += 1 + w
            else:
                used += 1
                current = w
        total += max(1, used)
    return max(1, total)


def estimate_text_height(text: Any, width_px: float, font_size_px: float, line_gap: float = 0) -> float:
    lines = estimate_wrapped_lines(text, width_px, font_size_px)
    return lines * (float(font_size_px or 12) * 1.22 + float(line_gap or 0))


def split_text_to_height(text: Any, width_px: float, font_size_px: float, max_height: float, line_gap: float = 0) -> List[str]:
    """Split a single very long text into chunks that fit approximately in max_height."""
    s = str(text or "")
    if not s.strip():
        return [s]
    line_height = float(font_size_px or 12) * 1.22 + float(line_gap or 0)
    max_lines = max(1, int(float(max_height or line_height) // max(1.0, line_height)))
    chars_per_line = approx_chars_per_line(width_px, font_size_px)
    max_chars = max(20, chars_per_line * max_lines)
    words = s.split()
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0
    for word in words:
        add_len = len(word) + (1 if current else 0)
        if current and current_len + add_len > max_chars:
            chunks.append(" ".join(current))
            current = [word]
            current_len = len(word)
        else:
            current.append(word)
            current_len += add_len
    if current:
        chunks.append(" ".join(current))
    return chunks or [s]


def split_bullet_items_by_height(obj: Dict[str, Any]) -> List[List[Dict[str, Any]]]:
    """Paginate bullet items by estimated rendered height, not just item count.

    Respects max_items_per_slide as an additional hard cap if provided. If a single
    bullet is too tall, its text is split into continuation bullet items.
    """
    items = obj.get("items", []) or []
    if not items:
        return [[]]
    x = safe_float(obj.get("x", 0), 0)
    w = safe_float(obj.get("w", 1000), 1000)
    h = safe_float(obj.get("h", 420), 420)
    size = safe_float(obj.get("size", 24), 24)
    icon_size = safe_float(obj.get("icon_size", size), size)
    icon_gap = safe_float(obj.get("icon_gap", 12), 12)
    line_gap = safe_float(obj.get("line_gap", 8), 8)
    max_items = safe_int(obj.get("max_items_per_slide", 0), 0)

    # Leave a little vertical breathing room to avoid borderline PPTX clipping.
    available_h = max(20.0, h - 8)
    text_w = max(20.0, w - icon_size - icon_gap - 4)
    pages: List[List[Dict[str, Any]]] = []
    current: List[Dict[str, Any]] = []
    used_h = 0.0

    def item_height(item: Dict[str, Any]) -> float:
        text = item.get("text", item) if isinstance(item, dict) else item
        return estimate_text_height(text, text_w, size, line_gap) + max(4.0, line_gap)

    for raw_item in items:
        item = dict(raw_item) if isinstance(raw_item, dict) else {"icon": "", "text": str(raw_item), "url": ""}
        ih = item_height(item)

        # If one item alone overflows the whole box, split its text by height.
        if ih > available_h * 1.05:
            text = str(item.get("text", ""))
            chunks = split_text_to_height(text, text_w, size, available_h * 0.92, line_gap)
            split_items = []
            for part_index, part in enumerate(chunks):
                new_item = dict(item)
                new_item["text"] = part
                if part_index > 0 and new_item.get("icon"):
                    # Avoid visually repeating an icon for continuation chunks unless desired.
                    new_item["icon"] = ""
                split_items.append(new_item)
        else:
            split_items = [item]

        for item2 in split_items:
            ih2 = item_height(item2)
            hit_item_cap = max_items > 0 and len(current) >= max_items
            hit_height_cap = current and (used_h + ih2 > available_h)
            if hit_item_cap or hit_height_cap:
                pages.append(current)
                current = []
                used_h = 0.0
            current.append(item2)
            used_h += ih2
    if current:
        pages.append(current)
    return pages or [[]]


def split_table_rows_by_height(obj: Dict[str, Any]) -> List[List[List[str]]]:
    """Paginate table rows by estimated row height.

    Header height is counted on every continuation slide because headers repeat.
    max_rows_per_slide is treated as an extra hard cap if present.
    """
    rows = obj.get("rows", []) or []
    if not rows:
        return [[]]
    headers = obj.get("headers", []) or []
    w = safe_float(obj.get("w", 1140), 1140)
    h = safe_float(obj.get("h", 430), 430)
    size = safe_float(obj.get("size", 18), 18)
    max_rows = safe_int(obj.get("max_rows_per_slide", 0), 0)
    col_count = max([len(headers)] + [len(r) for r in rows] + [1])
    cell_w = max(30.0, w / col_count - 14)
    base_line_h = size * 1.20 + 7
    header_h = 0.0
    if headers:
        header_h = max(base_line_h, max(estimate_text_height(c, cell_w, size, 2) for c in headers) + 10)
    available_h = max(30.0, h - header_h - 8)

    def row_height(row: List[Any]) -> float:
        if not isinstance(row, list):
            row = [row]
        return max(base_line_h, max(estimate_text_height(c, cell_w, size, 2) for c in row) + 10)

    pages: List[List[List[str]]] = []
    current: List[List[str]] = []
    used_h = 0.0
    for raw_row in rows:
        row = list(raw_row) if isinstance(raw_row, list) else [str(raw_row)]
        rh = row_height(row)
        hit_row_cap = max_rows > 0 and len(current) >= max_rows
        hit_height_cap = current and (used_h + rh > available_h)
        if hit_row_cap or hit_height_cap:
            pages.append(current)
            current = []
            used_h = 0.0
        current.append(row)
        used_h += rh
    if current:
        pages.append(current)
    return pages or [[]]


def mark_continued(slide_ir: Dict[str, Any], page_index: int) -> None:
    """Append a continuation marker to the first title-like text box."""
    if page_index <= 0:
        return
    suffix = f" (cont. {page_index + 1})"
    for obj in slide_ir.get("objects", []):
        if obj.get("type") == "TextBox" and str(obj.get("role", "")).lower() in {"title", "slide.title"}:
            text = str(obj.get("text", ""))
            if "(cont." not in text:
                obj["text"] = text + suffix
            return
    slide_ir["slide_title"] = str(slide_ir.get("slide_title", "")) + suffix



# -----------------------------------------------------------------------------
# PCL constraint application
# -----------------------------------------------------------------------------

def deep_merge_dict(base: Dict[str, Any], override: Dict[str, Any]) -> Dict[str, Any]:
    """Return a recursive merge without mutating inputs."""
    result = dict(base or {})
    for key, value in (override or {}).items():
        if isinstance(value, dict) and isinstance(result.get(key), dict):
            result[key] = deep_merge_dict(result[key], value)
        else:
            result[key] = value
    return result


def merged_constraint_for_slide(slide_ir: Dict[str, Any], constraints: Optional[ConstraintSet]) -> Dict[str, Any]:
    if constraints is None:
        return {}
    merged: Dict[str, Any] = dict(constraints.defaults or {})
    kind = slide_ir.get("kind")
    layout = slide_ir.get("layout", "")
    title = slide_ir.get("slide_title", "")

    if kind == "presentation":
        merged = deep_merge_dict(merged, constraints.presentation_layouts.get(layout, {}))
    elif kind == "section":
        merged = deep_merge_dict(merged, constraints.section_layouts.get(layout, {}))
    else:
        merged = deep_merge_dict(merged, constraints.layouts.get(layout, {}))

    merged = deep_merge_dict(merged, constraints.slides.get(title, {}))
    return merged


def constraint_block_for_object(obj: Dict[str, Any], slide_constraints: Dict[str, Any]) -> Dict[str, Any]:
    """Choose the PCL block relevant to a render object."""
    obj_type = obj.get("type")
    role = str(obj.get("role", "")).lower()
    role_block = slide_constraints.get(role, {}) if role else {}

    if obj_type == "BulletList":
        type_block = slide_constraints.get("bullets", {})
    elif obj_type == "TableBox":
        type_block = slide_constraints.get("table", {})
    elif obj_type == "ImageBox":
        type_block = slide_constraints.get("image", {})
    elif obj_type == "CardBox":
        type_block = slide_constraints.get("card", {})
    elif obj_type in {"TextBox", "Hyperlink"}:
        type_block = slide_constraints.get("text", {})
    else:
        type_block = {}

    if role.startswith("footer"):
        type_block = deep_merge_dict(type_block if isinstance(type_block, dict) else {}, slide_constraints.get("footer", {}))

    if isinstance(role_block, dict):
        return deep_merge_dict(type_block if isinstance(type_block, dict) else {}, role_block)
    return type_block if isinstance(type_block, dict) else {}


def apply_constraints_to_object(obj: Dict[str, Any], slide_constraints: Dict[str, Any]) -> Dict[str, Any]:
    new_obj = clone_obj(obj)
    block = constraint_block_for_object(new_obj, slide_constraints)
    if not block:
        return new_obj

    runtime_keys = {
        "overflow", "max_lines", "min_size", "max_size", "auto_shrink",
        "max_items_per_slide", "max_rows_per_slide", "line_gap",
        "mode", "fit", "align", "keep_inside", "avoid_overlap",
        "margin", "padding", "gap", "fill", "border", "border_color",
        "width", "height", "x", "y",
    }
    for key, value in block.items():
        if key in runtime_keys or key.startswith("x-"):
            if key == "fit":
                new_obj["mode"] = value
            else:
                new_obj[key] = value
    return new_obj


def apply_constraints(render_ir: Dict[str, Any], constraints: Optional[ConstraintSet]) -> Dict[str, Any]:
    """Apply PCL constraints to Render IR before pagination/rendering."""
    if constraints is None:
        return render_ir

    new_ir = dict(render_ir)
    new_slides: List[Dict[str, Any]] = []
    for slide_ir in render_ir.get("slides", []):
        sc = merged_constraint_for_slide(slide_ir, constraints)
        new_slide = dict(slide_ir)

        if isinstance(sc.get("slide"), dict):
            slide_block = sc["slide"]
            if "background" in slide_block:
                new_slide["background"] = slide_block["background"]
            if "background_image" in slide_block:
                new_slide["background_image"] = slide_block["background_image"]

        if isinstance(sc.get("background_image"), dict):
            bg = dict(new_slide.get("background_image") or {})
            bg.update(sc["background_image"])
            new_slide["background_image"] = bg

        new_slide["objects"] = [apply_constraints_to_object(obj, sc) for obj in slide_ir.get("objects", [])]
        new_slides.append(new_slide)

    new_ir["slides"] = new_slides
    return new_ir


def paginate_single_slide(slide_ir: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Split one render slide into multiple slides when content actually overflows.

    Supported objects:
    - BulletList with overflow: paginate
      Splits by estimated text height in the object box. max_items_per_slide remains
      a hard cap, but it is no longer the only criterion.
    - TableBox with overflow: paginate
      Splits by estimated row height. max_rows_per_slide remains a hard cap.

    This is still a heuristic engine, not a real PowerPoint layout measurement engine,
    but it prevents the common case where large text or long Vietnamese lines get
    clipped despite item count being small.
    """
    objects = slide_ir.get("objects", [])

    for idx, obj in enumerate(objects):
        if obj.get("type") == "BulletList" and obj.get("overflow") == "paginate":
            pages_items = split_bullet_items_by_height(obj)
            if len(pages_items) <= 1:
                return [slide_ir]
            pages = []
            for page_index, chunk in enumerate(pages_items):
                new_slide = dict(slide_ir)
                new_objects = [clone_obj(o) for o in objects]
                new_objects[idx]["items"] = chunk
                new_slide["objects"] = new_objects
                mark_continued(new_slide, page_index)
                pages.append(new_slide)
            return pages

        if obj.get("type") == "TableBox" and obj.get("overflow") == "paginate":
            pages_rows = split_table_rows_by_height(obj)
            if len(pages_rows) <= 1:
                return [slide_ir]
            pages = []
            for page_index, chunk in enumerate(pages_rows):
                new_slide = dict(slide_ir)
                new_objects = [clone_obj(o) for o in objects]
                new_objects[idx]["rows"] = chunk
                new_slide["objects"] = new_objects
                mark_continued(new_slide, page_index)
                pages.append(new_slide)
            return pages

    return [slide_ir]

def paginate_render_ir(render_ir: Dict[str, Any]) -> Dict[str, Any]:
    paginated: List[Dict[str, Any]] = []
    for slide_ir in render_ir.get("slides", []):
        paginated.extend(paginate_single_slide(slide_ir))
    render_ir = dict(render_ir)
    render_ir["slides"] = paginated
    return render_ir

def build_render_ir(doc: PresentationDoc, theme: Theme, constraints: Optional[ConstraintSet] = None) -> Dict[str, Any]:
    slides_ir: List[Dict[str, Any]] = []

    # Presentation cover/title slide
    cover_style = theme.presentation_layouts.get(doc.cover_layout, {})
    cover_objects = [make_text_object("title", doc.title, cover_style.get("title", {}), theme)]
    cover_objects.extend(
        build_text_objects_from_block(doc.cover, cover_style, theme, ["subtitle", "author", "date"])
    )
    slides_ir.append({
        "type": "RenderSlide",
        "kind": "presentation",
        "section": None,
        "slide_title": doc.title,
        "layout": doc.cover_layout,
        "background": style_value(cover_style, "background", theme, theme.page.get("background", "#FFFFFF")),
        "background_image": make_background_image_spec(doc.cover, cover_style),
        "objects": cover_objects,
        "notes": doc.cover.get("notes", "") if isinstance(doc.cover, dict) else "",
    })

    for section in doc.sections:
        # Section header slide
        header_style = theme.section_layouts.get(section.header_layout, {})
        header_objects = [make_text_object("title", section.title, header_style.get("title", {}), theme)]
        header_objects.extend(
            build_text_objects_from_block(section.header, header_style, theme, ["subtitle", "summary", "label"])
        )
        slides_ir.append({
            "type": "RenderSlide",
            "kind": "section",
            "section": section.title,
            "slide_title": section.title,
            "layout": section.header_layout,
            "background": style_value(header_style, "background", theme, theme.page.get("background", "#FFFFFF")),
            "background_image": make_background_image_spec(section.header, header_style),
            "objects": header_objects,
            "notes": section.header.get("notes", "") if isinstance(section.header, dict) else "",
        })

        # Content slides
        for slide in section.slides:
            layout_style = theme.layouts.get(slide.layout, {})
            objects: List[Dict[str, Any]] = []

            # Native image-caption layout: image + caption only.
            # It intentionally ignores title/subtitle/body/bullets even if present.
            if slide.layout == "image-caption":
                if "image" in slide.blocks:
                    objects.append(make_image_object(slide.blocks["image"], layout_style.get("image", {}), theme))
                caption_text = slide.blocks.get("caption", "")
                if not caption_text and isinstance(slide.blocks.get("image"), dict):
                    caption_text = slide.blocks["image"].get("alt", "")
                if caption_text:
                    objects.append(make_text_object("caption", caption_text, layout_style.get("caption", {}), theme))
            else:
                if "title" in slide.blocks:
                    objects.append(make_text_object("title", slide.blocks["title"], layout_style.get("title", {}), theme))

                if "subtitle" in slide.blocks:
                    objects.append(make_text_object("subtitle", slide.blocks["subtitle"], layout_style.get("subtitle", {}), theme))

                if "body" in slide.blocks:
                    objects.append(make_text_object("body", slide.blocks["body"], layout_style.get("body", {}), theme))

                if "bullets" in slide.blocks:
                    objects.append(make_bullet_object(slide.blocks["bullets"], layout_style.get("bullets", {}), theme))

                if "image" in slide.blocks:
                    objects.append(make_image_object(slide.blocks["image"], layout_style.get("image", {}), theme))

                if "link" in slide.blocks:
                    objects.append(make_link_object(slide.blocks["link"], layout_style.get("link", {}), theme))

                if "table" in slide.blocks:
                    objects.append(make_table_object(slide.blocks["table"], layout_style.get("table", {}), theme))

                # Native hero-image layout: one large image with a caption below.
                # Caption priority: explicit caption block -> image.alt fallback.
                if slide.layout == "hero-image" and "image" in slide.blocks:
                    caption_text = slide.blocks.get("caption", "")
                    if not caption_text and isinstance(slide.blocks["image"], dict):
                        caption_text = slide.blocks["image"].get("alt", "")
                    if caption_text:
                        objects.append(make_text_object("caption", caption_text, layout_style.get("caption", {}), theme))

            if "cells" in slide.blocks:
                objects.extend(make_grid_objects(slide.blocks["cells"], layout_style, theme, slide.layout))

            if slide.layout.startswith("numbered-columns") and ("columns" in slide.blocks or "cells" in slide.blocks or "items" in slide.blocks):
                objects.extend(make_numbered_columns_objects(slide.blocks, layout_style, theme, slide.layout))

            if slide.layout == "stair-progress" and ("steps" in slide.blocks or "items" in slide.blocks):
                objects.extend(make_stair_progress_objects(slide.blocks, layout_style, theme))

            if slide.layout == "stacked-stairs" and ("steps" in slide.blocks or "items" in slide.blocks):
                objects.extend(make_stacked_stairs_objects(slide.blocks, layout_style, theme))

            if slide.layout == "timeline" and ("milestones" in slide.blocks or "steps" in slide.blocks or "items" in slide.blocks):
                objects.extend(make_timeline_objects(slide.blocks, layout_style, theme))

            if slide.layout == "two-images":
                objects.extend(make_two_image_objects(slide.blocks, layout_style, theme))
            else:
                for side in ["left", "right"]:
                    if side in slide.blocks:
                        objects.extend(make_column_objects(side, slide.blocks[side], layout_style.get(side, {}), theme))

            # Footer objects are appended last so they appear above normal content.
            # Supported on all content slide layouts.
            objects.extend(make_footer_text_objects(slide.blocks, layout_style, theme))
            objects.extend(make_footer_image_objects(slide.blocks, layout_style, theme))

            slides_ir.append({
                "type": "RenderSlide",
                "kind": "slide",
                "section": section.title,
                "slide_title": slide.title,
                "layout": slide.layout,
                "background": style_value(layout_style, "background", theme, theme.page.get("background", "#FFFFFF")),
                "background_image": make_background_image_spec(slide.blocks, layout_style),
                "objects": objects,
                "notes": slide.blocks.get("notes", ""),
            })

    render_ir = {
        "type": "RenderPresentation",
        "title": doc.title,
        "page": theme.page,
        "slides": slides_ir,
    }
    render_ir = apply_constraints(render_ir, constraints)
    return paginate_render_ir(render_ir)



# -----------------------------------------------------------------------------
# Asset resolution helpers
# -----------------------------------------------------------------------------


def is_remote_src(src: str) -> bool:
    return str(src).startswith(("http://", "https://", "data:"))


def resolve_asset_path(src: Any, asset_base_dir: Optional[Path] = None) -> Optional[Path]:
    """Resolve a local image path.

    Why this exists:
    - PML usually writes paths relative to the DSL/script folder.
    - HTML/PPTX renderers may be called from another current working directory.
    - Without resolution, images silently disappear because src cannot be found.
    """
    if src is None:
        return None
    src_str = str(src).strip()
    if not src_str or is_remote_src(src_str):
        return None

    path = Path(src_str)
    if path.is_absolute():
        return path

    candidates = []
    if asset_base_dir is not None:
        candidates.append(Path(asset_base_dir) / path)
    candidates.append(Path.cwd() / path)

    # Return the first existing candidate, otherwise return the preferred base-dir candidate
    # so fallback messages can display a meaningful path.
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[0] if candidates else path


def html_src_for_image(src: Any, output_path: str, asset_base_dir: Optional[Path] = None) -> Tuple[str, bool]:
    """Return src value for HTML and whether the local asset exists.

    For local files, prefer a path relative to the HTML output directory. This makes
    the generated HTML portable when opened in a browser.
    """
    src_str = str(src or "").strip()
    if not src_str:
        return "", False
    if is_remote_src(src_str):
        return src_str, True

    resolved = resolve_asset_path(src_str, asset_base_dir)
    if resolved is None:
        return src_str, False

    exists = resolved.exists()
    try:
        rel = resolved.resolve().relative_to(Path(output_path).resolve().parent)
        return rel.as_posix(), exists
    except Exception:
        return resolved.resolve().as_posix(), exists


def create_demo_assets(asset_dir: Path) -> None:
    """Create tiny placeholder PNGs used by DEMO_PML.

    This is only for the bundled demo. In real use, the user should supply actual
    image files. The renderer itself does not require these exact filenames.
    """
    names = [
        "demo_bg.png", "corporate_bg.png", "viettel_logo.png", "architecture.png",
        "architecture_overview.png", "model_snapshot.png", "current_architecture.png",
        "target_architecture.png", "cover_bg.png",
    ]
    if PILImage is None:
        return
    asset_dir.mkdir(parents=True, exist_ok=True)
    for idx, name in enumerate(names):
        path = asset_dir / name
        if path.exists():
            continue
        w, h = (1280, 720) if "bg" in name or "cover" in name else (900, 500)
        img = PILImage.new("RGB", (w, h), (235 - idx * 7 % 80, 242 - idx * 11 % 80, 250 - idx * 13 % 80))
        # Avoid depending on ImageDraw fonts; a plain colored rectangle is enough
        # to verify that image placement works.
        img.save(path)

# -----------------------------------------------------------------------------
# Renderers
# -----------------------------------------------------------------------------

def hex_to_rgb(hex_color: str):
    if not hex_color:
        hex_color = "#000000"
    hex_color = str(hex_color).strip().lstrip("#")
    if len(hex_color) != 6:
        hex_color = "000000"
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def px_to_inches(value: float) -> float:
    return value / 96.0



def add_card(slide, obj: Dict[str, Any]) -> None:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if MSO_SHAPE is not None else 1
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 100))),
    )
    fill_value = obj.get("fill", "#FFFFFF")
    if fill_value in (None, "none", "transparent"):
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_value)
    border_value = obj.get("border", obj.get("border_color", "#D1D5DB"))
    if border_value in (None, "none", "transparent"):
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(border_value)


def add_linebox(slide, obj: Dict[str, Any]) -> None:
    shape = slide.shapes.add_shape(
        1,
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 4))),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(obj.get("fill", "#2563EB"))
    shape.line.fill.background()


def add_circle(slide, obj: Dict[str, Any]) -> None:
    shape_type = MSO_SHAPE.OVAL if MSO_SHAPE is not None else 9
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", obj.get("w", 22)))),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(obj.get("fill", "#2563EB"))
    shape.line.color.rgb = hex_to_rgb(obj.get("border", "#FFFFFF"))


def estimate_chars_per_line(width_px: Any, font_size: Any) -> int:
    """Very small heuristic for PPTX text fitting.

    It estimates average Latin/Vietnamese glyph width as ~0.55em.
    This is not typography-perfect, but prevents the worst slide overflows.
    """
    width = max(1.0, safe_float(width_px, 400.0))
    size = max(1.0, safe_float(font_size, 20.0))
    return max(8, int(width / (size * 0.55)))


def estimate_needed_lines(text: str, width_px: Any, font_size: Any) -> int:
    chars_per_line = estimate_chars_per_line(width_px, font_size)
    total = 0
    for raw_line in str(text).splitlines() or [""]:
        n = len(raw_line)
        total += max(1, (n + chars_per_line - 1) // chars_per_line)
    return total


def infer_max_lines(obj: Dict[str, Any]) -> int:
    explicit = safe_int(obj.get("max_lines", 0), 0)
    if explicit > 0:
        return explicit
    size = max(1.0, safe_float(obj.get("size", 20), 20.0))
    height = max(1.0, safe_float(obj.get("h", 100), 100.0))
    # line height about 1.25em
    return max(1, int(height / (size * 1.25)))


def ellipsize_to_lines(text: str, width_px: Any, font_size: Any, max_lines: int) -> str:
    if max_lines <= 0:
        return text
    chars_per_line = estimate_chars_per_line(width_px, font_size)
    limit = max(1, chars_per_line * max_lines)
    text = str(text)
    if len(text) <= limit:
        return text
    return text[: max(1, limit - 1)].rstrip() + "…"


def fitted_font_size_for_text(text: str, obj: Dict[str, Any]) -> int:
    size = safe_int(obj.get("size", 20), 20)
    min_size = safe_int(obj.get("min_size", 12), 12)
    max_lines = infer_max_lines(obj)
    while size > min_size and estimate_needed_lines(text, obj.get("w", 400), size) > max_lines:
        size -= 1
    return max(min_size, size)


def prepare_pptx_text(text: str, obj: Dict[str, Any]) -> Tuple[str, int]:
    """Return text and font size after applying overflow policy for PPTX."""
    mode = normalize_overflow(obj.get("overflow", "wrap"))
    size = safe_int(obj.get("size", 20), 20)
    max_lines = infer_max_lines(obj)

    if mode == "shrink":
        size = fitted_font_size_for_text(text, obj)
        # If even min_size still overflows, ellipsize as last-resort safety.
        if estimate_needed_lines(text, obj.get("w", 400), size) > max_lines:
            text = ellipsize_to_lines(text, obj.get("w", 400), size, max_lines)
    elif mode == "ellipsis":
        text = ellipsize_to_lines(text, obj.get("w", 400), size, max_lines)
    elif mode == "clip":
        # Same string truncation as ellipsis, but no ellipsis marker.
        clipped = ellipsize_to_lines(text, obj.get("w", 400), size, max_lines)
        text = clipped[:-1] if clipped.endswith("…") else clipped

    return text, size


def pptx_alignment(align: Any):
    if PP_ALIGN is None:
        return None
    align = normalize_align(align, "left")
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(align, PP_ALIGN.LEFT)


def apply_pptx_text_style(paragraph, obj: Dict[str, Any]) -> None:
    paragraph.alignment = pptx_alignment(obj.get("align", "left"))
    paragraph.font.name = obj["font"]
    paragraph.font.size = Pt(obj.get("_render_size", obj["size"]))
    paragraph.font.color.rgb = hex_to_rgb(obj["color"])
    paragraph.font.bold = parse_bool(obj.get("bold", False))
    paragraph.font.italic = parse_bool(obj.get("italic", False))
    paragraph.font.underline = parse_bool(obj.get("underline", False))

def bullet_item_to_plain_text(item: Any) -> str:
    if isinstance(item, dict):
        icon = str(item.get("icon", ""))
        text = str(item.get("text", ""))
        return f"{icon} {text}".strip()
    return str(item)


def apply_pptx_hyperlink(paragraph, url: str) -> None:
    """Apply hyperlink to the first run of a paragraph when possible."""
    if not url:
        return
    try:
        if paragraph.runs:
            paragraph.runs[0].hyperlink.address = url
    except Exception:
        # Keep rendering robust even if python-pptx cannot set hyperlink for a run.
        pass


def add_textbox(slide, obj: Dict[str, Any], bullet: bool = False) -> None:
    shape = slide.shapes.add_textbox(
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 100))),
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    if bullet:
        for idx, item in enumerate(obj["items"]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            text, render_size = prepare_pptx_text(bullet_item_to_plain_text(item), obj)
            p.text = text
            p.level = 0
            obj2 = dict(obj)
            obj2["_render_size"] = render_size
            apply_pptx_text_style(p, obj2)
            if isinstance(item, dict):
                apply_pptx_hyperlink(p, str(item.get("url", "")))
    else:
        p = tf.paragraphs[0]
        text, render_size = prepare_pptx_text(obj["text"], obj)
        p.text = text
        obj2 = dict(obj)
        obj2["_render_size"] = render_size
        apply_pptx_text_style(p, obj2)
        apply_pptx_hyperlink(p, str(obj.get("url", "")))


def add_table(slide, obj: Dict[str, Any]) -> None:
    headers = obj.get("headers", []) or []
    rows = obj.get("rows", []) or []
    total_rows = len(rows) + (1 if headers else 0)
    col_count = max([len(headers)] + [len(r) for r in rows] + [1])
    if total_rows <= 0:
        return

    shape = slide.shapes.add_table(
        total_rows,
        col_count,
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 300))),
    )
    table = shape.table

    def set_cell(cell, text: str, *, is_header: bool = False, row_idx: int = 0) -> None:
        cell.text = str(text)
        fill = cell.fill
        fill.solid()
        if is_header:
            fill.fore_color.rgb = hex_to_rgb(obj.get("header_fill", "#003A8C"))
            color = obj.get("header_color", "#FFFFFF")
            align = obj.get("header_align", obj.get("align", "center"))
            bold = obj.get("header_bold", True)
        else:
            fill.fore_color.rgb = hex_to_rgb(obj.get("alt_fill" if row_idx % 2 else "cell_fill", "#FFFFFF"))
            color = obj.get("color", "#111827")
            align = obj.get("align", "left")
            bold = obj.get("bold", False)

        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = pptx_alignment(align)
            paragraph.font.name = obj.get("font", "Aptos")
            paragraph.font.size = Pt(obj.get("size", 18))
            paragraph.font.color.rgb = hex_to_rgb(color)
            paragraph.font.bold = parse_bool(bold)

    r_offset = 0
    if headers:
        for c in range(col_count):
            set_cell(table.cell(0, c), headers[c] if c < len(headers) else "", is_header=True)
        r_offset = 1

    for r, row in enumerate(rows):
        for c in range(col_count):
            set_cell(table.cell(r + r_offset, c), row[c] if c < len(row) else "", is_header=False, row_idx=r)



def get_image_size_px(path: Path) -> Optional[Tuple[int, int]]:
    """Return image size in pixels when Pillow is available."""
    if PILImage is None:
        return None
    try:
        with PILImage.open(path) as img:
            return img.size
    except Exception:
        return None


def fitted_box_px(obj: Dict[str, Any], image_size: Optional[Tuple[int, int]]) -> Tuple[float, float, float, float]:
    """Compute x/y/w/h in DSL pixels for contain/original/stretch.

    Cover uses the original target box and crop properties in PPTX.
    """
    x, y, w, h = float(obj["x"]), float(obj["y"]), float(obj["w"]), float(obj.get("h", 100))
    mode = normalize_fit_mode(obj.get("mode", "stretch"))
    if not image_size or mode in {"stretch", "cover"}:
        return x, y, w, h

    iw, ih = image_size
    if iw <= 0 or ih <= 0:
        return x, y, w, h

    if mode == "original":
        new_w, new_h = float(iw), float(ih)
    else:  # contain
        scale = min(w / iw, h / ih)
        new_w, new_h = iw * scale, ih * scale

    return x + (w - new_w) / 2, y + (h - new_h) / 2, new_w, new_h


def apply_cover_crop(pic, box_w: float, box_h: float, image_size: Optional[Tuple[int, int]]) -> None:
    """Crop a PPTX picture so it behaves like CSS object-fit: cover."""
    if not image_size:
        return
    iw, ih = image_size
    if iw <= 0 or ih <= 0 or box_w <= 0 or box_h <= 0:
        return

    img_ar = iw / ih
    box_ar = box_w / box_h

    if img_ar > box_ar:
        visible_w = box_ar * ih
        crop_each = max(0.0, min(0.5, (iw - visible_w) / (2 * iw)))
        pic.crop_left = crop_each
        pic.crop_right = crop_each
    elif img_ar < box_ar:
        visible_h = iw / box_ar
        crop_each = max(0.0, min(0.5, (ih - visible_h) / (2 * ih)))
        pic.crop_top = crop_each
        pic.crop_bottom = crop_each


def add_image(slide, obj: Dict[str, Any], asset_base_dir: Optional[Path] = None) -> None:
    src = obj.get("src", "")
    path = resolve_asset_path(src, asset_base_dir)
    if path is None or not path.exists():
        # Graceful fallback: show missing path instead of crashing.
        missing = dict(obj)
        missing.update({
            "type": "TextBox",
            "role": "missing-image",
            "text": f"[Missing image: {src}]",
            "font": "Aptos",
            "size": 16,
            "color": "#AA0000",
        })
        add_textbox(slide, missing, bullet=False)
        return

    mode = normalize_fit_mode(obj.get("mode", "stretch"))
    image_size = get_image_size_px(path)
    x, y, w, h = fitted_box_px(obj, image_size)

    pic = slide.shapes.add_picture(
        str(path),
        Inches(px_to_inches(x)),
        Inches(px_to_inches(y)),
        width=Inches(px_to_inches(w)),
        height=Inches(px_to_inches(h)),
    )
    if mode == "cover":
        apply_cover_crop(pic, float(obj["w"]), float(obj.get("h", 100)), image_size)



def add_background_image(slide, spec: Optional[Dict[str, Any]], asset_base_dir: Optional[Path] = None) -> None:
    """Add a full-slide background image behind all objects.

    Supported PPTX modes:
      stretch  - fill slide, may distort.
      contain  - preserve aspect ratio; blank margins may remain.
      cover/crop - fill slide and crop overflow.
      original - use original pixel size, centered on slide.

    Note: python-pptx does not provide a simple public API for picture opacity.
    The opacity value is preserved in HTML output; PPTX uses the image as-is.
    """
    if not spec:
        return
    src = spec.get("src", "")
    path = resolve_asset_path(src, asset_base_dir)
    if path is None or not path.exists():
        return

    obj = {
        "x": 0,
        "y": 0,
        "w": 1280,
        "h": 720,
        "mode": normalize_fit_mode(spec.get("mode", "cover"), default="cover"),
    }
    image_size = get_image_size_px(path)
    x, y, w, h = fitted_box_px(obj, image_size)

    pic = slide.shapes.add_picture(
        str(path),
        Inches(px_to_inches(x)),
        Inches(px_to_inches(y)),
        width=Inches(px_to_inches(w)),
        height=Inches(px_to_inches(h)),
    )
    if obj["mode"] == "cover":
        apply_cover_crop(pic, 1280, 720, image_size)

def render_pptx(render_ir: Dict[str, Any], output_path: str, asset_base_dir: Optional[Path] = None) -> None:
    if PPTXPresentation is None:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = PPTXPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for slide_ir in render_ir["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_to_rgb(slide_ir.get("background", "#FFFFFF"))

        # Background image is drawn first so all normal objects stay on top.
        add_background_image(slide, slide_ir.get("background_image"), asset_base_dir)

        for obj in slide_ir["objects"]:
            if obj["type"] == "Card":
                add_card(slide, obj)
            elif obj["type"] == "LineBox":
                add_linebox(slide, obj)
            elif obj["type"] == "Circle":
                add_circle(slide, obj)
            elif obj["type"] == "TextBox":
                add_textbox(slide, obj, bullet=False)
            elif obj["type"] == "BulletList":
                add_textbox(slide, obj, bullet=True)
            elif obj["type"] == "TableBox":
                add_table(slide, obj)
            elif obj["type"] == "Image":
                add_image(slide, obj, asset_base_dir)

        notes = slide_ir.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    prs.save(output_path)



def html_text_css(obj: Dict[str, Any]) -> str:
    overflow = normalize_overflow(obj.get("overflow", "wrap"))
    max_lines = safe_int(obj.get("max_lines", 0), 0)
    css = (
        f"font-weight:{'700' if parse_bool(obj.get('bold', False)) else '400'}; "
        f"font-style:{'italic' if parse_bool(obj.get('italic', False)) else 'normal'}; "
        f"text-decoration:{'underline' if parse_bool(obj.get('underline', False)) else 'none'}; "
        f"text-align:{normalize_align(obj.get('align', 'left'))}; "
        "overflow-wrap:anywhere; word-break:normal; line-height:1.25; "
    )
    if overflow == "clip":
        css += "overflow:hidden; "
    elif overflow == "ellipsis":
        css += "overflow:hidden; "
        if max_lines > 0:
            css += f"display:-webkit-box; -webkit-line-clamp:{max_lines}; -webkit-box-orient:vertical; "
        else:
            css += "white-space:nowrap; text-overflow:ellipsis; "
    elif overflow == "shrink":
        # Browser-side true auto-shrink needs JavaScript. Keep wrapping and clipping safe.
        css += "overflow:hidden; "
        if max_lines > 0:
            css += f"display:-webkit-box; -webkit-line-clamp:{max_lines}; -webkit-box-orient:vertical; "
    else:
        css += "overflow-wrap:anywhere; "
    return css

def html_link_attrs(url: str) -> str:
    if not url:
        return ""
    safe = html.escape(str(url), quote=True)
    return f" href='{safe}' target='_blank' rel='noopener noreferrer'"


def html_linked_text(text: str, url: str, extra_style: str = "") -> str:
    escaped = html.escape(str(text))
    if not url:
        return escaped
    style_attr = f" style='{extra_style}'" if extra_style else ""
    return f"<a{html_link_attrs(url)}{style_attr}>{escaped}</a>"


def render_html(render_ir: Dict[str, Any], output_path: str, asset_base_dir: Optional[Path] = None) -> None:
    parts: List[str] = []
    parts.append("<!doctype html>")
    parts.append("<html><head><meta charset='utf-8'>")
    parts.append("<title>" + html.escape(render_ir["title"]) + "</title>")
    parts.append(
        """
<style>
body { margin: 0; font-family: Arial, sans-serif; background: #eee; }
.deck { display: flex; flex-direction: column; gap: 24px; padding: 24px; }
.slide { position: relative; width: 1280px; height: 720px; box-shadow: 0 8px 24px rgba(0,0,0,.18); overflow: hidden; }
.textbox, .bullets, .image, .card, .linebox, .circle, .bg-image, .tablebox { position: absolute; box-sizing: border-box; }
ul { margin: 0; padding-left: 1.2em; }
li { margin-bottom: .4em; }
img { display: block; }
table.dsl-table { border-collapse: collapse; width: 100%; height: 100%; table-layout: fixed; }
table.dsl-table th, table.dsl-table td { padding: 10px 12px; overflow-wrap: anywhere; vertical-align: top; }
a { color: inherit; text-decoration: inherit; }
</style>
"""
    )
    parts.append("</head><body><div class='deck'>")

    for slide in render_ir["slides"]:
        bg = slide.get("background", "#FFFFFF")
        parts.append(f"<section class='slide' style='background:{html.escape(str(bg))}'>")
        bg_img = slide.get("background_image")
        if bg_img and bg_img.get("src"):
            mode = normalize_fit_mode(bg_img.get("mode", "cover"), default="cover")
            fit_map = {"stretch": "fill", "contain": "contain", "cover": "cover", "original": "none"}
            fit = fit_map.get(mode, "cover")
            opacity = bg_img.get("opacity", 1.0)
            try:
                opacity = max(0.0, min(1.0, float(opacity)))
            except (TypeError, ValueError):
                opacity = 1.0
            src_raw = str(bg_img.get("src", ""))
            src_value, exists = html_src_for_image(src_raw, output_path, asset_base_dir)
            src = html.escape(src_value)
            alt = html.escape(str(bg_img.get("alt", "")))
            if exists or is_remote_src(src_raw):
                parts.append(
                    f"<img class='bg-image' src='{src}' alt='{alt}' "
                    f"style='left:0; top:0; width:1280px; height:720px; object-fit:{fit}; opacity:{opacity}; z-index:0;'>"
                )
        for obj in slide["objects"]:
            if obj["type"] == "Card":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 100)}px; "
                    f"background:{html.escape(str(obj.get('fill', '#FFFFFF')))}; "
                    f"border:1px solid {html.escape(str(obj.get('border', '#D1D5DB')))}; "
                    f"border-radius:{obj.get('radius', 16)}px;"
                )
                parts.append(f"<div class='card' style='{style}'></div>")
            elif obj["type"] == "LineBox":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 4)}px; "
                    f"background:{html.escape(str(obj.get('fill', '#2563EB')))}; "
                    f"border-radius:999px;"
                )
                parts.append(f"<div class='linebox' style='{style}'></div>")
            elif obj["type"] == "Circle":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', obj.get('w', 22))}px; "
                    f"background:{html.escape(str(obj.get('fill', '#2563EB')))}; "
                    f"border:3px solid {html.escape(str(obj.get('border', '#FFFFFF')))}; "
                    f"border-radius:999px;"
                )
                parts.append(f"<div class='circle' style='{style}'></div>")
            elif obj["type"] == "TextBox":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 100)}px; "
                    f"font-family:{html.escape(str(obj['font']))}; "
                    f"font-size:{obj['size']}px; color:{html.escape(str(obj['color']))};"
                    f"{html_text_css(obj)}"
                )
                url = str(obj.get("url", ""))
                parts.append(f"<div class='textbox' style='{style}'>" + html_linked_text(obj["text"], url) + "</div>")
            elif obj["type"] == "BulletList":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 420)}px; "
                    f"font-family:{html.escape(str(obj['font']))}; "
                    f"font-size:{obj['size']}px; color:{html.escape(str(obj['color']))};"
                    f"{html_text_css(obj)}"
                )
                icon_gap = obj.get("icon_gap", 12)
                line_gap = obj.get("line_gap", 8)
                icon_size = obj.get("icon_size", obj.get("size", 24))
                icon_color = html.escape(str(obj.get("icon_color", obj.get("color", "#000000"))))
                parts.append(f"<div class='bullets' style='{style}'>")
                for item in obj["items"]:
                    if isinstance(item, dict):
                        icon = html.escape(str(item.get("icon", "")))
                        text_raw = str(item.get("text", ""))
                        url = str(item.get("url", ""))
                    else:
                        icon = ""
                        text_raw = str(item)
                        url = ""
                    text_html = html_linked_text(text_raw, url)
                    if icon:
                        parts.append(
                            f"<div style='display:flex; align-items:flex-start; gap:{icon_gap}px; margin-bottom:{line_gap}px;'>"
                            f"<span style='color:{icon_color}; font-size:{icon_size}px; line-height:1.2; min-width:{icon_size}px; flex:0 0 {icon_size}px;'>"
                            f"{icon}</span><span style='line-height:1.25; flex:1; min-width:0; overflow-wrap:anywhere;'>{text_html}</span></div>"
                        )
                    else:
                        parts.append(f"<div style='margin-bottom:{line_gap}px; line-height:1.25;'>• {text_html}</div>")
                parts.append("</div>")
            elif obj["type"] == "TableBox":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 300)}px; "
                    f"font-family:{html.escape(str(obj.get('font', 'Arial')))}; "
                    f"font-size:{obj.get('size', 18)}px; color:{html.escape(str(obj.get('color', '#111827')))};"
                )
                border = html.escape(str(obj.get("border_color", "#CBD5E1")))
                header_fill = html.escape(str(obj.get("header_fill", "#003A8C")))
                header_color = html.escape(str(obj.get("header_color", "#FFFFFF")))
                cell_fill = html.escape(str(obj.get("cell_fill", "#FFFFFF")))
                alt_fill = html.escape(str(obj.get("alt_fill", "#F8FAFC")))
                align = html.escape(normalize_align(obj.get("align", "left")))
                header_align = html.escape(normalize_align(obj.get("header_align", "center")))
                parts.append(f"<div class='tablebox' style='{style}'><table class='dsl-table'>")
                if obj.get("headers"):
                    parts.append("<thead><tr>")
                    for cell in obj.get("headers", []):
                        parts.append(
                            f"<th style='background:{header_fill}; color:{header_color}; border:1px solid {border}; text-align:{header_align}; font-weight:700;'>"
                            f"{html.escape(str(cell))}</th>"
                        )
                    parts.append("</tr></thead>")
                parts.append("<tbody>")
                for idx, row in enumerate(obj.get("rows", [])):
                    fill = cell_fill if idx % 2 == 0 else alt_fill
                    parts.append("<tr>")
                    for cell in row:
                        parts.append(
                            f"<td style='background:{fill}; border:1px solid {border}; text-align:{align};'>"
                            f"{html.escape(str(cell))}</td>"
                        )
                    parts.append("</tr>")
                parts.append("</tbody></table></div>")
            elif obj["type"] == "Image":
                mode = normalize_fit_mode(obj.get("mode", "stretch"))
                fit_map = {"stretch": "fill", "contain": "contain", "cover": "cover", "original": "none"}
                fit = fit_map.get(mode, "fill")
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj['h']}px;"
                    f"object-fit:{fit};"
                )
                src_raw = str(obj.get("src", ""))
                src_value, exists = html_src_for_image(src_raw, output_path, asset_base_dir)
                src = html.escape(src_value)
                alt = html.escape(str(obj.get("alt", "")))
                if exists or is_remote_src(src_raw):
                    parts.append(f"<img class='image' src='{src}' alt='{alt}' style='{style}'>")
                else:
                    placeholder_style = style + "border:2px dashed #AA0000; color:#AA0000; padding:12px; font-family:Arial; font-size:16px;"
                    parts.append(f"<div class='image' style='{placeholder_style}'>Missing image: {html.escape(src_raw)}</div>")
        parts.append("</section>")

    parts.append("</div></body></html>")
    Path(output_path).write_text("\n".join(parts), encoding="utf-8")


# -----------------------------------------------------------------------------
# Demo input
# -----------------------------------------------------------------------------

DEMO_PML = '''
presentation "AI Strategy 2026":
  meta:
    author: "Viettel AI Lab"
    language: vi
    format: pptx

  use style: "corporate.psl"
  use constraints: "safe-layouts.pcl"

  cover_layout: title-slide
  cover:
    subtitle: Chiến lược tự động hóa tri thức doanh nghiệp
    author: Viettel AI Lab
    date: 2026
    background_image:
      src: "demo_bg.png"
      mode: cover
      opacity: 0.18

  section "Bối cảnh":
    header_layout: section-header
    header:
      subtitle: Từ dữ liệu phân tán đến insight có thể hành động

    slide "Vấn đề hiện tại":
      layout: title-bullets
      intent: explain_problem

      background_image:
        src: "corporate_bg.png"
        opacity: 0.12
        mode: cover

      title:
        Quá tải thông tin trong doanh nghiệp

      bullets:
        - Dữ liệu phân tán ở nhiều hệ thống
        - Báo cáo thủ công mất nhiều thời gian
        - Lãnh đạo cần insight nhanh hơn, nhưng nội dung bullet đôi khi rất dài nên cần cơ chế tự xuống dòng, co chữ hoặc cắt gọn để không phá bố cục slide.

      link:
        text: "Xem tài liệu pipeline"
        url: "https://example.com/pipeline"

      footer_text:
        Internal Use Only | Viettel AI Lab

      footer_image:
        src: "viettel_logo.png"
        alt: "Viettel logo"

      notes:
        Nhấn mạnh vấn đề không nằm ở thiếu dữ liệu, mà là thiếu khả năng tổng hợp.

    slide "Kiến trúc đề xuất":
      layout: title-image
      intent: show_architecture

      title:
        Kiến trúc tổng quan hệ thống

      subtitle:
        Dữ liệu → xử lý → tổng hợp → trình bày

      image:
        src: "architecture.png"
        alt: "Sơ đồ kiến trúc pipeline"

      notes:
        Có thể thay architecture.png bằng đường dẫn ảnh thật.

    slide "Mô hình triển khai":
      layout: text-image
      intent: explain_architecture

      title:
        Một cột nội dung, một cột hình minh họa

      left:
        heading: Thành phần chính
        bullets:
          - Thu thập dữ liệu từ nhiều nguồn
          - Chuẩn hóa và tạo chỉ mục
          - Truy xuất ngữ cảnh liên quan
          - Sinh báo cáo hoặc slide

      right:
        image:
          src: "architecture.png"
          alt: "Minh họa kiến trúc hệ thống"

      notes:
        Layout text-image dùng left cho chữ và right cho ảnh.

  section "Giải pháp":
    header_layout: section-header
    header:
      subtitle: Thiết kế pipeline và giao diện đầu ra

    slide "Luồng xử lý":
      layout: two-column
      intent: describe_pipeline

      title:
        Pipeline xử lý dữ liệu theo từng giai đoạn

      left:
        heading: Đầu vào
        bullets:
          - Văn bản
          - Báo cáo
          - Dữ liệu mạng xã hội

      right:
        heading: Đầu ra
        bullets:
          - Tóm tắt
          - Dashboard
          - Slide trình bày


    slide "Quy trình 4 bước có bullet":
      layout: numbered-columns-4
      intent: explain_process

      title:
        Quy trình xử lý theo 4 bước

      columns:
        - heading: Thu thập
          bullet_icon: check
          bullets:
            - API nội bộ
            - File CSV
            - Log hệ thống

        - heading: Làm sạch
          bullets:
            icon: diamond
            items:
              - Chuẩn hóa schema
              - Khử trùng lặp
              - Kiểm tra lỗi

        - heading: Truy xuất
          bullets:
            - icon: gear
              text: Tạo chỉ mục
            - icon: arrow
              text: Mở rộng ngữ cảnh

        - heading: Sinh đầu ra
          text: Tạo báo cáo hoặc slide.
          bullets:
            - Tóm tắt
            - Dashboard
            - PPTX


    slide "Kiến trúc tổng quan":
      layout: hero-image
      intent: show_overview_diagram

      title:
        Kiến trúc xử lý dữ liệu tổng quan

      subtitle:
        Một ảnh lớn ở trung tâm, chú thích ngắn ở phía dưới

      image:
        src: "architecture_overview.png"
        alt: "Sơ đồ kiến trúc tổng quan"

      caption:
        Hệ thống đi từ thu thập dữ liệu, chuẩn hóa, truy xuất tri thức đến sinh báo cáo và slide trình bày.

    slide "Ảnh minh họa độc lập":
      layout: image-caption
      intent: show_visual_only

      image:
        src: "model_snapshot.png"
        alt: "Ảnh minh họa mô hình tổng hợp tri thức"

      caption:
        Minh họa trực quan mô hình tổng hợp tri thức tự động, không cần tiêu đề phụ phía trên.

      footer:
        text: Confidential | Draft version
        align: center


    slide "So sánh hai phương án":
      layout: two-images
      intent: compare_visual_options

      title:
        So sánh kiến trúc hiện tại và kiến trúc đề xuất

      subtitle:
        Mỗi ảnh chiếm một cột độc lập để dễ đối chiếu trực quan

      images:
        - src: "current_architecture.png"
          alt: "Kiến trúc hiện tại"
        - src: "target_architecture.png"
          alt: "Kiến trúc đề xuất"

      notes:
        Nếu file ảnh chưa tồn tại, renderer PPTX sẽ hiển thị placeholder missing image.

    slide "6 năng lực chính":
      layout: grid-6
      intent: summarize_capabilities

      title:
        Sáu năng lực chính của hệ thống

      cells:
        - heading: Thu thập
          text: Lấy dữ liệu từ nhiều nguồn nội bộ và bên ngoài.
        - heading: Làm sạch
          text: Chuẩn hóa định dạng, khử trùng lặp và lọc nhiễu.
        - heading: Truy xuất
          text: Tìm các đoạn liên quan theo ngữ nghĩa và cấu trúc.
        - heading: Tổng hợp
          text: Kết hợp thông tin thành bản tóm tắt có kiểm soát.
        - heading: Trực quan
          text: Sinh dashboard, biểu đồ hoặc slide trình bày.
        - heading: Giám sát
          text: Theo dõi lỗi, chất lượng đầu ra và độ tin cậy.

      notes:
        Layout grid-6 dùng 3 cột x 2 hàng.


    slide "Lộ trình triển khai":
      layout: stair-progress

      title:
        Lộ trình triển khai theo bậc thang

      subtitle:
        Mỗi bậc thể hiện một mức trưởng thành của hệ thống

      steps:
        - heading: Khởi tạo
          text: Xác định phạm vi, nguồn dữ liệu và mục tiêu đầu ra.
        - heading: Chuẩn hóa
          text: Làm sạch, hợp nhất và kiểm soát chất lượng dữ liệu.
        - heading: Truy xuất
          text: Xây dựng chỉ mục, embedding và graph context.
        - heading: Sinh nội dung
          text: Tổng hợp báo cáo, dashboard và slide trình bày.
        - heading: Giám sát
          text: Theo dõi lỗi, feedback và cải tiến liên tục.

    slide "Các tầng năng lực xếp chồng":
      layout: stacked-stairs
      intent: show_maturity_levels

      title:
        Các tầng năng lực xếp chồng

      subtitle:
        Các bậc nhỏ dần và căn trái để thể hiện mức độ thu hẹp / trưởng thành

      steps:
        - heading: Nền tảng dữ liệu
          text: Thu thập, chuẩn hóa và quản trị nguồn dữ liệu.
        - heading: Truy xuất tri thức
          text: Tạo chỉ mục, metadata và quan hệ ngữ cảnh.
        - heading: Suy luận bằng LLM
          text: Tổng hợp, kiểm chứng và sinh nội dung có kiểm soát.
        - heading: Tự động hóa đầu ra
          text: Sinh dashboard, báo cáo và slide trình bày.
        - heading: Giám sát vận hành
          text: Theo dõi chất lượng, lỗi và feedback để cải tiến.

      footer_text:
        Stacked stairs demo | Internal Use Only

    slide "Ma trận 4x2 năng lực":
      layout: grid-4x2
      intent: show_capability_matrix

      title:
        Ma trận 4x2 năng lực hệ thống

      cells:
        - heading: Thu thập
          text: Kết nối nhiều nguồn dữ liệu.
        - heading: Làm sạch
          text: Chuẩn hóa, lọc nhiễu, khử trùng lặp.
        - heading: Lập chỉ mục
          text: Tạo embedding và metadata truy xuất.
        - heading: Truy xuất
          text: Chọn ngữ cảnh liên quan theo truy vấn.
        - heading: Suy luận
          text: Kết hợp bằng LLM và luật nghiệp vụ.
        - heading: Kiểm chứng
          text: Soát nguồn, phát hiện thiếu nhất quán.
        - heading: Trình bày
          text: Sinh báo cáo, dashboard hoặc slide.
        - heading: Giám sát
          text: Theo dõi lỗi, độ trễ và chất lượng.

      footer:
        text: Grid 4x2 demo | Internal Use Only
        align: center



    slide "Timeline triển khai":
      layout: timeline
      intent: show_roadmap

      title:
        Timeline triển khai theo các mốc chính

      milestones:
        - date: Q1
          heading: Khởi tạo
          text: Xác định phạm vi, nguồn dữ liệu và tiêu chí thành công.
        - date: Q2
          heading: Chuẩn hóa
          text: Làm sạch dữ liệu, thống nhất schema và metadata.
        - date: Q3
          heading: Tích hợp
          text: Kết nối retrieval, reasoning và dashboard vận hành.
        - date: Q4
          heading: Mở rộng
          text: Tối ưu chi phí, giám sát chất lượng và nhân rộng.

    slide "Icon bullets":
      layout: title-bullets
      intent: show_icon_bullets

      title:
        Bullet đặc biệt bằng icon

      bullets:
        icon: check
        items:
          - Chuẩn hóa dữ liệu đầu vào
          - Tạo metadata và chỉ mục
          - Sinh báo cáo tự động



    slide "Danh sách dài tự ngắt":
      layout: title-bullets
      intent: show_auto_pagination

      title:
        Danh sách dài tự ngắt sang slide sau

      bullets:
        icon: check
        overflow: paginate
        max_items_per_slide: 5
        items:
          - Thu thập dữ liệu từ nhiều nguồn nội bộ
          - Chuẩn hóa định dạng đầu vào
          - Loại bỏ bản ghi trùng lặp
          - Gắn metadata cho từng tài liệu
          - Tạo embedding và chỉ mục truy xuất
          - Mở rộng ngữ cảnh bằng đồ thị quan hệ
          - Sinh tóm tắt theo chuẩn báo cáo
          - Kiểm tra chất lượng đầu ra
          - Xuất dashboard và slide trình bày
          - Lưu log phục vụ kiểm toán

    slide "Per-item icon bullets":
      layout: title-bullets
      intent: show_per_item_icons

      title:
        Mỗi bullet có icon riêng

      bullets:
        - icon: database
          text: Thu thập dữ liệu từ nhiều nguồn
        - icon: gear
          text: Xử lý và chuẩn hóa
        - icon: star
          text: Tạo insight cho lãnh đạo
'''


DEMO_PCL = '''
pcl "safe-slide-constraints":
  default:
    text:
      overflow: shrink
      min_size: 13
      max_lines: 6
    bullets:
      overflow: paginate
      min_size: 14
      max_lines: 7
    table:
      overflow: paginate
      max_rows_per_slide: 8
    image:
      mode: contain
    background_image:
      mode: cover

  layout "title-bullets":
    bullets:
      overflow: paginate
      max_items_per_slide: 5

  layout "numbered-columns-4":
    card:
      keep_inside: true
      overflow: shrink
      min_size: 11

  slide "Danh sách dài tự ngắt theo constraint":
    bullets:
      overflow: paginate
      max_items_per_slide: 4
'''


DEMO_PSL = '''
theme "corporate":
  page:
    size: widescreen
    background: "#FFFFFF"

  colors:
    primary: "#003A8C"
    secondary: "#E6F0FF"
    text: "#1F1F1F"
    muted: "#666666"
    white: "#FFFFFF"

  fonts:
    heading: "Aptos Display"
    body: "Aptos"

  presentation.title-slide:
    background: primary
    background_image:
      src: "cover_bg.png"
      opacity: 0.20
      mode: cover

    title:
      font: heading
      size: 48
      color: white
      bold: true
      align: center
      position: [90, 220]
      width: 1100
      height: 90

    subtitle:
      font: body
      size: 26
      color: secondary
      italic: true
      align: center
      position: [95, 320]
      width: 1000
      height: 60

    author:
      font: body
      size: 20
      color: white
      position: [95, 430]
      width: 800
      height: 40

    date:
      font: body
      size: 18
      color: secondary
      position: [95, 470]
      width: 400
      height: 40

  section.section-header:
    background: secondary
    title:
      font: heading
      size: 44
      color: primary
      position: [90, 250]
      width: 1100
      height: 90

    subtitle:
      font: body
      size: 24
      color: text
      position: [95, 340]
      width: 1000
      height: 60

  slide.title-bullets:
    title:
      font: heading
      size: 36
      color: primary
      position: [60, 40]
      width: 1100
      height: 90

    bullets:
      font: body
      size: 24
      color: text
      position: [90, 170]
      width: 1000
      height: 420
      line_gap: 12
      overflow: shrink
      max_lines: 5
      min_size: 15

    link:
      font: body
      size: 18
      color: primary
      underline: true
      position: [90, 610]
      width: 700
      height: 36
      icon: arrow
      icon_color: primary
      icon_size: 24
      icon_gap: 12

    footer_text:
      font: body
      size: 14
      color: muted
      italic: true
      align: left
      position: [80, 685]
      width: 840
      height: 24

    footer_image:
      position: [1080, 665]
      width: 120
      height: 36

  slide.title-image:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 70

    subtitle:
      font: body
      size: 22
      color: muted
      position: [65, 105]
      width: 1080
      height: 50

    image:
      position: [170, 180]
      width: 940
      height: 450



  slide.text-image:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 155]
      width: 520
      height: 420
      heading_size: 25
      size: 21
      color: text

    right:
      position: [680, 155]
      width: 520
      height: 390
      image_width: 520
      image_height: 390

  slide.image-text:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 155]
      width: 520
      height: 390
      image_width: 520
      image_height: 390

    right:
      position: [680, 155]
      width: 520
      height: 420
      heading_size: 25
      size: 21
      color: text



  slide.hero-image:
    background: "#FFFFFF"

    title:
      font: heading
      size: 34
      color: primary
      position: [60, 36]
      width: 1120
      height: 70

    subtitle:
      font: body
      size: 19
      color: muted
      position: [62, 102]
      width: 1080
      height: 40

    image:
      position: [120, 155]
      width: 1040
      height: 410

    caption:
      font: body
      size: 17
      color: muted
      position: [140, 585]
      width: 1000
      height: 60

  slide.image-caption:
    background: "#FFFFFF"

    image:
      position: [120, 70]
      width: 1040
      height: 520

    caption:
      font: body
      size: 18
      color: muted
      position: [140, 610]
      width: 1000
      height: 60

    footer_text:
      font: body
      size: 14
      color: muted
      position: [140, 685]
      width: 1000
      height: 24


  slide.two-images:
    background: "#FFFFFF"

    title:
      font: heading
      size: 32
      color: primary
      position: [60, 36]
      width: 1120
      height: 70

    subtitle:
      font: body
      size: 19
      color: muted
      position: [62, 104]
      width: 1080
      height: 40

    left:
      position: [70, 165]
      width: 540
      height: 390
      image_width: 540
      image_height: 390
      caption_size: 16
      caption_color: muted

    right:
      position: [670, 165]
      width: 540
      height: 390
      image_width: 540
      image_height: 390
      caption_size: 16
      caption_color: muted

  slide.grid-3:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 155]
      width: 1140
      height: 430
      columns: 3
      gap: 24
      padding: 18
      fill: white
      border: secondary
      heading_size: 21
      size: 16
      color: text

  slide.grid-4:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 145]
      width: 1140
      height: 470
      columns: 2
      gap: 24
      padding: 18
      fill: white
      border: secondary
      heading_size: 21
      size: 16
      color: text

  slide.grid-5:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 145]
      width: 1140
      height: 470
      columns: 3
      gap: 22
      padding: 16
      fill: white
      border: secondary
      heading_size: 20
      size: 15
      color: text

  slide.grid-6:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 145]
      width: 1140
      height: 470
      columns: 3
      gap: 22
      padding: 16
      fill: white
      border: secondary
      heading_size: 20
      size: 15
      color: text


  slide.stair-progress:
    title:
      font: heading
      size: 34
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 60

    subtitle:
      font: body
      size: 18
      color: muted
      italic: true
      position: [65, 100]
      width: 1050
      height: 40

    stair:
      position: [80, 175]
      step_width: 255
      step_height: 118
      x_step: 195
      y_step: 58
      padding: 15
      fill: secondary
      border: primary
      number_color: primary
      heading_color: text
      color: muted
      heading_size: 18
      size: 14
      align: left

  slide.stacked-stairs:
    title:
      font: heading
      size: 34
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 60
      align: left

    subtitle:
      font: body
      size: 18
      color: muted
      italic: true
      position: [65, 100]
      width: 1050
      height: 40

    stacked_stairs:
      position: [120, 175]
      base_width: 980
      step_height: 90
      shrink: 95
      overlap: 24
      align_side: left
      padding_x: 22
      padding_y: 13
      fill: secondary
      border: primary
      number_color: primary
      heading_color: text
      color: muted
      heading_size: 18
      size: 14
      align: left

  slide.grid-4x2:
    title:
      font: heading
      size: 32
      color: primary
      position: [50, 35]
      width: 1180
      height: 70
      bold: true
      align: center

    grid:
      position: [50, 135]
      width: 1180
      height: 500
      columns: 4
      rows: 2
      gap: 18
      padding: 14
      fill: white
      border: secondary
      heading_font: heading
      heading_size: 17
      heading_color: primary
      size: 13
      color: text
      align: left

  slide.numbered-columns-4:
    title:
      font: heading
      size: 32
      color: primary
      bold: true
      align: center
      position: [60, 40]
      width: 1120
      height: 70

    numbered_columns:
      position: [70, 185]
      width: 1140
      height: 360
      gap: 24
      columns: 4
      circle_size: 58
      circle_offset_y: -32
      padding: 16
      fill: "#FFFFFF"
      border: "#CBD5E1"
      circle_fill: primary
      circle_color: "#FFFFFF"
      heading_size: 19
      heading_height: 40
      heading_align: center
      bullet_size: 13
      bullet_align: left
      bullet_icon: dot
      icon_color: primary
      line_gap: 4
      overflow: shrink
      max_lines: 5
      min_size: 10


  slide.title-table:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 70
      bold: true

    table:
      font: body
      size: 17
      color: text
      position: [70, 140]
      width: 1140
      height: 470
      header_fill: primary
      header_color: white
      cell_fill: "#FFFFFF"
      alt_fill: "#F8FAFC"
      border_color: "#CBD5E1"
      align: left
      header_align: center

  slide.two-column:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 160]
      width: 520

    right:
      position: [680, 160]
      width: 520


  slide.timeline:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 70
      bold: true
      align: center

    timeline:
      position: [90, 155]
      width: 1100
      height: 470
      axis_y: 365
      line_height: 5
      marker_size: 24
      connector_height: 58
      card_width: 230
      card_height: 118
      gap: 28
      alternate: true
      fill: white
      border: secondary
      line_color: primary
      marker_fill: primary
      marker_border: white
      date_color: primary
      heading_color: text
      color: muted
      date_size: 13
      heading_size: 17
      size: 13
'''


def resolve_relative_path(base_dir: Path, maybe_path: Optional[str]) -> Optional[Path]:
    """Resolve a file reference from PML relative to the script/current base dir."""
    if not maybe_path:
        return None
    path = Path(str(maybe_path))
    if not path.is_absolute():
        path = base_dir / path
    return path


def load_theme_for_doc(doc: PresentationDoc, base_dir: Path, fallback_text: str = DEMO_PSL) -> Theme:
    """Load PSL through `use style:`; fall back to embedded demo PSL if absent."""
    style_path = resolve_relative_path(base_dir, doc.style_file)
    if style_path and style_path.exists():
        return parse_psl(style_path.read_text(encoding="utf-8"))
    return parse_psl(fallback_text)


def load_constraints_for_doc(
    doc: PresentationDoc, base_dir: Path, fallback_text: Optional[str] = DEMO_PCL
) -> Optional[ConstraintSet]:
    """Load PCL through `use constraints:`; fall back to embedded demo PCL if absent."""
    constraint_path = resolve_relative_path(base_dir, doc.constraint_file)
    if constraint_path and constraint_path.exists():
        return parse_pcl(constraint_path.read_text(encoding="utf-8"))
    if fallback_text:
        return parse_pcl(fallback_text)
    return None


def write_demo_external_files(base_dir: Path) -> None:
    """Materialize files referenced by DEMO_PML so `use style` and `use constraints` work."""
    (base_dir / "corporate.psl").write_text(DEMO_PSL, encoding="utf-8")
    (base_dir / "safe-layouts.pcl").write_text(DEMO_PCL, encoding="utf-8")


def main() -> None:
    script_dir = Path(__file__).resolve().parent
    create_demo_assets(script_dir)
    write_demo_external_files(script_dir)

    doc = parse_pml(DEMO_PML)
    theme = load_theme_for_doc(doc, script_dir)
    constraints = load_constraints_for_doc(doc, script_dir)
    render_ir = build_render_ir(doc, theme, constraints)

    render_html(render_ir, "demo_output.html", asset_base_dir=script_dir)
    render_pptx(render_ir, "demo_output.pptx", asset_base_dir=script_dir)

    print(f"Style file: {doc.style_file}")
    print(f"Constraint file: {doc.constraint_file}")
    print("Generated demo_output.html")
    print("Generated demo_output.pptx")


if __name__ == "__main__":
    main()
