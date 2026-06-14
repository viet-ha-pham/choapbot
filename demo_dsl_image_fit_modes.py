"""
PML/PSL Renderer v1.0
=====================

Adds support for:
- presentation title slide
- section header slide
- content slides
- title + image slide layout
- text + image two-column layouts
- grid layouts: grid-3, grid-4, grid-5, grid-6
- two image layout: two-images
- hero image layout: hero-image (large image + caption)
- image caption layout: image-caption (image + caption only, no title/subtitle)
- footer images: footer_image / footer_images on content slides
- footer text: footer_text or footer.text on content slides
- background images: background_image or background.image on slides/layouts
- image fit modes: stretch, contain, cover/crop, original

Install:
    pip install python-pptx

Run:
    python demo_dsl_background_image.py

Outputs:
    demo_output.html
    demo_output.pptx
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import html
import re

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    PPTXPresentation = None
    Inches = Pt = RGBColor = None


# -----------------------------------------------------------------------------
# Data models
# -----------------------------------------------------------------------------

@dataclass
class Slide:
    title: str
    layout: str = "title-bullets"
    intent: Optional[str] = None
    blocks: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Section:
    title: str
    header_layout: str = "section-header"
    header: Dict[str, Any] = field(default_factory=dict)
    slides: List[Slide] = field(default_factory=list)


@dataclass
class PresentationDoc:
    title: str
    meta: Dict[str, Any] = field(default_factory=dict)
    style_file: Optional[str] = None
    cover_layout: str = "title-slide"
    cover: Dict[str, Any] = field(default_factory=dict)
    sections: List[Section] = field(default_factory=list)


@dataclass
class Theme:
    name: str
    page: Dict[str, Any] = field(default_factory=dict)
    colors: Dict[str, str] = field(default_factory=dict)
    fonts: Dict[str, str] = field(default_factory=dict)
    presentation_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    section_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)


# -----------------------------------------------------------------------------
# Low-level indentation parser
# -----------------------------------------------------------------------------

@dataclass
class Line:
    indent: int
    text: str
    lineno: int


def strip_comment(line: str) -> str:
    """Remove comments outside quotes."""
    in_quote = False
    quote_char = ""
    out = []

    for ch in line:
        if ch in ('"', "'"):
            if not in_quote:
                in_quote = True
                quote_char = ch
            elif quote_char == ch:
                in_quote = False
        if ch == "#" and not in_quote:
            break
        out.append(ch)

    return "".join(out).rstrip()


def tokenize_lines(text: str) -> List[Line]:
    lines: List[Line] = []
    for i, raw in enumerate(text.splitlines(), start=1):
        no_comment = strip_comment(raw)
        if not no_comment.strip():
            continue
        indent = len(no_comment) - len(no_comment.lstrip(" "))
        if "\t" in raw[:indent]:
            raise SyntaxError(f"Line {i}: tabs are not allowed for indentation")
        lines.append(Line(indent=indent, text=no_comment.strip(), lineno=i))
    return lines


def parse_scalar(value: str) -> Any:
    value = value.strip()

    if value == "":
        return ""

    if (value.startswith('"') and value.endswith('"')) or (
        value.startswith("'") and value.endswith("'")
    ):
        return value[1:-1]

    if value.lower() == "true":
        return True
    if value.lower() == "false":
        return False
    if value.lower() in {"null", "none"}:
        return None

    if value.startswith("[") and value.endswith("]"):
        inner = value[1:-1].strip()
        if not inner:
            return []
        return [parse_scalar(part.strip()) for part in inner.split(",")]

    try:
        if "." in value:
            return float(value)
        return int(value)
    except ValueError:
        return value


def split_key_value(text: str) -> Tuple[str, Optional[str]]:
    if ":" not in text:
        return text, None
    key, val = text.split(":", 1)
    return key.strip(), val.strip()


def parse_generic_block(lines: List[Line], start: int, indent: int) -> Tuple[Any, int]:
    result: Dict[str, Any] = {}
    items: List[Any] = []
    mode: Optional[str] = None
    i = start

    while i < len(lines):
        line = lines[i]
        if line.indent < indent:
            break
        if line.indent > indent:
            raise SyntaxError(
                f"Line {line.lineno}: unexpected indentation. Expected {indent}, got {line.indent}"
            )

        text = line.text

        if text.startswith("- "):
            if mode is None:
                mode = "list"
            elif mode != "list":
                raise SyntaxError(f"Line {line.lineno}: cannot mix list and mapping")

            item_text = text[2:].strip()

            # Supports list scalars:
            #   - item
            # Supports list mappings:
            #   - heading: A
            #     text: B
            # Supports nested named mappings:
            #   - image:
            #       src: a.png
            if item_text.endswith(":"):
                key = item_text[:-1].strip()
                child, new_i = parse_generic_block(lines, i + 1, indent + 2)
                items.append({key: child})
                i = new_i
            elif ":" in item_text:
                key, val = split_key_value(item_text)
                item: Dict[str, Any] = {key: parse_scalar(val or "")}
                new_i = i + 1
                if new_i < len(lines) and lines[new_i].indent > indent:
                    child, new_i = parse_generic_block(lines, new_i, lines[new_i].indent)
                    if isinstance(child, dict):
                        item.update(child)
                    else:
                        item["children"] = child
                items.append(item)
                i = new_i
            else:
                items.append(parse_scalar(item_text))
                i += 1
            continue

        if mode is None:
            mode = "dict"
        elif mode != "dict":
            raise SyntaxError(f"Line {line.lineno}: cannot mix mapping and list")

        key, val = split_key_value(text)
        if val is None:
            text_lines = []
            while i < len(lines) and lines[i].indent == indent:
                text_lines.append(lines[i].text)
                i += 1
            return "\n".join(text_lines), i

        if val != "":
            result[key] = parse_scalar(val)
            i += 1
        else:
            if i + 1 >= len(lines) or lines[i + 1].indent <= indent:
                result[key] = {}
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                result[key] = child
                i = new_i

    if mode == "list":
        return items, i
    return result, i


# -----------------------------------------------------------------------------
# PML / PSL parser
# -----------------------------------------------------------------------------

PRESENTATION_RE = re.compile(r'^presentation\s+"(.+)"\s*:$')
SECTION_RE = re.compile(r'^section\s+"(.+)"\s*:$')
SLIDE_RE = re.compile(r'^slide\s+"(.+)"\s*:$')
THEME_RE = re.compile(r'^theme\s+"(.+)"\s*:$')


def parse_pml(text: str) -> PresentationDoc:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PML")

    m = PRESENTATION_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PML must start with: presentation "Title":')

    doc = PresentationDoc(title=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        text = line.text
        key, val = split_key_value(text)

        if text == "meta:":
            block, i = parse_generic_block(lines, i + 1, 4)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: meta must be a mapping")
            doc.meta = block
            continue

        if key == "cover_layout":
            doc.cover_layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if text == "cover:":
            block, i = parse_generic_block(lines, i + 1, 4)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: cover must be a mapping")
            doc.cover = block
            continue

        if text.startswith("use style:"):
            _, value = split_key_value(text)
            doc.style_file = parse_scalar(value or "")
            i += 1
            continue

        sec_match = SECTION_RE.match(text)
        if sec_match:
            section, i = parse_section(lines, i)
            doc.sections.append(section)
            continue

        raise SyntaxError(f"Line {line.lineno}: unknown PML statement: {text}")

    return doc


def parse_section(lines: List[Line], start: int) -> Tuple[Section, int]:
    m = SECTION_RE.match(lines[start].text)
    assert m is not None
    section = Section(title=m.group(1))
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= lines[start].indent:
            break
        if line.indent != lines[start].indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected section child indent")

        key, val = split_key_value(line.text)

        if key == "header_layout":
            section.header_layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if line.text == "header:":
            block, i = parse_generic_block(lines, i + 1, line.indent + 2)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: section header must be a mapping")
            section.header = block
            continue

        slide_match = SLIDE_RE.match(line.text)
        if slide_match:
            slide, i = parse_slide(lines, i)
            section.slides.append(slide)
        else:
            raise SyntaxError(f"Line {line.lineno}: section only accepts header/header_layout/slide blocks")

    return section, i


def parse_slide(lines: List[Line], start: int) -> Tuple[Slide, int]:
    m = SLIDE_RE.match(lines[start].text)
    assert m is not None
    slide = Slide(title=m.group(1))
    base_indent = lines[start].indent
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= base_indent:
            break
        if line.indent != base_indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected slide child indent")

        key, val = split_key_value(line.text)

        if key == "layout":
            slide.layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if key == "intent":
            slide.intent = str(parse_scalar(val or ""))
            i += 1
            continue

        if val == "":
            if i + 1 >= len(lines) or lines[i + 1].indent <= line.indent:
                slide.blocks[key] = ""
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                slide.blocks[key] = child
                i = new_i
        else:
            slide.blocks[key] = parse_scalar(val or "")
            i += 1

    return slide, i


def parse_psl(text: str) -> Theme:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PSL")

    m = THEME_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PSL must start with: theme "Name":')

    theme = Theme(name=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        key, val = split_key_value(line.text)
        if val != "":
            raise SyntaxError(f"Line {line.lineno}: top-level PSL keys must be blocks")

        block, i = parse_generic_block(lines, i + 1, 4)
        if not isinstance(block, dict):
            raise SyntaxError(f"Line {line.lineno}: top-level PSL block must be a mapping")

        if key == "page":
            theme.page = block
        elif key == "colors":
            theme.colors = block
        elif key == "fonts":
            theme.fonts = block
        elif key.startswith("presentation."):
            name = key.split(".", 1)[1]
            theme.presentation_layouts[name] = block
        elif key.startswith("section."):
            name = key.split(".", 1)[1]
            theme.section_layouts[name] = block
        elif key.startswith("slide."):
            name = key.split(".", 1)[1]
            theme.layouts[name] = block
        else:
            theme.layouts[key] = block

    return theme


# -----------------------------------------------------------------------------
# Render IR construction
# -----------------------------------------------------------------------------

def resolve_token(value: Any, theme: Theme) -> Any:
    if isinstance(value, str):
        if value in theme.colors:
            return theme.colors[value]
        if value in theme.fonts:
            return theme.fonts[value]
    return value


def style_value(style: Dict[str, Any], key: str, theme: Theme, default: Any = None) -> Any:
    return resolve_token(style.get(key, default), theme)


def normalize_text_block(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, dict):
        if "text" in value:
            return str(value["text"])
        if "heading" in value:
            return str(value["heading"])
    return str(value)


def make_text_object(kind: str, value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [60, 40])
    return {
        "type": "TextBox",
        "role": kind,
        "text": normalize_text_block(value),
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 80),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 28),
        "color": style_value(style, "color", theme, "#000000"),
    }


def make_bullet_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [90, 160])
    if not isinstance(value, list):
        value = [value]
    return {
        "type": "BulletList",
        "items": [str(v) for v in value],
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 420),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 24),
        "color": style_value(style, "color", theme, "#000000"),
        "line_gap": style.get("line_gap", 8),
    }


def normalize_fit_mode(value: Any, default: str = "stretch") -> str:
    """Normalize image placement modes.

    Modes:
      stretch  - resize image to exactly fill the box; aspect ratio may change.
      contain  - fit entire image inside the box; aspect ratio is preserved.
      cover    - fill the box and crop overflow; aspect ratio is preserved.
      crop     - alias of cover.
      original - use original pixel size, centered in the box when possible.
    """
    mode = str(value or default).strip().lower()
    aliases = {
        "resize": "stretch",
        "fill": "stretch",
        "fit": "contain",
        "scale": "contain",
        "crop": "cover",
        "none": "original",
    }
    mode = aliases.get(mode, mode)
    if mode not in {"stretch", "contain", "cover", "original"}:
        return default
    return mode


def make_image_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    """
    Supported PML forms:

      image: "path/to/image.png"

      image:
        src: "path/to/image.png"
        alt: "description"
        mode: contain   # stretch | contain | cover/crop | original

    The image style in PSL may also define `mode`. PML overrides PSL.
    """
    pos = style.get("position", [700, 150])

    if isinstance(value, dict):
        src = str(value.get("src", value.get("path", "")))
        alt = str(value.get("alt", ""))
        mode = normalize_fit_mode(value.get("mode", style.get("mode", "stretch")))
    else:
        src = str(value)
        alt = ""
        mode = normalize_fit_mode(style.get("mode", "stretch"))

    return {
        "type": "Image",
        "src": src,
        "alt": alt,
        "mode": mode,
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 500),
        "h": style.get("height", 320),
    }



def normalize_background_image(value: Any) -> Dict[str, Any]:
    """Normalize background image declarations.

    Supported forms:

      background_image: "bg.png"

      background_image:
        src: "bg.png"
        opacity: 0.25
        mode: cover

      background:
        image: "bg.png"
        opacity: 0.25
        mode: cover
    """
    if value is None or value == "":
        return {}
    if isinstance(value, str):
        return {"src": value}
    if isinstance(value, dict):
        if "image" in value and "src" not in value:
            return {
                "src": str(value.get("image", "")),
                "opacity": value.get("opacity", 1.0),
                "mode": normalize_fit_mode(value.get("mode", "cover"), default="cover"),
            }
        return {
            "src": str(value.get("src", value.get("path", ""))),
            "opacity": value.get("opacity", 1.0),
            "mode": normalize_fit_mode(value.get("mode", "cover"), default="cover"),
        }
    return {"src": str(value)}


def make_background_image_spec(blocks: Dict[str, Any], layout_style: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Resolve background image from PML first, then PSL layout defaults.

    Priority:
      1. slide/cover/header block: background_image
      2. slide/cover/header block: background.image
      3. layout style: background_image
      4. layout style: background.image
    """
    candidates: List[Any] = []

    if isinstance(blocks, dict):
        if "background_image" in blocks:
            candidates.append(blocks.get("background_image"))
        bg_block = blocks.get("background")
        if isinstance(bg_block, dict) and "image" in bg_block:
            candidates.append(bg_block)

    if isinstance(layout_style, dict):
        if "background_image" in layout_style:
            candidates.append(layout_style.get("background_image"))
        bg_style = layout_style.get("background")
        if isinstance(bg_style, dict) and "image" in bg_style:
            candidates.append(bg_style)

    for candidate in candidates:
        spec = normalize_background_image(candidate)
        if spec.get("src"):
            spec.setdefault("opacity", 1.0)
            spec["mode"] = normalize_fit_mode(spec.get("mode", "cover"), default="cover")
            return spec

    return None

def make_footer_image_objects(slide_blocks: Dict[str, Any], layout_style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """Build footer image objects for a slide.

    Supported PML forms:

      footer_image:
        src: "logo.png"
        alt: "Logo"

      footer_images:
        - src: "logo_a.png"
          alt: "Logo A"
        - src: "logo_b.png"
          alt: "Logo B"

    Style can be declared per slide layout:

      slide.title-bullets:
        footer_image:
          position: [1080, 660]
          width: 120
          height: 36

        footer_images:
          position: [900, 665]
          width: 80
          height: 30
          gap: 16
    """
    objects: List[Dict[str, Any]] = []

    if "footer_image" in slide_blocks:
        style = layout_style.get("footer_image", layout_style.get("footer", {}))
        default_style = {
            "position": style.get("position", [1080, 660]),
            "width": style.get("width", 120),
            "height": style.get("height", 36),
        }
        obj = make_image_object(slide_blocks["footer_image"], default_style, theme)
        obj["role"] = "footer_image"
        objects.append(obj)

    if "footer_images" in slide_blocks:
        images = slide_blocks["footer_images"]
        if not isinstance(images, list):
            images = [images]

        style = layout_style.get("footer_images", layout_style.get("footer", {}))
        pos = style.get("position", [900, 665])
        width = style.get("width", 80)
        height = style.get("height", 30)
        gap = style.get("gap", 16)

        for idx, image in enumerate(images):
            item_style = {
                "position": [pos[0] + idx * (width + gap), pos[1]],
                "width": width,
                "height": height,
            }
            obj = make_image_object(normalize_image_item(image), item_style, theme)
            obj["role"] = "footer_image"
            objects.append(obj)

    return objects



def make_footer_text_objects(slide_blocks: Dict[str, Any], layout_style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """Build footer text objects for a slide.

    Supported PML forms:

      footer_text:
        Nguồn: Viettel AI Lab | Internal Use Only

      footer_text: "Internal Use Only"

      footer:
        text: "Confidential"
        align: center

    Supported PSL forms:

      slide.title-bullets:
        footer_text:
          position: [80, 685]
          width: 900
          height: 24
          font: body
          size: 14
          color: muted

        footer:
          text_position: [80, 685]
          text_width: 900
          text_height: 24
          font: body
          size: 14
          color: muted
    """
    footer_text: Any = None
    footer_block = slide_blocks.get("footer")

    if "footer_text" in slide_blocks:
        footer_text = slide_blocks["footer_text"]
    elif isinstance(footer_block, dict) and "text" in footer_block:
        footer_text = footer_block.get("text")

    if footer_text is None or footer_text == "":
        return []

    text = normalize_text_block(footer_text)
    footer_style = layout_style.get("footer", {})
    style = dict(layout_style.get("footer_text", {}))

    # Allow a generic footer block in PSL to provide defaults.
    if "position" not in style:
        style["position"] = footer_style.get("text_position", footer_style.get("position", [80, 685]))
    if "width" not in style:
        style["width"] = footer_style.get("text_width", 900)
    if "height" not in style:
        style["height"] = footer_style.get("text_height", 24)
    if "font" not in style:
        style["font"] = footer_style.get("font", "body")
    if "size" not in style:
        style["size"] = footer_style.get("size", 14)
    if "color" not in style:
        style["color"] = footer_style.get("color", "muted")

    # Optional quick alignment override in PML. This is simple positioning on a 1280px canvas.
    align = None
    if isinstance(footer_block, dict):
        align = footer_block.get("align")
    if isinstance(align, str):
        width = style.get("width", 900)
        y = style.get("position", [80, 685])[1]
        if align.lower() == "center":
            style["position"] = [(1280 - width) / 2, y]
        elif align.lower() == "right":
            style["position"] = [1280 - width - 80, y]
        elif align.lower() == "left":
            style["position"] = [80, y]

    obj = make_text_object("footer_text", text, style, theme)
    obj["role"] = "footer_text"
    return [obj]


def make_column_objects(role: str, value: Any, style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """
    Build objects for a left/right column.

    Supported forms:

      left:
        heading: Main points
        bullets:
          - A
          - B

      right:
        image:
          src: "diagram.png"
          alt: "Diagram"

      left:
        text: Free paragraph text

    A column can combine heading + bullets + image, but the common use case is
    text in one column and image in the other column.
    """
    pos = style.get("position", [70, 150])
    width = style.get("width", 520)
    height = style.get("height", 360)
    objects: List[Dict[str, Any]] = []

    text_y = pos[1]

    if isinstance(value, dict):
        if "heading" in value:
            objects.append({
                "type": "TextBox",
                "role": f"{role}.heading",
                "text": str(value["heading"]),
                "x": pos[0],
                "y": text_y,
                "w": width,
                "h": style.get("heading_height", 50),
                "font": style_value(style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": style.get("heading_size", 24),
                "color": style_value(style, "heading_color", theme, theme.colors.get("primary", "#000000")),
            })
            text_y += style.get("heading_height", 50) + style.get("heading_gap", 10)

        if "text" in value:
            objects.append({
                "type": "TextBox",
                "role": f"{role}.text",
                "text": str(value["text"]),
                "x": pos[0],
                "y": text_y,
                "w": width,
                "h": style.get("text_height", 220),
                "font": style_value(style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": style.get("size", 20),
                "color": style_value(style, "color", theme, theme.colors.get("text", "#000000")),
            })
            text_y += style.get("text_height", 220) + style.get("text_gap", 12)

        if "bullets" in value:
            objects.append({
                "type": "BulletList",
                "items": [str(x) for x in value["bullets"]],
                "x": pos[0],
                "y": text_y,
                "w": width,
                "h": style.get("bullets_height", 300),
                "font": style_value(style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": style.get("size", 20),
                "color": style_value(style, "color", theme, theme.colors.get("text", "#000000")),
                "line_gap": style.get("line_gap", 6),
            })

        if "image" in value:
            image_style = dict(style)
            image_style["position"] = style.get("image_position", pos)
            image_style["width"] = style.get("image_width", width)
            image_style["height"] = style.get("image_height", height)
            objects.append(make_image_object(value["image"], image_style, theme))

    else:
        objects.append({
            "type": "TextBox",
            "role": role,
            "text": str(value),
            "x": pos[0],
            "y": pos[1],
            "w": width,
            "h": height,
            "font": style_value(style, "font", theme, theme.fonts.get("body", "Aptos")),
            "size": style.get("size", 20),
            "color": style_value(style, "color", theme, theme.colors.get("text", "#000000")),
        })

    return objects


def normalize_image_item(value: Any) -> Dict[str, Any]:
    """Normalize an image entry for the two-images layout."""
    if isinstance(value, dict):
        if "image" in value and isinstance(value["image"], dict):
            return value["image"]
        return value
    return {"src": str(value), "alt": ""}


def make_two_image_objects(slide_blocks: Dict[str, Any], style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    """
    Build a native two-images layout.

    Supported PML forms:

      slide "Compare":
        layout: two-images
        images:
          - src: "before.png"
            alt: "Before"
          - src: "after.png"
            alt: "After"

      slide "Compare":
        layout: two-images
        left:
          image:
            src: "before.png"
        right:
          image:
            src: "after.png"
    """
    objects: List[Dict[str, Any]] = []

    if "images" in slide_blocks:
        images = slide_blocks["images"]
        if not isinstance(images, list):
            images = [images]
        images = images[:2]
        side_names = ["left", "right"]
        for idx, img in enumerate(images):
            side = side_names[idx]
            side_style = style.get(side, {})
            image_style = dict(side_style)
            image_style["position"] = side_style.get("image_position", side_style.get("position", [70 + idx * 610, 160]))
            image_style["width"] = side_style.get("image_width", side_style.get("width", 540))
            image_style["height"] = side_style.get("image_height", side_style.get("height", 380))
            objects.append(make_image_object(normalize_image_item(img), image_style, theme))
        return objects

    for side in ["left", "right"]:
        if side not in slide_blocks:
            continue
        value = slide_blocks[side]
        if isinstance(value, dict) and "image" in value:
            side_style = style.get(side, {})
            image_style = dict(side_style)
            image_style["position"] = side_style.get("image_position", side_style.get("position", [70, 160]))
            image_style["width"] = side_style.get("image_width", side_style.get("width", 540))
            image_style["height"] = side_style.get("image_height", side_style.get("height", 380))
            objects.append(make_image_object(value["image"], image_style, theme))

            if "caption" in value:
                cap_style = dict(side_style)
                cap_style["position"] = side_style.get("caption_position", [image_style["position"][0], image_style["position"][1] + image_style["height"] + 14])
                cap_style["width"] = side_style.get("caption_width", image_style["width"])
                cap_style["height"] = side_style.get("caption_height", 40)
                cap_style["size"] = side_style.get("caption_size", 16)
                cap_style["color"] = side_style.get("caption_color", theme.colors.get("muted", "#666666"))
                objects.append(make_text_object(f"{side}.caption", value["caption"], cap_style, theme))

    return objects


def normalize_cell(value: Any) -> Dict[str, Any]:
    """Normalize a grid cell to a dict with optional heading/text/bullets/image."""
    if isinstance(value, dict):
        if len(value) == 1 and "cell" in value and isinstance(value["cell"], dict):
            return value["cell"]
        return value
    return {"text": str(value)}


def make_grid_objects(value: Any, style: Dict[str, Any], theme: Theme, layout_name: str = "grid") -> List[Dict[str, Any]]:
    """
    Build card-like objects for grid layouts.

    Supported PML:
      cells:
        - heading: Thu thập
          text: Lấy dữ liệu từ nhiều nguồn
        - heading: Làm sạch
          bullets:
            - Chuẩn hóa
            - Khử trùng lặp
        - heading: Sơ đồ
          image:
            src: diagram.png
            alt: Diagram

    Supported PSL:
      slide.grid-6:
        grid:
          position: [70, 150]
          width: 1140
          height: 500
          columns: 3
          gap: 24
    """
    cells = value if isinstance(value, list) else [value]
    cells = [normalize_cell(cell) for cell in cells]

    grid_style = style.get("grid", style)
    pos = grid_style.get("position", [70, 150])
    total_w = grid_style.get("width", 1140)
    total_h = grid_style.get("height", 500)
    gap = grid_style.get("gap", 24)

    default_columns = 3 if layout_name in {"grid-5", "grid-6"} else 2
    columns = int(grid_style.get("columns", default_columns))
    columns = max(1, columns)
    rows = max(1, (len(cells) + columns - 1) // columns)

    cell_w = (total_w - gap * (columns - 1)) / columns
    cell_h = (total_h - gap * (rows - 1)) / rows

    objects: List[Dict[str, Any]] = []
    for idx, cell in enumerate(cells):
        row = idx // columns
        col = idx % columns
        x = pos[0] + col * (cell_w + gap)
        y = pos[1] + row * (cell_h + gap)

        objects.append({
            "type": "Card",
            "role": "grid.cell",
            "x": x,
            "y": y,
            "w": cell_w,
            "h": cell_h,
            "fill": style_value(grid_style, "fill", theme, "#FFFFFF"),
            "border": style_value(grid_style, "border", theme, "#D1D5DB"),
        })

        inner_pad = grid_style.get("padding", 18)
        text_x = x + inner_pad
        text_y = y + inner_pad
        text_w = cell_w - 2 * inner_pad

        if "heading" in cell:
            objects.append({
                "type": "TextBox",
                "role": "grid.heading",
                "text": str(cell["heading"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": grid_style.get("heading_height", 34),
                "font": style_value(grid_style, "heading_font", theme, theme.fonts.get("heading", "Aptos Display")),
                "size": grid_style.get("heading_size", 20),
                "color": style_value(grid_style, "heading_color", theme, theme.colors.get("primary", "#000000")),
            })
            text_y += grid_style.get("heading_height", 34) + grid_style.get("heading_gap", 8)

        if "text" in cell:
            objects.append({
                "type": "TextBox",
                "role": "grid.text",
                "text": str(cell["text"]),
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": grid_style.get("text_height", max(45, cell_h - (text_y - y) - inner_pad)),
                "font": style_value(grid_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": grid_style.get("size", 16),
                "color": style_value(grid_style, "color", theme, theme.colors.get("text", "#000000")),
            })
            text_y += grid_style.get("text_height", 70) + grid_style.get("text_gap", 8)

        if "bullets" in cell:
            bullets = cell["bullets"] if isinstance(cell["bullets"], list) else [cell["bullets"]]
            objects.append({
                "type": "BulletList",
                "items": [str(x) for x in bullets],
                "x": text_x,
                "y": text_y,
                "w": text_w,
                "h": grid_style.get("bullets_height", max(45, cell_h - (text_y - y) - inner_pad)),
                "font": style_value(grid_style, "font", theme, theme.fonts.get("body", "Aptos")),
                "size": grid_style.get("size", 15),
                "color": style_value(grid_style, "color", theme, theme.colors.get("text", "#000000")),
                "line_gap": grid_style.get("line_gap", 4),
            })
            text_y += grid_style.get("bullets_height", 80) + grid_style.get("text_gap", 8)

        if "image" in cell:
            image_style = dict(grid_style)
            image_style["position"] = [text_x, text_y]
            image_style["width"] = text_w
            image_style["height"] = max(50, cell_h - (text_y - y) - inner_pad)
            objects.append(make_image_object(cell["image"], image_style, theme))

    return objects


def build_text_objects_from_block(
    block: Dict[str, Any],
    layout_style: Dict[str, Any],
    theme: Theme,
    allowed_keys: Optional[List[str]] = None,
) -> List[Dict[str, Any]]:
    objects: List[Dict[str, Any]] = []
    keys = allowed_keys or list(block.keys())

    for key in keys:
        if key not in block:
            continue
        style = layout_style.get(key, {})
        objects.append(make_text_object(key, block[key], style, theme))

    return objects


def build_render_ir(doc: PresentationDoc, theme: Theme) -> Dict[str, Any]:
    slides_ir: List[Dict[str, Any]] = []

    # Presentation cover/title slide
    cover_style = theme.presentation_layouts.get(doc.cover_layout, {})
    cover_objects = [make_text_object("title", doc.title, cover_style.get("title", {}), theme)]
    cover_objects.extend(
        build_text_objects_from_block(doc.cover, cover_style, theme, ["subtitle", "author", "date"])
    )
    slides_ir.append({
        "type": "RenderSlide",
        "kind": "presentation",
        "section": None,
        "slide_title": doc.title,
        "layout": doc.cover_layout,
        "background": style_value(cover_style, "background", theme, theme.page.get("background", "#FFFFFF")),
        "background_image": make_background_image_spec(doc.cover, cover_style),
        "objects": cover_objects,
        "notes": doc.cover.get("notes", "") if isinstance(doc.cover, dict) else "",
    })

    for section in doc.sections:
        # Section header slide
        header_style = theme.section_layouts.get(section.header_layout, {})
        header_objects = [make_text_object("title", section.title, header_style.get("title", {}), theme)]
        header_objects.extend(
            build_text_objects_from_block(section.header, header_style, theme, ["subtitle", "summary", "label"])
        )
        slides_ir.append({
            "type": "RenderSlide",
            "kind": "section",
            "section": section.title,
            "slide_title": section.title,
            "layout": section.header_layout,
            "background": style_value(header_style, "background", theme, theme.page.get("background", "#FFFFFF")),
            "background_image": make_background_image_spec(section.header, header_style),
            "objects": header_objects,
            "notes": section.header.get("notes", "") if isinstance(section.header, dict) else "",
        })

        # Content slides
        for slide in section.slides:
            layout_style = theme.layouts.get(slide.layout, {})
            objects: List[Dict[str, Any]] = []

            # Native image-caption layout: image + caption only.
            # It intentionally ignores title/subtitle/body/bullets even if present.
            if slide.layout == "image-caption":
                if "image" in slide.blocks:
                    objects.append(make_image_object(slide.blocks["image"], layout_style.get("image", {}), theme))
                caption_text = slide.blocks.get("caption", "")
                if not caption_text and isinstance(slide.blocks.get("image"), dict):
                    caption_text = slide.blocks["image"].get("alt", "")
                if caption_text:
                    objects.append(make_text_object("caption", caption_text, layout_style.get("caption", {}), theme))
            else:
                if "title" in slide.blocks:
                    objects.append(make_text_object("title", slide.blocks["title"], layout_style.get("title", {}), theme))

                if "subtitle" in slide.blocks:
                    objects.append(make_text_object("subtitle", slide.blocks["subtitle"], layout_style.get("subtitle", {}), theme))

                if "body" in slide.blocks:
                    objects.append(make_text_object("body", slide.blocks["body"], layout_style.get("body", {}), theme))

                if "bullets" in slide.blocks:
                    objects.append(make_bullet_object(slide.blocks["bullets"], layout_style.get("bullets", {}), theme))

                if "image" in slide.blocks:
                    objects.append(make_image_object(slide.blocks["image"], layout_style.get("image", {}), theme))

                # Native hero-image layout: one large image with a caption below.
                # Caption priority: explicit caption block -> image.alt fallback.
                if slide.layout == "hero-image" and "image" in slide.blocks:
                    caption_text = slide.blocks.get("caption", "")
                    if not caption_text and isinstance(slide.blocks["image"], dict):
                        caption_text = slide.blocks["image"].get("alt", "")
                    if caption_text:
                        objects.append(make_text_object("caption", caption_text, layout_style.get("caption", {}), theme))

            if "cells" in slide.blocks:
                objects.extend(make_grid_objects(slide.blocks["cells"], layout_style, theme, slide.layout))

            if slide.layout == "two-images":
                objects.extend(make_two_image_objects(slide.blocks, layout_style, theme))
            else:
                for side in ["left", "right"]:
                    if side in slide.blocks:
                        objects.extend(make_column_objects(side, slide.blocks[side], layout_style.get(side, {}), theme))

            # Footer objects are appended last so they appear above normal content.
            # Supported on all content slide layouts.
            objects.extend(make_footer_text_objects(slide.blocks, layout_style, theme))
            objects.extend(make_footer_image_objects(slide.blocks, layout_style, theme))

            slides_ir.append({
                "type": "RenderSlide",
                "kind": "slide",
                "section": section.title,
                "slide_title": slide.title,
                "layout": slide.layout,
                "background": style_value(layout_style, "background", theme, theme.page.get("background", "#FFFFFF")),
                "background_image": make_background_image_spec(slide.blocks, layout_style),
                "objects": objects,
                "notes": slide.blocks.get("notes", ""),
            })

    return {
        "type": "RenderPresentation",
        "title": doc.title,
        "page": theme.page,
        "slides": slides_ir,
    }


# -----------------------------------------------------------------------------
# Renderers
# -----------------------------------------------------------------------------

def hex_to_rgb(hex_color: str):
    if not hex_color:
        hex_color = "#000000"
    hex_color = str(hex_color).strip().lstrip("#")
    if len(hex_color) != 6:
        hex_color = "000000"
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def px_to_inches(value: float) -> float:
    return value / 96.0



def add_card(slide, obj: Dict[str, Any]) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE without importing enum; stable enough for prototype
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 100))),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(obj.get("fill", "#FFFFFF"))
    shape.line.color.rgb = hex_to_rgb(obj.get("border", "#D1D5DB"))

def add_textbox(slide, obj: Dict[str, Any], bullet: bool = False) -> None:
    shape = slide.shapes.add_textbox(
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 100))),
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    if bullet:
        for idx, item in enumerate(obj["items"]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = obj["font"]
            p.font.size = Pt(obj["size"])
            p.font.color.rgb = hex_to_rgb(obj["color"])
    else:
        p = tf.paragraphs[0]
        p.text = obj["text"]
        p.font.name = obj["font"]
        p.font.size = Pt(obj["size"])
        p.font.color.rgb = hex_to_rgb(obj["color"])


def get_image_size_px(path: Path) -> Optional[Tuple[int, int]]:
    """Return image size in pixels when Pillow is available."""
    if PILImage is None:
        return None
    try:
        with PILImage.open(path) as img:
            return img.size
    except Exception:
        return None


def fitted_box_px(obj: Dict[str, Any], image_size: Optional[Tuple[int, int]]) -> Tuple[float, float, float, float]:
    """Compute x/y/w/h in DSL pixels for contain/original/stretch.

    Cover uses the original target box and crop properties in PPTX.
    """
    x, y, w, h = float(obj["x"]), float(obj["y"]), float(obj["w"]), float(obj.get("h", 100))
    mode = normalize_fit_mode(obj.get("mode", "stretch"))
    if not image_size or mode in {"stretch", "cover"}:
        return x, y, w, h

    iw, ih = image_size
    if iw <= 0 or ih <= 0:
        return x, y, w, h

    if mode == "original":
        new_w, new_h = float(iw), float(ih)
    else:  # contain
        scale = min(w / iw, h / ih)
        new_w, new_h = iw * scale, ih * scale

    return x + (w - new_w) / 2, y + (h - new_h) / 2, new_w, new_h


def apply_cover_crop(pic, box_w: float, box_h: float, image_size: Optional[Tuple[int, int]]) -> None:
    """Crop a PPTX picture so it behaves like CSS object-fit: cover."""
    if not image_size:
        return
    iw, ih = image_size
    if iw <= 0 or ih <= 0 or box_w <= 0 or box_h <= 0:
        return

    img_ar = iw / ih
    box_ar = box_w / box_h

    if img_ar > box_ar:
        visible_w = box_ar * ih
        crop_each = max(0.0, min(0.5, (iw - visible_w) / (2 * iw)))
        pic.crop_left = crop_each
        pic.crop_right = crop_each
    elif img_ar < box_ar:
        visible_h = iw / box_ar
        crop_each = max(0.0, min(0.5, (ih - visible_h) / (2 * ih)))
        pic.crop_top = crop_each
        pic.crop_bottom = crop_each


def add_image(slide, obj: Dict[str, Any]) -> None:
    src = obj.get("src", "")
    path = Path(src)
    if not path.exists():
        # Graceful fallback: show missing path instead of crashing.
        missing = dict(obj)
        missing.update({
            "type": "TextBox",
            "role": "missing-image",
            "text": f"[Missing image: {src}]",
            "font": "Aptos",
            "size": 16,
            "color": "#AA0000",
        })
        add_textbox(slide, missing, bullet=False)
        return

    mode = normalize_fit_mode(obj.get("mode", "stretch"))
    image_size = get_image_size_px(path)
    x, y, w, h = fitted_box_px(obj, image_size)

    pic = slide.shapes.add_picture(
        str(path),
        Inches(px_to_inches(x)),
        Inches(px_to_inches(y)),
        width=Inches(px_to_inches(w)),
        height=Inches(px_to_inches(h)),
    )
    if mode == "cover":
        apply_cover_crop(pic, float(obj["w"]), float(obj.get("h", 100)), image_size)



def add_background_image(slide, spec: Optional[Dict[str, Any]]) -> None:
    """Add a full-slide background image behind all objects.

    Supported PPTX modes:
      stretch  - fill slide, may distort.
      contain  - preserve aspect ratio; blank margins may remain.
      cover/crop - fill slide and crop overflow.
      original - use original pixel size, centered on slide.

    Note: python-pptx does not provide a simple public API for picture opacity.
    The opacity value is preserved in HTML output; PPTX uses the image as-is.
    """
    if not spec:
        return
    src = spec.get("src", "")
    path = Path(src)
    if not path.exists():
        return

    obj = {
        "x": 0,
        "y": 0,
        "w": 1280,
        "h": 720,
        "mode": normalize_fit_mode(spec.get("mode", "cover"), default="cover"),
    }
    image_size = get_image_size_px(path)
    x, y, w, h = fitted_box_px(obj, image_size)

    pic = slide.shapes.add_picture(
        str(path),
        Inches(px_to_inches(x)),
        Inches(px_to_inches(y)),
        width=Inches(px_to_inches(w)),
        height=Inches(px_to_inches(h)),
    )
    if obj["mode"] == "cover":
        apply_cover_crop(pic, 1280, 720, image_size)

def render_pptx(render_ir: Dict[str, Any], output_path: str) -> None:
    if PPTXPresentation is None:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = PPTXPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for slide_ir in render_ir["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_to_rgb(slide_ir.get("background", "#FFFFFF"))

        # Background image is drawn first so all normal objects stay on top.
        add_background_image(slide, slide_ir.get("background_image"))

        for obj in slide_ir["objects"]:
            if obj["type"] == "Card":
                add_card(slide, obj)
            elif obj["type"] == "TextBox":
                add_textbox(slide, obj, bullet=False)
            elif obj["type"] == "BulletList":
                add_textbox(slide, obj, bullet=True)
            elif obj["type"] == "Image":
                add_image(slide, obj)

        notes = slide_ir.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    prs.save(output_path)


def render_html(render_ir: Dict[str, Any], output_path: str) -> None:
    parts: List[str] = []
    parts.append("<!doctype html>")
    parts.append("<html><head><meta charset='utf-8'>")
    parts.append("<title>" + html.escape(render_ir["title"]) + "</title>")
    parts.append(
        """
<style>
body { margin: 0; font-family: Arial, sans-serif; background: #eee; }
.deck { display: flex; flex-direction: column; gap: 24px; padding: 24px; }
.slide { position: relative; width: 1280px; height: 720px; box-shadow: 0 8px 24px rgba(0,0,0,.18); overflow: hidden; }
.textbox, .bullets, .image, .card, .bg-image { position: absolute; box-sizing: border-box; }
ul { margin: 0; padding-left: 1.2em; }
li { margin-bottom: .4em; }
img { display: block; }
</style>
"""
    )
    parts.append("</head><body><div class='deck'>")

    for slide in render_ir["slides"]:
        bg = slide.get("background", "#FFFFFF")
        parts.append(f"<section class='slide' style='background:{html.escape(str(bg))}'>")
        bg_img = slide.get("background_image")
        if bg_img and bg_img.get("src"):
            mode = normalize_fit_mode(bg_img.get("mode", "cover"), default="cover")
            fit_map = {"stretch": "fill", "contain": "contain", "cover": "cover", "original": "none"}
            fit = fit_map.get(mode, "cover")
            opacity = bg_img.get("opacity", 1.0)
            try:
                opacity = max(0.0, min(1.0, float(opacity)))
            except (TypeError, ValueError):
                opacity = 1.0
            src = html.escape(str(bg_img.get("src", "")))
            alt = html.escape(str(bg_img.get("alt", "")))
            parts.append(
                f"<img class='bg-image' src='{src}' alt='{alt}' "
                f"style='left:0; top:0; width:1280px; height:720px; object-fit:{fit}; opacity:{opacity}; z-index:0;'>"
            )
        for obj in slide["objects"]:
            if obj["type"] == "Card":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 100)}px; "
                    f"background:{html.escape(str(obj.get('fill', '#FFFFFF')))}; "
                    f"border:1px solid {html.escape(str(obj.get('border', '#D1D5DB')))}; "
                    f"border-radius:16px;"
                )
                parts.append(f"<div class='card' style='{style}'></div>")
            elif obj["type"] == "TextBox":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 100)}px; "
                    f"font-family:{html.escape(str(obj['font']))}; "
                    f"font-size:{obj['size']}px; color:{html.escape(str(obj['color']))};"
                )
                parts.append(f"<div class='textbox' style='{style}'>{html.escape(obj['text'])}</div>")
            elif obj["type"] == "BulletList":
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj.get('h', 420)}px; "
                    f"font-family:{html.escape(str(obj['font']))}; "
                    f"font-size:{obj['size']}px; color:{html.escape(str(obj['color']))};"
                )
                parts.append(f"<div class='bullets' style='{style}'><ul>")
                for item in obj["items"]:
                    parts.append(f"<li>{html.escape(item)}</li>")
                parts.append("</ul></div>")
            elif obj["type"] == "Image":
                mode = normalize_fit_mode(obj.get("mode", "stretch"))
                fit_map = {"stretch": "fill", "contain": "contain", "cover": "cover", "original": "none"}
                fit = fit_map.get(mode, "fill")
                style = (
                    f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; height:{obj['h']}px;"
                    f"object-fit:{fit};"
                )
                src = html.escape(str(obj.get("src", "")))
                alt = html.escape(str(obj.get("alt", "")))
                parts.append(f"<img class='image' src='{src}' alt='{alt}' style='{style}'>")
        parts.append("</section>")

    parts.append("</div></body></html>")
    Path(output_path).write_text("\n".join(parts), encoding="utf-8")


# -----------------------------------------------------------------------------
# Demo input
# -----------------------------------------------------------------------------

DEMO_PML = '''
presentation "AI Strategy 2026":
  meta:
    author: "Viettel AI Lab"
    language: vi
    format: pptx

  use style: "corporate.psl"

  cover_layout: title-slide
  cover:
    subtitle: Chiến lược tự động hóa tri thức doanh nghiệp
    author: Viettel AI Lab
    date: 2026
    background_image:
      src: "demo_bg.png"
      mode: cover
      opacity: 0.18

  section "Bối cảnh":
    header_layout: section-header
    header:
      subtitle: Từ dữ liệu phân tán đến insight có thể hành động

    slide "Vấn đề hiện tại":
      layout: title-bullets
      intent: explain_problem

      background_image:
        src: "corporate_bg.png"
        opacity: 0.12
        mode: cover

      title:
        Quá tải thông tin trong doanh nghiệp

      bullets:
        - Dữ liệu phân tán ở nhiều hệ thống
        - Báo cáo thủ công mất nhiều thời gian
        - Lãnh đạo cần insight nhanh hơn

      footer_text:
        Internal Use Only | Viettel AI Lab

      footer_image:
        src: "viettel_logo.png"
        alt: "Viettel logo"

      notes:
        Nhấn mạnh vấn đề không nằm ở thiếu dữ liệu, mà là thiếu khả năng tổng hợp.

    slide "Kiến trúc đề xuất":
      layout: title-image
      intent: show_architecture

      title:
        Kiến trúc tổng quan hệ thống

      subtitle:
        Dữ liệu → xử lý → tổng hợp → trình bày

      image:
        src: "architecture.png"
        alt: "Sơ đồ kiến trúc pipeline"

      notes:
        Có thể thay architecture.png bằng đường dẫn ảnh thật.

    slide "Mô hình triển khai":
      layout: text-image
      intent: explain_architecture

      title:
        Một cột nội dung, một cột hình minh họa

      left:
        heading: Thành phần chính
        bullets:
          - Thu thập dữ liệu từ nhiều nguồn
          - Chuẩn hóa và tạo chỉ mục
          - Truy xuất ngữ cảnh liên quan
          - Sinh báo cáo hoặc slide

      right:
        image:
          src: "architecture.png"
          alt: "Minh họa kiến trúc hệ thống"

      notes:
        Layout text-image dùng left cho chữ và right cho ảnh.

  section "Giải pháp":
    header_layout: section-header
    header:
      subtitle: Thiết kế pipeline và giao diện đầu ra

    slide "Luồng xử lý":
      layout: two-column
      intent: describe_pipeline

      title:
        Pipeline xử lý dữ liệu theo từng giai đoạn

      left:
        heading: Đầu vào
        bullets:
          - Văn bản
          - Báo cáo
          - Dữ liệu mạng xã hội

      right:
        heading: Đầu ra
        bullets:
          - Tóm tắt
          - Dashboard
          - Slide trình bày


    slide "Kiến trúc tổng quan":
      layout: hero-image
      intent: show_overview_diagram

      title:
        Kiến trúc xử lý dữ liệu tổng quan

      subtitle:
        Một ảnh lớn ở trung tâm, chú thích ngắn ở phía dưới

      image:
        src: "architecture_overview.png"
        alt: "Sơ đồ kiến trúc tổng quan"

      caption:
        Hệ thống đi từ thu thập dữ liệu, chuẩn hóa, truy xuất tri thức đến sinh báo cáo và slide trình bày.

    slide "Ảnh minh họa độc lập":
      layout: image-caption
      intent: show_visual_only

      image:
        src: "model_snapshot.png"
        alt: "Ảnh minh họa mô hình tổng hợp tri thức"

      caption:
        Minh họa trực quan mô hình tổng hợp tri thức tự động, không cần tiêu đề phụ phía trên.

      footer:
        text: Confidential | Draft version
        align: center


    slide "So sánh hai phương án":
      layout: two-images
      intent: compare_visual_options

      title:
        So sánh kiến trúc hiện tại và kiến trúc đề xuất

      subtitle:
        Mỗi ảnh chiếm một cột độc lập để dễ đối chiếu trực quan

      images:
        - src: "current_architecture.png"
          alt: "Kiến trúc hiện tại"
        - src: "target_architecture.png"
          alt: "Kiến trúc đề xuất"

      notes:
        Nếu file ảnh chưa tồn tại, renderer PPTX sẽ hiển thị placeholder missing image.

    slide "6 năng lực chính":
      layout: grid-6
      intent: summarize_capabilities

      title:
        Sáu năng lực chính của hệ thống

      cells:
        - heading: Thu thập
          text: Lấy dữ liệu từ nhiều nguồn nội bộ và bên ngoài.
        - heading: Làm sạch
          text: Chuẩn hóa định dạng, khử trùng lặp và lọc nhiễu.
        - heading: Truy xuất
          text: Tìm các đoạn liên quan theo ngữ nghĩa và cấu trúc.
        - heading: Tổng hợp
          text: Kết hợp thông tin thành bản tóm tắt có kiểm soát.
        - heading: Trực quan
          text: Sinh dashboard, biểu đồ hoặc slide trình bày.
        - heading: Giám sát
          text: Theo dõi lỗi, chất lượng đầu ra và độ tin cậy.

      notes:
        Layout grid-6 dùng 3 cột x 2 hàng.
'''


DEMO_PSL = '''
theme "corporate":
  page:
    size: widescreen
    background: "#FFFFFF"

  colors:
    primary: "#003A8C"
    secondary: "#E6F0FF"
    text: "#1F1F1F"
    muted: "#666666"
    white: "#FFFFFF"

  fonts:
    heading: "Aptos Display"
    body: "Aptos"

  presentation.title-slide:
    background: primary
    background_image:
      src: "E:\\SPIRIT\\PHOTO\\43.png"
      opacity: 0.20
      mode: original

    title:
      font: heading
      size: 48
      color: white
      position: [90, 220]
      width: 1100
      height: 90

    subtitle:
      font: body
      size: 26
      color: secondary
      position: [95, 320]
      width: 1000
      height: 60

    author:
      font: body
      size: 20
      color: white
      position: [95, 430]
      width: 800
      height: 40

    date:
      font: body
      size: 18
      color: secondary
      position: [95, 470]
      width: 400
      height: 40

  section.section-header:
    background: secondary
    title:
      font: heading
      size: 44
      color: primary
      position: [90, 250]
      width: 1100
      height: 90

    subtitle:
      font: body
      size: 24
      color: text
      position: [95, 340]
      width: 1000
      height: 60

  slide.title-bullets:
    title:
      font: heading
      size: 36
      color: primary
      position: [60, 40]
      width: 1100
      height: 90

    bullets:
      font: body
      size: 24
      color: text
      position: [90, 170]
      width: 1000
      height: 420
      line_gap: 12

    footer_text:
      font: body
      size: 14
      color: muted
      position: [80, 685]
      width: 840
      height: 24

    footer_image:
      position: [1080, 665]
      width: 120
      height: 36

  slide.title-image:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 70

    subtitle:
      font: body
      size: 22
      color: muted
      position: [65, 105]
      width: 1080
      height: 50

    image:
      position: [170, 180]
      width: 940
      height: 450



  slide.text-image:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 155]
      width: 520
      height: 420
      heading_size: 25
      size: 21
      color: text

    right:
      position: [680, 155]
      width: 520
      height: 390
      image_width: 520
      image_height: 390

  slide.image-text:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 155]
      width: 520
      height: 390
      image_width: 520
      image_height: 390

    right:
      position: [680, 155]
      width: 520
      height: 420
      heading_size: 25
      size: 21
      color: text



  slide.hero-image:
    background: "#FFFFFF"

    title:
      font: heading
      size: 34
      color: primary
      position: [60, 36]
      width: 1120
      height: 70

    subtitle:
      font: body
      size: 19
      color: muted
      position: [62, 102]
      width: 1080
      height: 40

    image:
      position: [120, 155]
      width: 1040
      height: 410

    caption:
      font: body
      size: 17
      color: muted
      position: [140, 585]
      width: 1000
      height: 60

  slide.image-caption:
    background: "#FFFFFF"

    image:
      position: [120, 70]
      width: 1040
      height: 520

    caption:
      font: body
      size: 18
      color: muted
      position: [140, 610]
      width: 1000
      height: 60

    footer_text:
      font: body
      size: 14
      color: muted
      position: [140, 685]
      width: 1000
      height: 24


  slide.two-images:
    background: "#FFFFFF"

    title:
      font: heading
      size: 32
      color: primary
      position: [60, 36]
      width: 1120
      height: 70

    subtitle:
      font: body
      size: 19
      color: muted
      position: [62, 104]
      width: 1080
      height: 40

    left:
      position: [70, 165]
      width: 540
      height: 390
      image_width: 540
      image_height: 390
      caption_size: 16
      caption_color: muted

    right:
      position: [670, 165]
      width: 540
      height: 390
      image_width: 540
      image_height: 390
      caption_size: 16
      caption_color: muted

  slide.grid-3:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 155]
      width: 1140
      height: 430
      columns: 3
      gap: 24
      padding: 18
      fill: white
      border: secondary
      heading_size: 21
      size: 16
      color: text

  slide.grid-4:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 145]
      width: 1140
      height: 470
      columns: 2
      gap: 24
      padding: 18
      fill: white
      border: secondary
      heading_size: 21
      size: 16
      color: text

  slide.grid-5:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 145]
      width: 1140
      height: 470
      columns: 3
      gap: 22
      padding: 16
      fill: white
      border: secondary
      heading_size: 20
      size: 15
      color: text

  slide.grid-6:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    grid:
      position: [70, 145]
      width: 1140
      height: 470
      columns: 3
      gap: 22
      padding: 16
      fill: white
      border: secondary
      heading_size: 20
      size: 15
      color: text


  slide.two-column:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 160]
      width: 520

    right:
      position: [680, 160]
      width: 520
'''


def main() -> None:
    doc = parse_pml(DEMO_PML)
    theme = parse_psl(DEMO_PSL)
    render_ir = build_render_ir(doc, theme)

    render_html(render_ir, "demo_output.html")
    render_pptx(render_ir, "demo_output.pptx")

    print("Generated demo_output.html")
    print("Generated demo_output.pptx")


if __name__ == "__main__":
    main()
