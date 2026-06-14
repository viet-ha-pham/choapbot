"""
PML/PSL Renderer v0.1
=====================

A tiny prototype for:

    PML: Presentation Markup Language
    PSL: Presentation Styling Language

It parses simple indentation-based PML/PSL files, builds AST-like Python objects,
merges content with style, and renders to:

    - PowerPoint (.pptx) via python-pptx
    - HTML (.html)

Install:

    pip install python-pptx

Run demo:

    python pml_psl_renderer_v0_1.py

This will generate:

    demo_output.pptx
    demo_output.html
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import html
import re

try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:  # allows HTML-only use if python-pptx is not installed
    PPTXPresentation = None


# -----------------------------------------------------------------------------
# Data models
# -----------------------------------------------------------------------------

@dataclass
class Slide:
    title: str
    layout: str = "title-bullets"
    intent: Optional[str] = None
    blocks: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Section:
    title: str
    slides: List[Slide] = field(default_factory=list)


@dataclass
class PresentationDoc:
    title: str
    meta: Dict[str, Any] = field(default_factory=dict)
    style_file: Optional[str] = None
    sections: List[Section] = field(default_factory=list)


@dataclass
class Theme:
    name: str
    page: Dict[str, Any] = field(default_factory=dict)
    colors: Dict[str, str] = field(default_factory=dict)
    fonts: Dict[str, str] = field(default_factory=dict)
    layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)


# -----------------------------------------------------------------------------
# Low-level indentation parser
# -----------------------------------------------------------------------------

@dataclass
class Line:
    indent: int
    text: str
    lineno: int


def strip_comment(line: str) -> str:
    """Remove comments outside quotes."""
    in_quote = False
    quote_char = ""
    out = []

    for ch in line:
        if ch in ('"', "'"):
            if not in_quote:
                in_quote = True
                quote_char = ch
            elif quote_char == ch:
                in_quote = False
        if ch == "#" and not in_quote:
            break
        out.append(ch)

    return "".join(out).rstrip()


def tokenize_lines(text: str) -> List[Line]:
    lines: List[Line] = []
    for i, raw in enumerate(text.splitlines(), start=1):
        no_comment = strip_comment(raw)
        if not no_comment.strip():
            continue
        indent = len(no_comment) - len(no_comment.lstrip(" "))
        if "\t" in raw[:indent]:
            raise SyntaxError(f"Line {i}: tabs are not allowed for indentation")
        lines.append(Line(indent=indent, text=no_comment.strip(), lineno=i))
    return lines


def parse_scalar(value: str) -> Any:
    value = value.strip()

    if value == "":
        return ""

    if (value.startswith('"') and value.endswith('"')) or (
        value.startswith("'") and value.endswith("'")
    ):
        return value[1:-1]

    if value.lower() == "true":
        return True
    if value.lower() == "false":
        return False
    if value.lower() in {"null", "none"}:
        return None

    if value.startswith("[") and value.endswith("]"):
        inner = value[1:-1].strip()
        if not inner:
            return []
        return [parse_scalar(part.strip()) for part in inner.split(",")]

    # number
    try:
        if "." in value:
            return float(value)
        return int(value)
    except ValueError:
        return value


def split_key_value(text: str) -> Tuple[str, Optional[str]]:
    """Split 'key: value' or 'key:' into key/value."""
    if ":" not in text:
        return text, None
    key, val = text.split(":", 1)
    return key.strip(), val.strip()


def parse_generic_block(lines: List[Line], start: int, indent: int) -> Tuple[Any, int]:
    """
    Parse a generic block into dict/list/scalar.

    Supports:

        key: value
        key:
          nested: value
        list:
          - item
          - item

    Also supports multiline scalar:

        title:
          Some text here
    """
    result: Dict[str, Any] = {}
    items: List[Any] = []
    mode: Optional[str] = None
    i = start

    while i < len(lines):
        line = lines[i]
        if line.indent < indent:
            break
        if line.indent > indent:
            raise SyntaxError(
                f"Line {line.lineno}: unexpected indentation. Expected {indent}, got {line.indent}"
            )

        text = line.text

        if text.startswith("- "):
            if mode is None:
                mode = "list"
            elif mode != "list":
                raise SyntaxError(f"Line {line.lineno}: cannot mix list and mapping")

            item_text = text[2:].strip()
            if item_text.endswith(":"):
                key = item_text[:-1].strip()
                child, new_i = parse_generic_block(lines, i + 1, indent + 2)
                items.append({key: child})
                i = new_i
            else:
                items.append(parse_scalar(item_text))
                i += 1
            continue

        if mode is None:
            mode = "dict"
        elif mode != "dict":
            raise SyntaxError(f"Line {line.lineno}: cannot mix mapping and list")

        key, val = split_key_value(text)
        if val is None:
            # bare scalar line in a block: collect as text until dedent
            text_lines = []
            while i < len(lines) and lines[i].indent == indent:
                text_lines.append(lines[i].text)
                i += 1
            return "\n".join(text_lines), i

        if val != "":
            result[key] = parse_scalar(val)
            i += 1
        else:
            if i + 1 >= len(lines) or lines[i + 1].indent <= indent:
                result[key] = {}
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                result[key] = child
                i = new_i

    if mode == "list":
        return items, i
    return result, i


# -----------------------------------------------------------------------------
# PML parser
# -----------------------------------------------------------------------------

PRESENTATION_RE = re.compile(r'^presentation\s+"(.+)"\s*:$')
SECTION_RE = re.compile(r'^section\s+"(.+)"\s*:$')
SLIDE_RE = re.compile(r'^slide\s+"(.+)"\s*:$')
THEME_RE = re.compile(r'^theme\s+"(.+)"\s*:$')


def parse_pml(text: str) -> PresentationDoc:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PML")

    m = PRESENTATION_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PML must start with: presentation "Title":')

    doc = PresentationDoc(title=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        text = line.text

        if text == "meta:":
            block, i = parse_generic_block(lines, i + 1, 4)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: meta must be a mapping")
            doc.meta = block
            continue

        if text.startswith("use style:"):
            _, value = split_key_value(text)
            doc.style_file = parse_scalar(value or "")
            i += 1
            continue

        sec_match = SECTION_RE.match(text)
        if sec_match:
            section, i = parse_section(lines, i)
            doc.sections.append(section)
            continue

        raise SyntaxError(f"Line {line.lineno}: unknown PML statement: {text}")

    return doc


def parse_section(lines: List[Line], start: int) -> Tuple[Section, int]:
    m = SECTION_RE.match(lines[start].text)
    assert m is not None
    section = Section(title=m.group(1))
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= lines[start].indent:
            break
        if line.indent != lines[start].indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected section child indent")

        slide_match = SLIDE_RE.match(line.text)
        if slide_match:
            slide, i = parse_slide(lines, i)
            section.slides.append(slide)
        else:
            raise SyntaxError(f"Line {line.lineno}: section only accepts slide blocks")

    return section, i


def parse_slide(lines: List[Line], start: int) -> Tuple[Slide, int]:
    m = SLIDE_RE.match(lines[start].text)
    assert m is not None
    slide = Slide(title=m.group(1))
    base_indent = lines[start].indent
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= base_indent:
            break
        if line.indent != base_indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected slide child indent")

        key, val = split_key_value(line.text)

        if key == "layout":
            slide.layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if key == "intent":
            slide.intent = str(parse_scalar(val or ""))
            i += 1
            continue

        if val == "":
            if i + 1 >= len(lines) or lines[i + 1].indent <= line.indent:
                slide.blocks[key] = ""
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                slide.blocks[key] = child
                i = new_i
        else:
            slide.blocks[key] = parse_scalar(val or "")
            i += 1

    return slide, i


# -----------------------------------------------------------------------------
# PSL parser
# -----------------------------------------------------------------------------

def parse_psl(text: str) -> Theme:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PSL")

    m = THEME_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PSL must start with: theme "Name":')

    theme = Theme(name=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        key, val = split_key_value(line.text)
        if val != "":
            raise SyntaxError(f"Line {line.lineno}: top-level PSL keys must be blocks")

        block, i = parse_generic_block(lines, i + 1, 4)

        if key == "page":
            theme.page = block
        elif key == "colors":
            theme.colors = block
        elif key == "fonts":
            theme.fonts = block
        elif key.startswith("slide."):
            layout_name = key.split(".", 1)[1]
            theme.layouts[layout_name] = block
        else:
            # keep unknown blocks as layouts/extensions if desired
            theme.layouts[key] = block

    return theme


# -----------------------------------------------------------------------------
# Render IR construction
# -----------------------------------------------------------------------------

def resolve_token(value: Any, theme: Theme) -> Any:
    if isinstance(value, str):
        if value in theme.colors:
            return theme.colors[value]
        if value in theme.fonts:
            return theme.fonts[value]
    return value


def normalize_text_block(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, dict):
        if "text" in value:
            return str(value["text"])
        if "heading" in value:
            return str(value["heading"])
    return str(value)


def build_render_ir(doc: PresentationDoc, theme: Theme) -> Dict[str, Any]:
    slides_ir = []

    for section in doc.sections:
        for slide in section.slides:
            layout_style = theme.layouts.get(slide.layout, {})
            objects = []

            # Title block
            if "title" in slide.blocks:
                style = layout_style.get("title", {})
                objects.append(make_text_object("Title", slide.blocks["title"], style, theme))

            # Subtitle block
            if "subtitle" in slide.blocks:
                style = layout_style.get("subtitle", layout_style.get("title", {}))
                obj = make_text_object("Subtitle", slide.blocks["subtitle"], style, theme)
                objects.append(obj)

            # Bullets block
            if "bullets" in slide.blocks:
                style = layout_style.get("bullets", {})
                objects.append(make_bullet_object(slide.blocks["bullets"], style, theme))

            # Two-column blocks
            for side in ["left", "right"]:
                if side in slide.blocks:
                    style = layout_style.get(side, {})
                    objects.extend(make_column_objects(side, slide.blocks[side], style, theme))

            slides_ir.append(
                {
                    "type": "RenderSlide",
                    "section": section.title,
                    "slide_title": slide.title,
                    "layout": slide.layout,
                    "background": theme.page.get("background", "#FFFFFF"),
                    "objects": objects,
                    "notes": slide.blocks.get("notes", ""),
                }
            )

    return {
        "type": "RenderPresentation",
        "title": doc.title,
        "page": theme.page,
        "slides": slides_ir,
    }


def style_value(style: Dict[str, Any], key: str, theme: Theme, default: Any = None) -> Any:
    return resolve_token(style.get(key, default), theme)


def make_text_object(kind: str, value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [60, 40])
    return {
        "type": "TextBox",
        "role": kind,
        "text": normalize_text_block(value),
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 80),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 28),
        "color": style_value(style, "color", theme, "#000000"),
    }


def make_bullet_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [90, 160])
    if not isinstance(value, list):
        value = [value]
    return {
        "type": "BulletList",
        "items": [str(v) for v in value],
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 420),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 24),
        "color": style_value(style, "color", theme, "#000000"),
        "line_gap": style.get("line_gap", 8),
    }


def make_column_objects(role: str, value: Any, style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    pos = style.get("position", [70, 150])
    width = style.get("width", 520)
    objects = []

    if isinstance(value, dict):
        if "heading" in value:
            objects.append(
                {
                    "type": "TextBox",
                    "role": f"{role}.heading",
                    "text": str(value["heading"]),
                    "x": pos[0],
                    "y": pos[1],
                    "w": width,
                    "h": 50,
                    "font": theme.fonts.get("heading", "Aptos Display"),
                    "size": 24,
                    "color": theme.colors.get("primary", "#000000"),
                }
            )
        if "bullets" in value:
            objects.append(
                {
                    "type": "BulletList",
                    "items": [str(x) for x in value["bullets"]],
                    "x": pos[0],
                    "y": pos[1] + 60,
                    "w": width,
                    "h": 300,
                    "font": theme.fonts.get("body", "Aptos"),
                    "size": 20,
                    "color": theme.colors.get("text", "#000000"),
                    "line_gap": 6,
                }
            )
    else:
        objects.append(
            {
                "type": "TextBox",
                "role": role,
                "text": str(value),
                "x": pos[0],
                "y": pos[1],
                "w": width,
                "h": 300,
                "font": theme.fonts.get("body", "Aptos"),
                "size": 20,
                "color": theme.colors.get("text", "#000000"),
            }
        )

    return objects


# -----------------------------------------------------------------------------
# PowerPoint renderer
# -----------------------------------------------------------------------------

def hex_to_rgb(hex_color: str) -> RGBColor:
    if not hex_color:
        hex_color = "#000000"
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) != 6:
        hex_color = "000000"
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def px_to_inches(value: float) -> float:
    """
    Treat DSL coordinates as a 1280x720 design canvas.
    PowerPoint widescreen is 13.333 x 7.5 inches.
    """
    return value / 96.0


def add_textbox(slide, obj: Dict[str, Any], bullet: bool = False) -> None:
    shape = slide.shapes.add_textbox(
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 100))),
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    if bullet:
        for idx, item in enumerate(obj["items"]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = obj["font"]
            p.font.size = Pt(obj["size"])
            p.font.color.rgb = hex_to_rgb(obj["color"])
    else:
        p = tf.paragraphs[0]
        p.text = obj["text"]
        p.font.name = obj["font"]
        p.font.size = Pt(obj["size"])
        p.font.color.rgb = hex_to_rgb(obj["color"])


def render_pptx(render_ir: Dict[str, Any], output_path: str) -> None:
    if PPTXPresentation is None:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = PPTXPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for slide_ir in render_ir["slides"]:
        slide = prs.slides.add_slide(blank_layout)

        # background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_to_rgb(slide_ir.get("background", "#FFFFFF"))

        for obj in slide_ir["objects"]:
            if obj["type"] == "TextBox":
                add_textbox(slide, obj, bullet=False)
            elif obj["type"] == "BulletList":
                add_textbox(slide, obj, bullet=True)

        notes = slide_ir.get("notes")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = str(notes)

    prs.save(output_path)


# -----------------------------------------------------------------------------
# HTML renderer
# -----------------------------------------------------------------------------

def render_html(render_ir: Dict[str, Any], output_path: str) -> None:
    parts = []
    parts.append("<!doctype html>")
    parts.append("<html><head><meta charset='utf-8'>")
    parts.append("<title>" + html.escape(render_ir["title"]) + "</title>")
    parts.append(
        """
<style>
body { margin: 0; font-family: Arial, sans-serif; background: #eee; }
.deck { display: flex; flex-direction: column; gap: 24px; padding: 24px; }
.slide { position: relative; width: 1280px; height: 720px; box-shadow: 0 8px 24px rgba(0,0,0,.18); overflow: hidden; }
.textbox, .bullets { position: absolute; }
ul { margin: 0; padding-left: 1.2em; }
li { margin-bottom: .4em; }
</style>
"""
    )
    parts.append("</head><body><div class='deck'>")

    for slide in render_ir["slides"]:
        bg = slide.get("background", "#FFFFFF")
        parts.append(f"<section class='slide' style='background:{html.escape(bg)}'>")
        for obj in slide["objects"]:
            common_style = (
                f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; "
                f"font-family:{html.escape(str(obj['font']))}; "
                f"font-size:{obj['size']}px; color:{html.escape(str(obj['color']))};"
            )
            if obj["type"] == "TextBox":
                parts.append(
                    f"<div class='textbox' style='{common_style}'>"
                    f"{html.escape(obj['text'])}</div>"
                )
            elif obj["type"] == "BulletList":
                parts.append(f"<div class='bullets' style='{common_style}'><ul>")
                for item in obj["items"]:
                    parts.append(f"<li>{html.escape(item)}</li>")
                parts.append("</ul></div>")
        parts.append("</section>")

    parts.append("</div></body></html>")
    Path(output_path).write_text("\n".join(parts), encoding="utf-8")


# -----------------------------------------------------------------------------
# Demo input
# -----------------------------------------------------------------------------

DEMO_PML = '''
presentation "AI Strategy 2026":
  meta:
    author: "Viettel AI Lab"
    language: vi
    format: pptx

  use style: "corporate.psl"

  section "Bối cảnh":

    slide "Vấn đề hiện tại":
      layout: title-bullets
      intent: explain_problem

      title:
        Quá tải thông tin trong doanh nghiệp

      bullets:
        - Dữ liệu phân tán ở nhiều hệ thống
        - Báo cáo thủ công mất nhiều thời gian
        - Lãnh đạo cần insight nhanh hơn

      notes:
        Nhấn mạnh vấn đề không nằm ở thiếu dữ liệu, mà là thiếu khả năng tổng hợp.

    slide "Mục tiêu":
      layout: two-column
      intent: propose_solution

      title:
        Xây dựng hệ thống tổng hợp tri thức tự động

      left:
        heading: Đầu vào
        bullets:
          - Văn bản
          - Báo cáo
          - Dữ liệu mạng xã hội

      right:
        heading: Đầu ra
        bullets:
          - Tóm tắt
          - Dashboard
          - Slide trình bày
'''


DEMO_PSL = '''
theme "corporate":
  page:
    size: widescreen
    background: "#FFFFFF"

  colors:
    primary: "#003A8C"
    secondary: "#E6F0FF"
    text: "#1F1F1F"
    muted: "#666666"

  fonts:
    heading: "Aptos Display"
    body: "Aptos"

  slide.title-bullets:
    title:
      font: heading
      size: 36
      color: primary
      position: [60, 40]
      width: 1100
      height: 90

    bullets:
      font: body
      size: 24
      color: text
      position: [90, 170]
      width: 1000
      height: 420
      line_gap: 12

  slide.two-column:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 160]
      width: 520

    right:
      position: [680, 160]
      width: 520
'''

DEMO_PML2 = '''
presentation "Data Pipeline Overview":
  meta:
    author: "Engineering Team"
    language: vi
    format: pptx

  use style: "tech.psl"

  section "Tổng quan":

    slide "Bài toán":
      layout: title-bullets
      intent: explain_context

      title:
        Tự động hóa xử lý dữ liệu nội bộ

      bullets:
        - Dữ liệu đến từ nhiều nguồn khác nhau
        - Cần làm sạch trước khi phân tích
        - Pipeline phải dễ mở rộng và kiểm soát lỗi

      notes:
        Nhấn mạnh mục tiêu là biến dữ liệu thô thành dữ liệu sẵn sàng khai thác.

    slide "Luồng xử lý":
      layout: two-column
      intent: describe_pipeline

      title:
        Pipeline xử lý dữ liệu theo từng giai đoạn

      left:
        heading: Input
        bullets:
          - File CSV
          - API nội bộ
          - Log hệ thống

      right:
        heading: Output
        bullets:
          - Dataset sạch
          - Báo cáo thống kê
          - Cảnh báo lỗi
'''

DEMO_PSL2 = '''
theme "tech":
  page:
    size: widescreen
    background: "#F8FAFC"

  colors:
    primary: "#0F766E"
    secondary: "#CCFBF1"
    text: "#111827"
    muted: "#6B7280"

  fonts:
    heading: "Time New Roman"
    body: "Arial"

  slide.title-bullets:
    title:
      font: heading
      size: 38
      color: primary
      position: [60, 45]
      width: 1100
      height: 90

    bullets:
      font: body
      size: 25
      color: text
      position: [90, 175]
      width: 1000
      height: 420
      line_gap: 12

  slide.two-column:
    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 165]
      width: 520

    right:
      position: [680, 165]
      width: 520
'''


def main() -> None:
    doc = parse_pml(DEMO_PML)
    theme = parse_psl(DEMO_PSL)
    render_ir = build_render_ir(doc, theme)

    render_html(render_ir, "demo_output.html")
    render_pptx(render_ir, "demo_output.pptx")

    print("Generated demo_output.html")
    print("Generated demo_output.pptx")

    doc = parse_pml(DEMO_PML2)
    theme = parse_psl(DEMO_PSL2)
    render_ir = build_render_ir(doc, theme)

    render_html(render_ir, "demo_output2.html")
    render_pptx(render_ir, "demo_output2.pptx")

    print("Generated demo_output2.html")
    print("Generated demo_output2.pptx")


if __name__ == "__main__":
    main()
