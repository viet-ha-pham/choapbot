from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def safe_text(s: str) -> str:
    return re.sub(r"\s+", " ", s or "").strip()


def emu_to_inch(x) -> float | None:
    if x is None:
        return None
    return round(x / 914400, 3)


def get_shape_kind(shape) -> str:
    t = shape.shape_type
    if t == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if t == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if t == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if t == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "has_text_frame", False):
        return "text"
    return str(t)


def iter_shapes_recursive(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(shape.shapes)


def extract_shape_for_markdown(shape) -> dict[str, Any]:
    item = {
        "name": shape.name,
        "kind": get_shape_kind(shape),
        "left_in": emu_to_inch(shape.left),
        "top_in": emu_to_inch(shape.top),
        "width_in": emu_to_inch(shape.width),
        "height_in": emu_to_inch(shape.height),
        "text": "",
    }

    if getattr(shape, "has_text_frame", False):
        item["text"] = safe_text(shape.text)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        item["image"] = {
            "filename": getattr(shape.image, "filename", None),
            "content_type": getattr(shape.image, "content_type", None),
            "ext": getattr(shape.image, "ext", None),
            "size": getattr(shape.image, "size", None),
        }

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        rows = []
        for row in shape.table.rows:
            rows.append([safe_text(cell.text) for cell in row.cells])
        item["table"] = rows

    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        try:
            item["chart"] = {
                "chart_type": str(shape.chart.chart_type),
                "has_legend": shape.chart.has_legend,
            }
        except Exception:
            item["chart"] = {}

    return item


def extract_title(slide_data: dict[str, Any]) -> str | None:
    candidates = [
        sh for sh in slide_data["shapes"]
        if sh["kind"] == "text" and sh.get("text")
    ]
    if not candidates:
        return None

    candidates.sort(key=lambda x: (x.get("top_in") or 999, x.get("left_in") or 999))
    return candidates[0]["text"]


def shape_to_markdown(shape: dict[str, Any]) -> str:
    kind = shape["kind"]

    if kind == "text":
        if not shape.get("text"):
            return ""
        return f"- **TextBox** `{shape['name']}`: {shape['text']}"

    if kind == "picture":
        img = shape.get("image", {})
        label = img.get("filename") or img.get("content_type") or ""
        return f"- **Image** `{shape['name']}`: {label}"

    if kind == "table":
        rows = shape.get("table", [])
        if not rows:
            return f"- **Table** `{shape['name']}`"

        lines = [f"- **Table** `{shape['name']}`:"]
        header = rows[0]
        lines.append("  " + "| " + " | ".join(header) + " |")
        lines.append("  " + "| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            lines.append("  " + "| " + " | ".join(row) + " |")
        return "\n".join(lines)

    if kind == "chart":
        return f"- **Chart** `{shape['name']}`: {shape.get('chart', {})}"

    return f"- **Shape** `{shape['name']}`: kind={kind}"


def export_pptx_to_markdown(pptx_path: str | Path) -> str:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    lines = [
        "# PPTX symbolic markdown",
        "",
        f"- File: `{pptx_path}`",
        f"- Slides: {len(prs.slides)}",
        f"- Size: {emu_to_inch(prs.slide_width)} × {emu_to_inch(prs.slide_height)} in",
        "",
    ]

    for i, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "index": i,
            "layout_name": slide.slide_layout.name,
            "shapes": [
                extract_shape_for_markdown(sh)
                for sh in iter_shapes_recursive(slide.shapes)
            ],
        }
        title = extract_title(slide_data)

        lines.append(f"## Slide {i}: {title or '(untitled)'}")
        lines.append("")
        lines.append(f"- Layout: `{slide_data['layout_name']}`")
        lines.append("")

        for shape in slide_data["shapes"]:
            md = shape_to_markdown(shape)
            if md:
                lines.append(md)

        lines.append("")

    return "\n".join(lines)


def save_markdown(pptx_path: str | Path, out_md: str | Path) -> None:
    md = export_pptx_to_markdown(pptx_path)
    Path(out_md).write_text(md, encoding="utf-8")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("--out", default="deck.md")
    args = parser.parse_args()

    save_markdown(args.pptx, args.pptx.replace(".pptx", ".md"))