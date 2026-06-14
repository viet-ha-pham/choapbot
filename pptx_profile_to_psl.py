#!/usr/bin/env python3
"""Infer a renderer-compatible PSL theme from a sample PPTX profile."""
from __future__ import annotations

import json
import re
from collections import Counter
from pathlib import Path
from typing import Any, Dict, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    from pptx_profile_inferer import infer_pptx_profile
except Exception:  # pragma: no cover
    infer_pptx_profile = None

DEFAULT_COLORS = {
    "primary": "#003A8C",
    "secondary": "#E6F0FF",
    "text": "#1F1F1F",
    "muted": "#666666",
    "background": "#FFFFFF",
    "white": "#FFFFFF",
}


def normalize_hex(value: Any, default: str = "#000000") -> str:
    if not value:
        return default
    s = str(value).strip()
    if not s:
        return default
    if not s.startswith("#"):
        s = "#" + s
    if re.fullmatch(r"#[0-9A-Fa-f]{6}", s):
        return s.upper()
    return default


def hex_to_rgb_tuple(hex_color: str) -> Tuple[int, int, int]:
    h = normalize_hex(hex_color).lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def rel_luminance(hex_color: str) -> float:
    r, g, b = [x / 255.0 for x in hex_to_rgb_tuple(hex_color)]
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def choose_palette(profile: Dict[str, Any]) -> Dict[str, str]:
    theme_colors = list((profile.get("theme") or {}).get("colors", {}).values())
    text_colors = list((profile.get("style_stats") or {}).get("text_colors", {}).keys())
    raw = [normalize_hex(c, "") for c in (text_colors + theme_colors)]
    colors = [c for c in raw if c and re.fullmatch(r"#[0-9A-F]{6}", c)]
    if not colors:
        return dict(DEFAULT_COLORS)

    counts = Counter(colors)
    sorted_colors = [c for c, _ in counts.most_common()]
    dark = [c for c in sorted_colors if rel_luminance(c) < 0.45]
    light = [c for c in sorted_colors if rel_luminance(c) >= 0.75]
    accent = [c for c in sorted_colors if c not in {"#000000", "#FFFFFF"} and 0.12 <= rel_luminance(c) <= 0.78]

    return {
        "primary": accent[0] if accent else (dark[0] if dark else sorted_colors[0]),
        "secondary": light[1] if len(light) > 1 else DEFAULT_COLORS["secondary"],
        "text": dark[0] if dark else DEFAULT_COLORS["text"],
        "muted": dark[1] if len(dark) > 1 else DEFAULT_COLORS["muted"],
        "background": light[0] if light else DEFAULT_COLORS["background"],
        "white": "#FFFFFF",
    }


def choose_fonts(profile: Dict[str, Any]) -> Dict[str, str]:
    inference = profile.get("inference") or {}
    theme = profile.get("theme") or {}
    stats = profile.get("style_stats") or {}
    candidates = []
    candidates.extend(inference.get("likely_font_family") or [])
    candidates.extend(list((stats.get("font_names") or {}).keys()))
    candidates.extend(theme.get("major_fonts") or [])
    candidates.extend(theme.get("minor_fonts") or [])
    clean = [str(f).strip() for f in candidates if str(f).strip()]
    heading = clean[0] if clean else "Aptos Display"
    body = clean[1] if len(clean) > 1 else (clean[0] if clean else "Aptos")
    return {"heading": heading, "body": body}


def choose_sizes(profile: Dict[str, Any]) -> Dict[str, int]:
    sizes = (profile.get("style_stats") or {}).get("font_sizes_pt") or {}
    median = sizes.get("median") or 20
    max_size = sizes.get("max") or 36
    body = int(max(18, min(24, round(float(median)))))
    title = int(max(30, min(44, round(float(max_size)))))
    return {
        "title": title,
        "subtitle": int(max(20, min(28, body + 2))),
        "body": body,
        "caption": int(max(14, min(18, body - 4))),
        "table": int(max(14, min(18, body - 3))),
    }


def image_ext_from_content_type(content_type: str) -> str:
    if "png" in content_type:
        return ".png"
    if "jpeg" in content_type or "jpg" in content_type:
        return ".jpg"
    return ".png"


def extract_largest_background_asset(pptx_path: str | Path, out_dir: str | Path, min_area_ratio: float = 0.35) -> Optional[Path]:
    if Presentation is None or MSO_SHAPE_TYPE is None:
        return None
    pptx_path = Path(pptx_path)
    out_dir = Path(out_dir)
    asset_dir = out_dir / "template_assets"
    asset_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx_path))
    slide_area = float(prs.slide_width * prs.slide_height) if prs.slide_width and prs.slide_height else 1.0
    best = None
    best_area = 0.0
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            area = float(shape.width * shape.height)
            if area / slide_area >= min_area_ratio and area > best_area:
                best = (slide_index, shape)
                best_area = area
    if not best:
        return None
    slide_index, shape = best
    image = shape.image
    ext = image_ext_from_content_type(getattr(image, "content_type", ""))
    out_path = asset_dir / f"template_background_s{slide_index}{ext}"
    out_path.write_bytes(image.blob)
    return out_path


def profile_to_psl(profile: Dict[str, Any], theme_name: str = "inferred-template", background_image_path: str | Path | None = None) -> str:
    colors = choose_palette(profile)
    fonts = choose_fonts(profile)
    sizes = choose_sizes(profile)
    bg_page_block = ""
    bg_layout_block = ""
    if background_image_path:
        bg = str(Path(background_image_path).resolve()).replace("\\", "/")
        bg_layout_block = f'''    background_image:\n      src: "{bg}"\n      mode: cover\n      opacity: 0.14\n'''
        bg_page_block = f'''\n    background_image:\n      src: "{bg}"\n      mode: cover\n      opacity: 0.14\n'''

    return f'''theme "{theme_name}":
  page:
    size: widescreen
    background: "{colors['background']}"{bg_page_block}

  colors:
    primary: "{colors['primary']}"
    secondary: "{colors['secondary']}"
    text: "{colors['text']}"
    muted: "{colors['muted']}"
    white: "{colors['white']}"

  fonts:
    heading: "{fonts['heading']}"
    body: "{fonts['body']}"

  presentation.title-slide:
    background: background
{bg_layout_block}    title:
      font: heading
      size: {max(sizes['title'] + 4, 38)}
      color: primary
      bold: true
      align: center
      position: [90, 210]
      width: 1100
      height: 95
    subtitle:
      font: body
      size: {sizes['subtitle']}
      color: text
      align: center
      position: [140, 315]
      width: 1000
      height: 60
    author:
      font: body
      size: {sizes['caption']}
      color: muted
      align: center
      position: [160, 430]
      width: 960
      height: 35

  section.section-header:
    background: secondary
    title:
      font: heading
      size: {max(sizes['title'], 34)}
      color: primary
      bold: true
      align: center
      position: [100, 250]
      width: 1080
      height: 80
    subtitle:
      font: body
      size: {sizes['subtitle']}
      color: text
      align: center
      position: [150, 345]
      width: 980
      height: 60

  slide.title-bullets:
    title:
      font: heading
      size: {sizes['title']}
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 80
    subtitle:
      font: body
      size: {sizes['caption'] + 2}
      color: muted
      position: [65, 105]
      width: 1080
      height: 40
    bullets:
      font: body
      size: {sizes['body']}
      color: text
      position: [90, 165]
      width: 1040
      height: 390
      line_gap: 10
      overflow: paginate
      min_size: 14
    conclusion:
      font: body
      size: {sizes['caption'] + 2}
      color: text
      fill: secondary
      border: primary
      bold: true
      align: center
      position: [80, 585]
      width: 1120
      height: 70

  slide.two-column:
    title:
      font: heading
      size: {max(sizes['title'] - 2, 30)}
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 80
    left:
      position: [70, 160]
      width: 520
      font: body
      size: {max(sizes['body'] - 2, 18)}
      color: text
    right:
      position: [680, 160]
      width: 520
      font: body
      size: {max(sizes['body'] - 2, 18)}
      color: text

  slide.title-table:
    title:
      font: heading
      size: {max(sizes['title'] - 2, 30)}
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 80
    table:
      position: [70, 150]
      width: 1140
      height: 460
      font: body
      size: {sizes['table']}
      align: center
      header_fill: primary
      header_color: white
      border_color: muted
      cell_fill: "#FFFFFF"

  slide.timeline:
    title:
      font: heading
      size: {max(sizes['title'] - 2, 30)}
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 80
    timeline:
      position: [90, 180]
      width: 1100
      height: 370
      line_color: primary
      card_fill: "#FFFFFF"
      card_border: secondary
      text_color: text

  slide.grid-4:
    title:
      font: heading
      size: {max(sizes['title'] - 2, 30)}
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 80
    grid:
      position: [70, 150]
      width: 1140
      height: 470
      columns: 2
      gap: 24
    card:
      fill: "#FFFFFF"
      border: secondary
      heading_color: primary
      text_color: text

  slide.numbered-columns-4:
    title:
      font: heading
      size: {max(sizes['title'] - 2, 30)}
      color: primary
      bold: true
      position: [60, 40]
      width: 1120
      height: 80
    columns:
      position: [60, 170]
      width: 1160
      height: 410
      gap: 20
    card:
      fill: "#FFFFFF"
      border: secondary
      heading_color: primary
      text_color: text
    number:
      fill: primary
      color: white
'''


def infer_template_psl(pptx_path: str | Path, out_dir: str | Path, theme_name: str = "inferred-template", extract_background: bool = True) -> Tuple[str, Dict[str, Any], Optional[Path]]:
    if infer_pptx_profile is None:
        raise RuntimeError("Cannot import pptx_profile_inferer.infer_pptx_profile")
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    profile = infer_pptx_profile(pptx_path)
    bg_path = extract_largest_background_asset(pptx_path, out_dir) if extract_background else None
    psl = profile_to_psl(profile, theme_name=theme_name, background_image_path=bg_path)
    (out_dir / "template_profile.json").write_text(json.dumps(profile, ensure_ascii=False, indent=2), encoding="utf-8")
    (out_dir / "template_inferred.psl").write_text(psl, encoding="utf-8")
    return psl, profile, bg_path


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("--out-dir", default="template_style_out")
    parser.add_argument("--name", default="inferred-template")
    parser.add_argument("--no-background", action="store_true")
    args = parser.parse_args()
    psl, profile, bg = infer_template_psl(args.pptx, args.out_dir, theme_name=args.name, extract_background=not args.no_background)
    print(psl)
    print(f"\nProfile: {Path(args.out_dir) / 'template_profile.json'}")
    if bg:
        print(f"Background asset: {bg}")
