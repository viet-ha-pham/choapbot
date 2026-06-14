# pptx_miner.py
from __future__ import annotations

import json
import re
import statistics
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def safe_text(s: str) -> str:
    return re.sub(r"\s+", " ", s or "").strip()


def rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    try:
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return str(rgb)


def emu_to_pt(x) -> float | None:
    if x is None:
        return None
    return round(x / 12700, 2)


def emu_to_inch(x) -> float | None:
    if x is None:
        return None
    return round(x / 914400, 3)


def get_shape_kind(shape) -> str:
    t = shape.shape_type
    if t == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if t == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if t == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if t == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "has_text_frame", False):
        return "text"
    return str(t)


def iter_shapes_recursive(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(shape.shapes)


def extract_run_style(run) -> dict[str, Any]:
    font = run.font
    style = {
        "text": run.text,
        "font_name": font.name,
        "font_size_pt": emu_to_pt(font.size) if font.size else None,
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color": None,
    }

    try:
        if font.color and font.color.rgb:
            style["color"] = f"#{font.color.rgb}"
    except Exception:
        pass

    return style


def extract_paragraph_style(paragraph) -> dict[str, Any]:
    return {
        "text": safe_text(paragraph.text),
        "level": paragraph.level,
        "alignment": str(paragraph.alignment) if paragraph.alignment else None,
        "runs": [extract_run_style(r) for r in paragraph.runs if r.text],
    }


def extract_shape(shape) -> dict[str, Any]:
    item = {
        "name": shape.name,
        "kind": get_shape_kind(shape),
        "left_in": emu_to_inch(shape.left),
        "top_in": emu_to_inch(shape.top),
        "width_in": emu_to_inch(shape.width),
        "height_in": emu_to_inch(shape.height),
        "text": "",
        "paragraphs": [],
    }

    if getattr(shape, "has_text_frame", False):
        item["text"] = safe_text(shape.text)
        item["paragraphs"] = [
            extract_paragraph_style(p)
            for p in shape.text_frame.paragraphs
            if safe_text(p.text)
        ]

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        item["image"] = {
            "content_type": getattr(shape.image, "content_type", None),
            "filename": getattr(shape.image, "filename", None),
            "ext": getattr(shape.image, "ext", None),
            "size": getattr(shape.image, "size", None),
        }

    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        rows = []
        for row in shape.table.rows:
            rows.append([safe_text(cell.text) for cell in row.cells])
        item["table"] = rows

    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        try:
            item["chart"] = {
                "chart_type": str(shape.chart.chart_type),
                "has_legend": shape.chart.has_legend,
            }
        except Exception:
            item["chart"] = {}

    return item


def extract_theme_from_zip(pptx_path: str | Path) -> dict[str, Any]:
    """
    python-pptx không expose đầy đủ theme.
    Hàm này đọc trực tiếp ppt/theme/theme*.xml trong file pptx.
    """
    pptx_path = Path(pptx_path)
    result = {
        "theme_files": [],
        "theme_name": None,
        "colors": {},
        "major_fonts": [],
        "minor_fonts": [],
    }

    with zipfile.ZipFile(pptx_path) as z:
        theme_files = [n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
        result["theme_files"] = theme_files

        if not theme_files:
            return result

        xml = z.read(theme_files[0])
        root = ET.fromstring(xml)

        result["theme_name"] = root.attrib.get("name")

        # Color scheme
        clr_scheme = root.find(".//a:clrScheme", NS)
        if clr_scheme is not None:
            for child in clr_scheme:
                key = child.tag.split("}")[-1]
                srgb = child.find(".//a:srgbClr", NS)
                sysclr = child.find(".//a:sysClr", NS)
                if srgb is not None:
                    result["colors"][key] = "#" + srgb.attrib.get("val", "")
                elif sysclr is not None:
                    result["colors"][key] = "#" + sysclr.attrib.get("lastClr", "")

        # Font scheme
        major = root.findall(".//a:fontScheme/a:majorFont//a:latin", NS)
        minor = root.findall(".//a:fontScheme/a:minorFont//a:latin", NS)
        result["major_fonts"] = [x.attrib.get("typeface") for x in major if x.attrib.get("typeface")]
        result["minor_fonts"] = [x.attrib.get("typeface") for x in minor if x.attrib.get("typeface")]

    return result


def collect_style_stats(slides: list[dict[str, Any]]) -> dict[str, Any]:
    font_names = Counter()
    font_sizes = []
    colors = Counter()
    bold_count = 0
    italic_count = 0
    run_count = 0

    for slide in slides:
        for shape in slide["shapes"]:
            for p in shape.get("paragraphs", []):
                for r in p.get("runs", []):
                    run_count += 1
                    if r.get("font_name"):
                        font_names[r["font_name"]] += 1
                    if r.get("font_size_pt"):
                        font_sizes.append(r["font_size_pt"])
                    if r.get("color"):
                        colors[r["color"]] += 1
                    if r.get("bold"):
                        bold_count += 1
                    if r.get("italic"):
                        italic_count += 1

    return {
        "font_names": dict(font_names.most_common()),
        "font_sizes_pt": {
            "all": sorted(set(font_sizes)),
            "median": statistics.median(font_sizes) if font_sizes else None,
            "min": min(font_sizes) if font_sizes else None,
            "max": max(font_sizes) if font_sizes else None,
        },
        "text_colors": dict(colors.most_common()),
        "bold_ratio": round(bold_count / run_count, 3) if run_count else None,
        "italic_ratio": round(italic_count / run_count, 3) if run_count else None,
    }


def infer_deck_style(profile: dict[str, Any]) -> dict[str, Any]:
    stats = profile["style_stats"]
    theme = profile["theme"]

    colors = list(stats.get("text_colors", {}).keys()) or list(theme.get("colors", {}).values())
    font_names = list(stats.get("font_names", {}).keys())
    median_size = stats.get("font_sizes_pt", {}).get("median")

    inference = {
        "likely_font_family": font_names[:3] or theme.get("major_fonts") or theme.get("minor_fonts"),
        "likely_palette": colors[:8],
        "density": None,
        "visual_style": [],
        "notes": [],
    }

    shape_counts = [len(s["shapes"]) for s in profile["slides"]]
    avg_shapes = statistics.mean(shape_counts) if shape_counts else 0

    text_lengths = [
        len(safe_text(" ".join(sh.get("text", "") for sh in s["shapes"])))
        for s in profile["slides"]
    ]
    avg_text = statistics.mean(text_lengths) if text_lengths else 0

    if avg_text < 120:
        inference["density"] = "low_text / presentation-oriented"
    elif avg_text < 350:
        inference["density"] = "medium_text"
    else:
        inference["density"] = "high_text / report-like"

    if avg_shapes > 8:
        inference["visual_style"].append("shape-heavy")
    if median_size and median_size >= 24:
        inference["visual_style"].append("large-readable-fonts")
    if len(colors) <= 3:
        inference["visual_style"].append("limited-color-palette")
    elif len(colors) >= 6:
        inference["visual_style"].append("varied-color-palette")

    if not font_names:
        inference["notes"].append("Many font values are inherited from theme/master and not directly visible through python-pptx runs.")

    return inference


def shape_to_markdown(shape: dict[str, Any]) -> str:
    kind = shape["kind"]

    if kind == "text":
        txt = shape.get("text", "")
        if not txt:
            return ""
        return f"- **TextBox** `{shape['name']}`: {txt}"

    if kind == "picture":
        img = shape.get("image", {})
        return f"- **Image** `{shape['name']}`: {img.get('filename') or img.get('content_type') or ''}"

    if kind == "table":
        rows = shape.get("table", [])
        if not rows:
            return f"- **Table** `{shape['name']}`"
        md = [f"- **Table** `{shape['name']}`:"]
        header = rows[0]
        md.append("  " + "| " + " | ".join(header) + " |")
        md.append("  " + "| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            md.append("  " + "| " + " | ".join(row) + " |")
        return "\n".join(md)

    if kind == "chart":
        return f"- **Chart** `{shape['name']}`: {shape.get('chart', {})}"

    return f"- **Shape** `{shape['name']}`: kind={kind}"


def build_markdown(profile: dict[str, Any]) -> str:
    lines = []
    lines.append(f"# PPTX symbolic markdown")
    lines.append("")
    lines.append(f"- File: `{profile['file']}`")
    lines.append(f"- Slides: {len(profile['slides'])}")
    lines.append(f"- Size: {profile['page_size']['width_in']} × {profile['page_size']['height_in']} in")
    lines.append("")

    for slide in profile["slides"]:
        lines.append(f"## Slide {slide['index']}: {slide.get('title') or '(untitled)'}")
        lines.append("")
        lines.append(f"- Layout: `{slide.get('layout_name')}`")
        lines.append(f"- Notes: {slide.get('notes') or ''}")
        lines.append("")
        for sh in slide["shapes"]:
            md = shape_to_markdown(sh)
            if md:
                lines.append(md)
        lines.append("")

    lines.append("## Extracted / inferred style profile")
    lines.append("")
    lines.append("```json")
    compact = {
        "theme": profile["theme"],
        "style_stats": profile["style_stats"],
        "inference": profile["inference"],
    }
    lines.append(json.dumps(compact, ensure_ascii=False, indent=2))
    lines.append("```")
    return "\n".join(lines)


def read_slide_notes(slide) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        notes = slide.notes_slide.notes_text_frame.text
        return safe_text(notes)
    except Exception:
        return ""


def extract_title(slide_data: dict[str, Any]) -> str | None:
    candidates = []
    for sh in slide_data["shapes"]:
        if sh["kind"] == "text" and sh.get("text"):
            top = sh.get("top_in") or 999
            size_max = 0
            for p in sh.get("paragraphs", []):
                for r in p.get("runs", []):
                    if r.get("font_size_pt"):
                        size_max = max(size_max, r["font_size_pt"])
            candidates.append((top, -size_max, sh["text"]))
    if not candidates:
        return None
    candidates.sort()
    return candidates[0][2]


def parse_pptx(
    pptx_path: str | Path,
    llm_style_reasoner: Callable[[dict[str, Any]], dict[str, Any]] | None = None,
) -> dict[str, Any]:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "index": i,
            "layout_name": slide.slide_layout.name,
            "notes": read_slide_notes(slide),
            "shapes": [extract_shape(sh) for sh in iter_shapes_recursive(slide.shapes)],
        }
        slide_data["title"] = extract_title(slide_data)
        slides.append(slide_data)

    profile = {
        "file": str(pptx_path),
        "page_size": {
            "width_in": emu_to_inch(prs.slide_width),
            "height_in": emu_to_inch(prs.slide_height),
        },
        "theme": extract_theme_from_zip(pptx_path),
        "slides": slides,
    }

    profile["style_stats"] = collect_style_stats(slides)
    profile["inference"] = infer_deck_style(profile)

    if llm_style_reasoner:
        profile["llm_inference"] = llm_style_reasoner(profile)

    profile["markdown"] = build_markdown(profile)
    return profile


def export_pptx_markdown_and_profile(
    pptx_path: str | Path,
    out_md: str | Path | None = None,
    out_json: str | Path | None = None,
) -> dict[str, Any]:
    profile = parse_pptx(pptx_path)

    if out_md:
        Path(out_md).write_text(profile["markdown"], encoding="utf-8")

    if out_json:
        data = {k: v for k, v in profile.items() if k != "markdown"}
        Path(out_json).write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

    return profile


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("--md", default="deck.md")
    parser.add_argument("--json", default="deck_profile.json")
    args = parser.parse_args()

    profile = export_pptx_markdown_and_profile(args.pptx, args.md, args.json)
    print(profile["markdown"])