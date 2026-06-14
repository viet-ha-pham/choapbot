"""
PML/PSL Renderer v0.3
=====================

Adds first-class support for:

- Presentation title / cover slide via `cover:` and `presentation.<layout>`
- Section header / divider slide via `header:` and `section.<layout>`
- Normal content slides via `slide.<layout>`

Run:
    python demo_dsl_title_section.py

Outputs:
    demo_output_v3.html
    demo_output_v3.pptx
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
import html
import re

try:
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:  # HTML-only mode still works
    PPTXPresentation = None
    RGBColor = None


# -----------------------------------------------------------------------------
# Data models
# -----------------------------------------------------------------------------

@dataclass
class Slide:
    title: str
    layout: str = "title-bullets"
    intent: Optional[str] = None
    blocks: Dict[str, Any] = field(default_factory=dict)


@dataclass
class Section:
    title: str
    header_layout: str = "section-header"
    header: Dict[str, Any] = field(default_factory=dict)
    slides: List[Slide] = field(default_factory=list)


@dataclass
class PresentationDoc:
    title: str
    meta: Dict[str, Any] = field(default_factory=dict)
    style_file: Optional[str] = None
    cover_layout: str = "title-slide"
    cover: Dict[str, Any] = field(default_factory=dict)
    sections: List[Section] = field(default_factory=list)


@dataclass
class Theme:
    name: str
    page: Dict[str, Any] = field(default_factory=dict)
    colors: Dict[str, str] = field(default_factory=dict)
    fonts: Dict[str, str] = field(default_factory=dict)
    presentation_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    section_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)
    slide_layouts: Dict[str, Dict[str, Any]] = field(default_factory=dict)


# -----------------------------------------------------------------------------
# Low-level indentation parser
# -----------------------------------------------------------------------------

@dataclass
class Line:
    indent: int
    text: str
    lineno: int


def strip_comment(line: str) -> str:
    """Remove comments outside quotes."""
    in_quote = False
    quote_char = ""
    out: List[str] = []

    for ch in line:
        if ch in ('"', "'"):
            if not in_quote:
                in_quote = True
                quote_char = ch
            elif quote_char == ch:
                in_quote = False
        if ch == "#" and not in_quote:
            break
        out.append(ch)

    return "".join(out).rstrip()


def tokenize_lines(text: str) -> List[Line]:
    lines: List[Line] = []
    for i, raw in enumerate(text.splitlines(), start=1):
        no_comment = strip_comment(raw)
        if not no_comment.strip():
            continue
        indent = len(no_comment) - len(no_comment.lstrip(" "))
        if "\t" in raw[:indent]:
            raise SyntaxError(f"Line {i}: tabs are not allowed for indentation")
        lines.append(Line(indent=indent, text=no_comment.strip(), lineno=i))
    return lines


def parse_scalar(value: str) -> Any:
    value = value.strip()

    if value == "":
        return ""

    if (value.startswith('"') and value.endswith('"')) or (
        value.startswith("'") and value.endswith("'")
    ):
        return value[1:-1]

    lower = value.lower()
    if lower == "true":
        return True
    if lower == "false":
        return False
    if lower in {"null", "none"}:
        return None

    if value.startswith("[") and value.endswith("]"):
        inner = value[1:-1].strip()
        if not inner:
            return []
        return [parse_scalar(part.strip()) for part in inner.split(",")]

    try:
        if "." in value:
            return float(value)
        return int(value)
    except ValueError:
        return value


def split_key_value(text: str) -> Tuple[str, Optional[str]]:
    """Split 'key: value' or 'key:' into key/value."""
    if ":" not in text:
        return text, None
    key, val = text.split(":", 1)
    return key.strip(), val.strip()


def parse_generic_block(lines: List[Line], start: int, indent: int) -> Tuple[Any, int]:
    """
    Parse a generic indentation-based block into dict/list/scalar.

    Supports:
        key: value
        key:
          nested: value
        list:
          - item
          - item
        title:
          Multiline scalar text
    """
    result: Dict[str, Any] = {}
    items: List[Any] = []
    mode: Optional[str] = None
    i = start

    while i < len(lines):
        line = lines[i]
        if line.indent < indent:
            break
        if line.indent > indent:
            raise SyntaxError(
                f"Line {line.lineno}: unexpected indentation. Expected {indent}, got {line.indent}"
            )

        text = line.text

        if text.startswith("- "):
            if mode is None:
                mode = "list"
            elif mode != "list":
                raise SyntaxError(f"Line {line.lineno}: cannot mix list and mapping")

            item_text = text[2:].strip()
            if item_text.endswith(":"):
                key = item_text[:-1].strip()
                child, new_i = parse_generic_block(lines, i + 1, indent + 2)
                items.append({key: child})
                i = new_i
            else:
                items.append(parse_scalar(item_text))
                i += 1
            continue

        if mode is None:
            mode = "dict"
        elif mode != "dict":
            raise SyntaxError(f"Line {line.lineno}: cannot mix mapping and list")

        key, val = split_key_value(text)
        if val is None:
            text_lines = []
            while i < len(lines) and lines[i].indent == indent:
                text_lines.append(lines[i].text)
                i += 1
            return "\n".join(text_lines), i

        if val != "":
            result[key] = parse_scalar(val)
            i += 1
        else:
            if i + 1 >= len(lines) or lines[i + 1].indent <= indent:
                result[key] = {}
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                result[key] = child
                i = new_i

    if mode == "list":
        return items, i
    return result, i


# -----------------------------------------------------------------------------
# PML parser
# -----------------------------------------------------------------------------

PRESENTATION_RE = re.compile(r'^presentation\s+"(.+)"\s*:$')
SECTION_RE = re.compile(r'^section\s+"(.+)"\s*:$')
SLIDE_RE = re.compile(r'^slide\s+"(.+)"\s*:$')
THEME_RE = re.compile(r'^theme\s+"(.+)"\s*:$')


def parse_pml(text: str) -> PresentationDoc:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PML")

    m = PRESENTATION_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PML must start with: presentation "Title":')

    doc = PresentationDoc(title=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        text = line.text
        key, val = split_key_value(text)

        if text == "meta:":
            block, i = parse_generic_block(lines, i + 1, 4)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: meta must be a mapping")
            doc.meta = block
            continue

        if key == "use style":
            doc.style_file = parse_scalar(val or "")
            i += 1
            continue

        if key == "cover_layout":
            doc.cover_layout = str(parse_scalar(val or "title-slide"))
            i += 1
            continue

        if key == "cover" and val == "":
            block, i = parse_generic_block(lines, i + 1, 4)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: cover must be a mapping")
            doc.cover = block
            continue

        sec_match = SECTION_RE.match(text)
        if sec_match:
            section, i = parse_section(lines, i)
            doc.sections.append(section)
            continue

        raise SyntaxError(f"Line {line.lineno}: unknown PML statement: {text}")

    return doc


def parse_section(lines: List[Line], start: int) -> Tuple[Section, int]:
    m = SECTION_RE.match(lines[start].text)
    assert m is not None
    section = Section(title=m.group(1))
    base_indent = lines[start].indent
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= base_indent:
            break
        if line.indent != base_indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected section child indent")

        key, val = split_key_value(line.text)

        if key == "header_layout":
            section.header_layout = str(parse_scalar(val or "section-header"))
            i += 1
            continue

        if key == "header" and val == "":
            block, i = parse_generic_block(lines, i + 1, line.indent + 2)
            if not isinstance(block, dict):
                raise SyntaxError(f"Line {line.lineno}: section header must be a mapping")
            section.header = block
            continue

        slide_match = SLIDE_RE.match(line.text)
        if slide_match:
            slide, i = parse_slide(lines, i)
            section.slides.append(slide)
            continue

        raise SyntaxError(
            f"Line {line.lineno}: section accepts header_layout, header, or slide blocks only"
        )

    return section, i


def parse_slide(lines: List[Line], start: int) -> Tuple[Slide, int]:
    m = SLIDE_RE.match(lines[start].text)
    assert m is not None
    slide = Slide(title=m.group(1))
    base_indent = lines[start].indent
    i = start + 1

    while i < len(lines):
        line = lines[i]
        if line.indent <= base_indent:
            break
        if line.indent != base_indent + 2:
            raise SyntaxError(f"Line {line.lineno}: expected slide child indent")

        key, val = split_key_value(line.text)

        if key == "layout":
            slide.layout = str(parse_scalar(val or ""))
            i += 1
            continue

        if key == "intent":
            slide.intent = str(parse_scalar(val or ""))
            i += 1
            continue

        if val == "":
            if i + 1 >= len(lines) or lines[i + 1].indent <= line.indent:
                slide.blocks[key] = ""
                i += 1
            else:
                child, new_i = parse_generic_block(lines, i + 1, lines[i + 1].indent)
                slide.blocks[key] = child
                i = new_i
        else:
            slide.blocks[key] = parse_scalar(val or "")
            i += 1

    return slide, i


# -----------------------------------------------------------------------------
# PSL parser
# -----------------------------------------------------------------------------

def parse_psl(text: str) -> Theme:
    lines = tokenize_lines(text)
    if not lines:
        raise SyntaxError("Empty PSL")

    m = THEME_RE.match(lines[0].text)
    if not m:
        raise SyntaxError('PSL must start with: theme "Name":')

    theme = Theme(name=m.group(1))
    i = 1

    while i < len(lines):
        line = lines[i]
        if line.indent != 2:
            raise SyntaxError(f"Line {line.lineno}: expected top-level indent 2")

        key, val = split_key_value(line.text)
        if val != "":
            raise SyntaxError(f"Line {line.lineno}: top-level PSL keys must be blocks")

        block, i = parse_generic_block(lines, i + 1, 4)
        if not isinstance(block, dict):
            raise SyntaxError(f"Line {line.lineno}: PSL block must be a mapping")

        if key == "page":
            theme.page = block
        elif key == "colors":
            theme.colors = block
        elif key == "fonts":
            theme.fonts = block
        elif key.startswith("presentation."):
            layout_name = key.split(".", 1)[1]
            theme.presentation_layouts[layout_name] = block
        elif key.startswith("section."):
            layout_name = key.split(".", 1)[1]
            theme.section_layouts[layout_name] = block
        elif key.startswith("slide."):
            layout_name = key.split(".", 1)[1]
            theme.slide_layouts[layout_name] = block
        else:
            # Unknown top-level blocks are treated as slide layouts for compatibility.
            theme.slide_layouts[key] = block

    return theme


# -----------------------------------------------------------------------------
# Render IR construction
# -----------------------------------------------------------------------------

def resolve_token(value: Any, theme: Theme) -> Any:
    if isinstance(value, str):
        if value in theme.colors:
            return theme.colors[value]
        if value in theme.fonts:
            return theme.fonts[value]
    return value


def style_value(style: Dict[str, Any], key: str, theme: Theme, default: Any = None) -> Any:
    return resolve_token(style.get(key, default), theme)


def normalize_text_block(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, dict):
        if "text" in value:
            return str(value["text"])
        if "heading" in value:
            return str(value["heading"])
    return str(value)


def make_text_object(kind: str, value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [60, 40])
    return {
        "type": "TextBox",
        "role": kind,
        "text": normalize_text_block(value),
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 80),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 28),
        "color": style_value(style, "color", theme, "#000000"),
    }


def make_bullet_object(value: Any, style: Dict[str, Any], theme: Theme) -> Dict[str, Any]:
    pos = style.get("position", [90, 160])
    if not isinstance(value, list):
        value = [value]
    return {
        "type": "BulletList",
        "items": [str(v) for v in value],
        "x": pos[0],
        "y": pos[1],
        "w": style.get("width", 1000),
        "h": style.get("height", 420),
        "font": style_value(style, "font", theme, "Aptos"),
        "size": style.get("size", 24),
        "color": style_value(style, "color", theme, "#000000"),
        "line_gap": style.get("line_gap", 8),
    }


def make_column_objects(role: str, value: Any, style: Dict[str, Any], theme: Theme) -> List[Dict[str, Any]]:
    pos = style.get("position", [70, 150])
    width = style.get("width", 520)
    objects: List[Dict[str, Any]] = []

    if isinstance(value, dict):
        if "heading" in value:
            objects.append(
                {
                    "type": "TextBox",
                    "role": f"{role}.heading",
                    "text": str(value["heading"]),
                    "x": pos[0],
                    "y": pos[1],
                    "w": width,
                    "h": 50,
                    "font": theme.fonts.get("heading", "Aptos Display"),
                    "size": 24,
                    "color": theme.colors.get("primary", "#000000"),
                }
            )
        if "bullets" in value:
            objects.append(
                {
                    "type": "BulletList",
                    "items": [str(x) for x in value["bullets"]],
                    "x": pos[0],
                    "y": pos[1] + 60,
                    "w": width,
                    "h": 300,
                    "font": theme.fonts.get("body", "Aptos"),
                    "size": 20,
                    "color": theme.colors.get("text", "#000000"),
                    "line_gap": 6,
                }
            )
    else:
        objects.append(
            {
                "type": "TextBox",
                "role": role,
                "text": str(value),
                "x": pos[0],
                "y": pos[1],
                "w": width,
                "h": 300,
                "font": theme.fonts.get("body", "Aptos"),
                "size": 20,
                "color": theme.colors.get("text", "#000000"),
            }
        )

    return objects


def build_cover_slide_ir(doc: PresentationDoc, theme: Theme) -> Dict[str, Any]:
    layout_style = theme.presentation_layouts.get(doc.cover_layout, {})
    objects: List[Dict[str, Any]] = []

    # Main presentation title always comes from `presentation "Title"`.
    title_style = layout_style.get("title", {})
    objects.append(make_text_object("Title", doc.title, title_style, theme))

    # Other cover blocks come from `cover:`.
    for key, value in doc.cover.items():
        style = layout_style.get(key, {})
        if key == "bullets":
            objects.append(make_bullet_object(value, style, theme))
        else:
            objects.append(make_text_object(key, value, style, theme))

    background = style_value(
        layout_style,
        "background",
        theme,
        theme.page.get("background", "#FFFFFF"),
    )

    return {
        "type": "RenderSlide",
        "kind": "presentation-cover",
        "section": None,
        "slide_title": doc.title,
        "layout": doc.cover_layout,
        "background": background,
        "objects": objects,
        "notes": doc.cover.get("notes", ""),
    }


def build_section_header_ir(section: Section, theme: Theme) -> Dict[str, Any]:
    layout_style = theme.section_layouts.get(section.header_layout, {})
    objects: List[Dict[str, Any]] = []

    title_style = layout_style.get("title", {})
    objects.append(make_text_object("SectionTitle", section.title, title_style, theme))

    for key, value in section.header.items():
        style = layout_style.get(key, {})
        if key == "bullets":
            objects.append(make_bullet_object(value, style, theme))
        else:
            objects.append(make_text_object(key, value, style, theme))

    background = style_value(
        layout_style,
        "background",
        theme,
        theme.page.get("background", "#FFFFFF"),
    )

    return {
        "type": "RenderSlide",
        "kind": "section-header",
        "section": section.title,
        "slide_title": section.title,
        "layout": section.header_layout,
        "background": background,
        "objects": objects,
        "notes": section.header.get("notes", ""),
    }


def build_content_slide_ir(section: Section, slide: Slide, theme: Theme) -> Dict[str, Any]:
    layout_style = theme.slide_layouts.get(slide.layout, {})
    objects: List[Dict[str, Any]] = []

    if "title" in slide.blocks:
        style = layout_style.get("title", {})
        objects.append(make_text_object("Title", slide.blocks["title"], style, theme))

    if "subtitle" in slide.blocks:
        style = layout_style.get("subtitle", layout_style.get("title", {}))
        objects.append(make_text_object("Subtitle", slide.blocks["subtitle"], style, theme))

    if "bullets" in slide.blocks:
        style = layout_style.get("bullets", {})
        objects.append(make_bullet_object(slide.blocks["bullets"], style, theme))

    for side in ["left", "right"]:
        if side in slide.blocks:
            style = layout_style.get(side, {})
            objects.extend(make_column_objects(side, slide.blocks[side], style, theme))

    background = style_value(
        layout_style,
        "background",
        theme,
        theme.page.get("background", "#FFFFFF"),
    )

    return {
        "type": "RenderSlide",
        "kind": "content",
        "section": section.title,
        "slide_title": slide.title,
        "layout": slide.layout,
        "background": background,
        "objects": objects,
        "notes": slide.blocks.get("notes", ""),
    }


def build_render_ir(doc: PresentationDoc, theme: Theme) -> Dict[str, Any]:
    slides_ir: List[Dict[str, Any]] = []

    # 1. Presentation title / cover slide
    slides_ir.append(build_cover_slide_ir(doc, theme))

    # 2. Section header + content slides
    for section in doc.sections:
        slides_ir.append(build_section_header_ir(section, theme))
        for slide in section.slides:
            slides_ir.append(build_content_slide_ir(section, slide, theme))

    return {
        "type": "RenderPresentation",
        "title": doc.title,
        "page": theme.page,
        "slides": slides_ir,
    }


# -----------------------------------------------------------------------------
# PowerPoint renderer
# -----------------------------------------------------------------------------

def hex_to_rgb(hex_color: str):
    if RGBColor is None:
        raise RuntimeError("python-pptx is not installed")
    if not hex_color:
        hex_color = "#000000"
    hex_color = str(hex_color).strip().lstrip("#")
    if len(hex_color) != 6:
        hex_color = "000000"
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def px_to_inches(value: float) -> float:
    """Treat DSL coordinates as pixel-like values on a 1280x720 canvas."""
    return value / 96.0


def add_textbox(slide, obj: Dict[str, Any], bullet: bool = False) -> None:
    shape = slide.shapes.add_textbox(
        Inches(px_to_inches(obj["x"])),
        Inches(px_to_inches(obj["y"])),
        Inches(px_to_inches(obj["w"])),
        Inches(px_to_inches(obj.get("h", 100))),
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    if bullet:
        for idx, item in enumerate(obj["items"]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = obj["font"]
            p.font.size = Pt(obj["size"])
            p.font.color.rgb = hex_to_rgb(obj["color"])
    else:
        p = tf.paragraphs[0]
        p.text = obj["text"]
        p.font.name = obj["font"]
        p.font.size = Pt(obj["size"])
        p.font.color.rgb = hex_to_rgb(obj["color"])


def render_pptx(render_ir: Dict[str, Any], output_path: str) -> None:
    if PPTXPresentation is None:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = PPTXPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for slide_ir in render_ir["slides"]:
        slide = prs.slides.add_slide(blank_layout)

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_to_rgb(slide_ir.get("background", "#FFFFFF"))

        for obj in slide_ir["objects"]:
            if obj["type"] == "TextBox":
                add_textbox(slide, obj, bullet=False)
            elif obj["type"] == "BulletList":
                add_textbox(slide, obj, bullet=True)

        notes = slide_ir.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    prs.save(output_path)


# -----------------------------------------------------------------------------
# HTML renderer
# -----------------------------------------------------------------------------

def render_html(render_ir: Dict[str, Any], output_path: str) -> None:
    parts: List[str] = []
    parts.append("<!doctype html>")
    parts.append("<html><head><meta charset='utf-8'>")
    parts.append("<title>" + html.escape(render_ir["title"]) + "</title>")
    parts.append(
        """
<style>
body { margin: 0; font-family: Arial, sans-serif; background: #eee; }
.deck { display: flex; flex-direction: column; gap: 24px; padding: 24px; }
.slide { position: relative; width: 1280px; height: 720px; box-shadow: 0 8px 24px rgba(0,0,0,.18); overflow: hidden; }
.textbox, .bullets { position: absolute; white-space: pre-wrap; }
ul { margin: 0; padding-left: 1.2em; }
li { margin-bottom: .4em; }
.kind { position: absolute; right: 20px; bottom: 14px; font-size: 14px; opacity: .4; }
</style>
"""
    )
    parts.append("</head><body><div class='deck'>")

    for slide in render_ir["slides"]:
        bg = slide.get("background", "#FFFFFF")
        parts.append(f"<section class='slide' style='background:{html.escape(str(bg))}'>")
        for obj in slide["objects"]:
            common_style = (
                f"left:{obj['x']}px; top:{obj['y']}px; width:{obj['w']}px; "
                f"font-family:{html.escape(str(obj['font']))}; "
                f"font-size:{obj['size']}px; color:{html.escape(str(obj['color']))};"
            )
            if obj["type"] == "TextBox":
                parts.append(
                    f"<div class='textbox' style='{common_style}'>"
                    f"{html.escape(str(obj['text']))}</div>"
                )
            elif obj["type"] == "BulletList":
                parts.append(f"<div class='bullets' style='{common_style}'><ul>")
                for item in obj["items"]:
                    parts.append(f"<li>{html.escape(str(item))}</li>")
                parts.append("</ul></div>")
        parts.append(f"<div class='kind'>{html.escape(str(slide.get('kind', '')))}</div>")
        parts.append("</section>")

    parts.append("</div></body></html>")
    Path(output_path).write_text("\n".join(parts), encoding="utf-8")


# -----------------------------------------------------------------------------
# Demo input
# -----------------------------------------------------------------------------

DEMO_PML3 = '''
presentation "AI Strategy 2026":
  meta:
    author: "Viettel AI Lab"
    language: vi
    format: pptx

  use style: "corporate_v3.psl"

  cover_layout: title-slide
  cover:
    subtitle: Chiến lược ứng dụng AI trong doanh nghiệp
    author: Viettel AI Lab
    date: 2026
    notes:
      Đây là slide bìa của toàn bộ deck.

  section "Bối cảnh":
    header_layout: section-header
    header:
      subtitle: Vì sao doanh nghiệp cần hệ thống tổng hợp tri thức?
      notes:
        Slide này đóng vai trò ngăn chương.

    slide "Vấn đề hiện tại":
      layout: title-bullets
      intent: explain_problem

      title:
        Quá tải thông tin trong doanh nghiệp

      bullets:
        - Dữ liệu phân tán ở nhiều hệ thống
        - Báo cáo thủ công mất nhiều thời gian
        - Lãnh đạo cần insight nhanh hơn

      notes:
        Nhấn mạnh vấn đề không nằm ở thiếu dữ liệu, mà là thiếu khả năng tổng hợp.

    slide "Mục tiêu":
      layout: two-column
      intent: propose_solution

      title:
        Xây dựng hệ thống tổng hợp tri thức tự động

      left:
        heading: Đầu vào
        bullets:
          - Văn bản
          - Báo cáo
          - Dữ liệu mạng xã hội

      right:
        heading: Đầu ra
        bullets:
          - Tóm tắt
          - Dashboard
          - Slide trình bày

  section "Giải pháp":
    header_layout: section-header
    header:
      subtitle: Thiết kế pipeline từ dữ liệu thô đến tri thức hành động

    slide "Pipeline đề xuất":
      layout: title-bullets
      intent: describe_pipeline

      title:
        Pipeline xử lý và tổng hợp tri thức

      bullets:
        - Thu thập dữ liệu từ nhiều nguồn
        - Làm sạch và chuẩn hóa văn bản
        - Trích xuất insight bằng mô hình ngôn ngữ
        - Render báo cáo, dashboard và slide
'''


DEMO_PSL3 = '''
theme "corporate-v3":
  page:
    size: widescreen
    background: "#FFFFFF"

  colors:
    primary: "#003A8C"
    secondary: "#E6F0FF"
    dark: "#0F172A"
    text: "#1F1F1F"
    muted: "#666666"
    white: "#FFFFFF"

  fonts:
    heading: "Aptos Display"
    body: "Aptos"

  presentation.title-slide:
    background: primary

    title:
      font: heading
      size: 48
      color: white
      position: [90, 210]
      width: 1100
      height: 90

    subtitle:
      font: body
      size: 26
      color: secondary
      position: [95, 315]
      width: 1050
      height: 60

    author:
      font: body
      size: 21
      color: white
      position: [95, 430]
      width: 700
      height: 40

    date:
      font: body
      size: 21
      color: white
      position: [95, 470]
      width: 500
      height: 40

  section.section-header:
    background: secondary

    title:
      font: heading
      size: 44
      color: primary
      position: [90, 250]
      width: 1080
      height: 90

    subtitle:
      font: body
      size: 24
      color: dark
      position: [95, 345]
      width: 1000
      height: 70

  slide.title-bullets:
    background: white

    title:
      font: heading
      size: 36
      color: primary
      position: [60, 40]
      width: 1100
      height: 90

    bullets:
      font: body
      size: 24
      color: text
      position: [90, 170]
      width: 1000
      height: 420
      line_gap: 12

  slide.two-column:
    background: white

    title:
      font: heading
      size: 34
      color: primary
      position: [60, 40]
      width: 1120
      height: 80

    left:
      position: [70, 160]
      width: 520

    right:
      position: [680, 160]
      width: 520
'''


def main() -> None:
    doc = parse_pml(DEMO_PML3)
    theme = parse_psl(DEMO_PSL3)
    render_ir = build_render_ir(doc, theme)

    render_html(render_ir, "demo_output_v3.html")

    if PPTXPresentation is not None:
        render_pptx(render_ir, "demo_output_v3.pptx")
        print("Generated demo_output_v3.pptx")
    else:
        print("Skipped PPTX: python-pptx is not installed")

    print("Generated demo_output_v3.html")


if __name__ == "__main__":
    main()
