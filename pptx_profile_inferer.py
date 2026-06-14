from __future__ import annotations

import json
import re
import statistics
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def safe_text(s: str) -> str:
    return re.sub(r"\s+", " ", s or "").strip()


def emu_to_pt(x) -> float | None:
    if x is None:
        return None
    return round(x / 12700, 2)


def emu_to_inch(x) -> float | None:
    if x is None:
        return None
    return round(x / 914400, 3)


def get_shape_kind(shape) -> str:
    t = shape.shape_type
    if t == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if t == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if t == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if t == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "has_text_frame", False):
        return "text"
    return str(t)


def iter_shapes_recursive(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_recursive(shape.shapes)


def extract_run_style(run) -> dict[str, Any]:
    font = run.font

    style = {
        "text": run.text,
        "font_name": font.name,
        "font_size_pt": emu_to_pt(font.size) if font.size else None,
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color": None,
    }

    try:
        if font.color and font.color.rgb:
            style["color"] = f"#{font.color.rgb}"
    except Exception:
        pass

    return style


def extract_paragraph_style(paragraph) -> dict[str, Any]:
    return {
        "text": safe_text(paragraph.text),
        "level": paragraph.level,
        "alignment": str(paragraph.alignment) if paragraph.alignment else None,
        "runs": [
            extract_run_style(run)
            for run in paragraph.runs
            if run.text
        ],
    }


def extract_shape_style(shape) -> dict[str, Any]:
    item = {
        "name": shape.name,
        "kind": get_shape_kind(shape),
        "left_in": emu_to_inch(shape.left),
        "top_in": emu_to_inch(shape.top),
        "width_in": emu_to_inch(shape.width),
        "height_in": emu_to_inch(shape.height),
        "text": "",
        "paragraphs": [],
    }

    if getattr(shape, "has_text_frame", False):
        item["text"] = safe_text(shape.text)
        item["paragraphs"] = [
            extract_paragraph_style(p)
            for p in shape.text_frame.paragraphs
            if safe_text(p.text)
        ]

    return item


def extract_theme_from_zip(pptx_path: str | Path) -> dict[str, Any]:
    pptx_path = Path(pptx_path)

    result = {
        "theme_files": [],
        "theme_name": None,
        "colors": {},
        "major_fonts": [],
        "minor_fonts": [],
    }

    with zipfile.ZipFile(pptx_path) as z:
        theme_files = [
            n for n in z.namelist()
            if n.startswith("ppt/theme/theme") and n.endswith(".xml")
        ]
        result["theme_files"] = theme_files

        if not theme_files:
            return result

        root = ET.fromstring(z.read(theme_files[0]))
        result["theme_name"] = root.attrib.get("name")

        clr_scheme = root.find(".//a:clrScheme", NS)
        if clr_scheme is not None:
            for child in clr_scheme:
                key = child.tag.split("}")[-1]

                srgb = child.find(".//a:srgbClr", NS)
                sysclr = child.find(".//a:sysClr", NS)

                if srgb is not None:
                    result["colors"][key] = "#" + srgb.attrib.get("val", "")
                elif sysclr is not None:
                    result["colors"][key] = "#" + sysclr.attrib.get("lastClr", "")

        major = root.findall(".//a:fontScheme/a:majorFont//a:latin", NS)
        minor = root.findall(".//a:fontScheme/a:minorFont//a:latin", NS)

        result["major_fonts"] = [
            x.attrib.get("typeface")
            for x in major
            if x.attrib.get("typeface")
        ]
        result["minor_fonts"] = [
            x.attrib.get("typeface")
            for x in minor
            if x.attrib.get("typeface")
        ]

    return result


def collect_style_stats(slides: list[dict[str, Any]]) -> dict[str, Any]:
    font_names = Counter()
    font_sizes = []
    colors = Counter()
    shape_kinds = Counter()

    bold_count = 0
    italic_count = 0
    run_count = 0

    for slide in slides:
        for shape in slide["shapes"]:
            shape_kinds[shape["kind"]] += 1

            for p in shape.get("paragraphs", []):
                for r in p.get("runs", []):
                    run_count += 1

                    if r.get("font_name"):
                        font_names[r["font_name"]] += 1

                    if r.get("font_size_pt"):
                        font_sizes.append(r["font_size_pt"])

                    if r.get("color"):
                        colors[r["color"]] += 1

                    if r.get("bold"):
                        bold_count += 1

                    if r.get("italic"):
                        italic_count += 1

    return {
        "font_names": dict(font_names.most_common()),
        "font_sizes_pt": {
            "all": sorted(set(font_sizes)),
            "median": statistics.median(font_sizes) if font_sizes else None,
            "min": min(font_sizes) if font_sizes else None,
            "max": max(font_sizes) if font_sizes else None,
        },
        "text_colors": dict(colors.most_common()),
        "shape_kinds": dict(shape_kinds.most_common()),
        "bold_ratio": round(bold_count / run_count, 3) if run_count else None,
        "italic_ratio": round(italic_count / run_count, 3) if run_count else None,
    }


def infer_profile(profile: dict[str, Any]) -> dict[str, Any]:
    stats = profile["style_stats"]
    theme = profile["theme"]

    font_names = list(stats.get("font_names", {}).keys())
    text_colors = list(stats.get("text_colors", {}).keys())
    theme_colors = list(theme.get("colors", {}).values())

    median_size = stats.get("font_sizes_pt", {}).get("median")

    shape_counts = [len(s["shapes"]) for s in profile["slides"]]
    avg_shapes = statistics.mean(shape_counts) if shape_counts else 0

    text_lengths = [
        len(safe_text(" ".join(sh.get("text", "") for sh in s["shapes"])))
        for s in profile["slides"]
    ]
    avg_text = statistics.mean(text_lengths) if text_lengths else 0

    if avg_text < 120:
        density = "low_text / presentation-oriented"
    elif avg_text < 350:
        density = "medium_text"
    else:
        density = "high_text / report-like"

    visual_style = []

    if avg_shapes > 8:
        visual_style.append("shape-heavy")

    if median_size and median_size >= 24:
        visual_style.append("large-readable-fonts")

    palette = text_colors or theme_colors

    if len(palette) <= 3:
        visual_style.append("limited-color-palette")
    elif len(palette) >= 6:
        visual_style.append("varied-color-palette")

    notes = []

    if not font_names:
        notes.append(
            "Font values may be inherited from theme/master and are not always exposed directly by python-pptx."
        )

    return {
        "likely_font_family": (
            font_names[:3]
            or theme.get("major_fonts")
            or theme.get("minor_fonts")
        ),
        "likely_palette": palette[:8],
        "density": density,
        "avg_shapes_per_slide": round(avg_shapes, 2),
        "avg_text_chars_per_slide": round(avg_text, 2),
        "visual_style": visual_style,
        "notes": notes,
    }


def infer_pptx_profile(
    pptx_path: str | Path,
    llm_style_reasoner: Callable[[dict[str, Any]], dict[str, Any]] | None = None,
) -> dict[str, Any]:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "index": i,
            "layout_name": slide.slide_layout.name,
            "shapes": [
                extract_shape_style(sh)
                for sh in iter_shapes_recursive(slide.shapes)
            ],
        }
        slides.append(slide_data)

    profile = {
        "file": str(pptx_path),
        "page_size": {
            "width_in": emu_to_inch(prs.slide_width),
            "height_in": emu_to_inch(prs.slide_height),
        },
        "theme": extract_theme_from_zip(pptx_path),
        "slides": slides,
    }

    profile["style_stats"] = collect_style_stats(slides)
    profile["inference"] = infer_profile(profile)

    if llm_style_reasoner:
        profile["llm_inference"] = llm_style_reasoner(profile)

    return profile


def save_profile_json(pptx_path: str | Path, out_json: str | Path) -> None:
    profile = infer_pptx_profile(pptx_path)
    Path(out_json).write_text(
        json.dumps(profile, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("--out", default="deck_profile.json")
    args = parser.parse_args()

    save_profile_json(args.pptx, args.out)